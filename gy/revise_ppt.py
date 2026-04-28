import shutil
from pptx import Presentation

shutil.copy("gy/gy_ppt.pptx", "gy/backup/gy_ppt_pre_revision.pptx")
prs = Presentation("gy/gy_ppt.pptx")


def get_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def replace_run(shape, para_idx, run_idx, old, new):
    tf = shape.text_frame
    para = tf.paragraphs[para_idx]
    run = para.runs[run_idx]
    run.text = run.text.replace(old, new)


def set_run(shape, para_idx, run_idx, text):
    tf = shape.text_frame
    para = tf.paragraphs[para_idx]
    para.runs[run_idx].text = text


# ── Slide 10: competitor comparison ──────────────────────────────────────────
s10 = prs.slides[9]

# Calypso/Murex 信创适配: "✗" → "✗ 失格"
set_run(get_shape(s10, "VItem0_5"), 1, 0, "✗ 失格")

# Takeaway: add 9/9 reference
set_run(get_shape(s10, "Takeaway"), 0, 1,
        "  |  唯一同时满足：FICC前台基因 + 买方资管平台 + 信创合规认证 + 一体化统一部署  |  综合评分 9/9")

# ── Slide 12: why now ─────────────────────────────────────────────────────────
s12 = prs.slides[11]

# Subtitle: 18个月 → 24个月
replace_run(get_shape(s12, "TextBox 3"), 0, 0, "18个月黄金建设窗口", "24个月黄金建设窗口，M6 Quick Win验收节点确认价值")

# Takeaway: two runs both contain 18个月
replace_run(get_shape(s12, "Takeaway"), 0, 0, "18个月追赶头部3-5年", "24个月追赶头部3-5年")
replace_run(get_shape(s12, "Takeaway"), 0, 1, "国元18个月达到同等水平", "国元24个月达到同等水平")

# Supp: Phase 1 → M6 Quick Win
replace_run(get_shape(s12, "Supp"), 0, 0, "Phase 1在6个月内", "M6 Quick Win在6个月内")

# ── Slide 33: roadmap ────────────────────────────────────────────────────────
s33 = prs.slides[32]

# Title: 18个月 → 24个月
replace_run(get_shape(s33, "TextBox 2"), 0, 0, "三期18个月实施路径", "三期24个月实施路径")

# Subtitle: update Phase 1 reference
replace_run(get_shape(s33, "TextBox 3"), 0, 0,
            "Phase 1六个月灯塔效应", "M6 Quick Win 验收触发首期付款")

# Phase 1 header: 0-6月 → M1-M12 筑基期
set_run(get_shape(s33, "TextBox 6"), 0, 0, "Phase 1（M1 – M12  筑基期）")

# Phase 1 bullet 5: update last bullet to cover full M7-M12 scope
set_run(get_shape(s33, "TextBox 8"), 5, 0,
        "▸ ★ M6 Quick Win：一屏可见 · M7-M12量化策略+绩效归因+指令管理+投前合规")

# Phase 1 milestone
set_run(get_shape(s33, "TextBox 10"), 0, 0,
        "里程碑：M6 Quick Win验收 · M12千亿全资产实时一屏可见")

# Phase 2 header: 6-12月 → M13-M18 赋能期
set_run(get_shape(s33, "TextBox 14"), 0, 0, "Phase 2（M13 – M18  赋能期）")

# Phase 2 milestone
set_run(get_shape(s33, "TextBox 18"), 0, 0, "里程碑：M18 风控+量化策略全上线")

# Phase 3 header: 12-18月 → M19-M24 进化期
set_run(get_shape(s33, "TextBox 22"), 0, 0, "Phase 3（M19 – M24  进化期）")

# ── Slide 34: day-in-life ────────────────────────────────────────────────────
s34 = prs.slides[33]

# Title: Phase 1灯塔效应 → M6 Quick Win效应
replace_run(get_shape(s34, "TextBox 2"), 0, 0, "Phase 1灯塔效应", "M6 Quick Win效应")

# Footer: Phase 1交付承诺 → M6 Quick Win交付承诺
replace_run(get_shape(s34, "TextBox 6"), 0, 0, "Phase 1交付承诺", "M6 Quick Win交付承诺")
replace_run(get_shape(s34, "TextBox 6"), 1, 0, "Phase 1本身就是", "M6 Quick Win本身就是")

prs.save("gy/gy_ppt.pptx")
print("✓ gy_ppt.pptx saved — 4 slides updated (10, 12, 33, 34)")
