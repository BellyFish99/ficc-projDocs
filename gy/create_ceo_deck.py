"""
Generate CEO Deck v2.1 PowerPoint from content markdown.
Creates a professional 20-slide deck with all content, tables, and speaker notes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === Color Palette (Professional dark blue theme) ===
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MEDIUM_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
ACCENT_GOLD = RGBColor(0xC4, 0x9A, 0x2A)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=14,
                bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """Add a textbox with text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(text_frame, text, font_size=14, bold=False, color=DARK_GRAY,
             alignment=PP_ALIGN.LEFT, space_before=Pt(4), font_name="Microsoft YaHei"):
    """Add a paragraph to existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_table_slide(slide, left, top, width, row_height, headers, rows,
                    header_color=DARK_BLUE, font_size=11):
    """Add a table to a slide."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width,
                                          Emu(row_height * num_rows))
    table = table_shape.table

    # Set column widths proportionally
    col_width = width // num_cols
    for i in range(num_cols):
        table.columns[i].width = col_width

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Microsoft YaHei"

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK_GRAY
                p.font.name = "Microsoft YaHei"

    return table_shape


def set_notes(slide, text):
    """Set speaker notes."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def make_title_bar(slide, title_text, subtitle_text=""):
    """Add a dark blue title bar at top."""
    add_shape(slide, Inches(0), Inches(0), W, Inches(1.2), fill_color=DARK_BLUE)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
                title_text, font_size=24, bold=True, color=WHITE)
    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.4),
                    subtitle_text, font_size=13, color=ACCENT_GOLD)


def make_act_divider(slide, act_text, act_subtitle):
    """Create an act divider slide."""
    add_bg(slide, DARK_BLUE)
    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
                act_text, font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
                act_subtitle, font_size=18, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1 — Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK_BLUE)

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
            "国元证券全资产POMS平台解决方案",
            font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
            "让千亿自营资金拥有世界级的组合管理能力",
            font_size=22, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

# Separator line
add_shape(slide, Inches(4.5), Inches(4.2), Inches(4), Inches(0.03), fill_color=ACCENT_GOLD)

add_textbox(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.5),
            "华锐技术  |  2026年4月",
            font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.4),
            "— 保密文件，仅供国元证券内部决策使用 —",
            font_size=11, color=RGBColor(0x88, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)

set_notes(slide, "封面页。开场前确认参会人员名单和预计时间（约55分钟+Q&A）。")


# ============================================================
# SLIDE 2 — Executive Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "执行摘要",
               "一页看懂：千亿自营资金需要什么、我们提供什么、值多少")

headers = ["维度", "要点"]
rows = [
    ["国元现状", "自营已成第一大收入来源（2025H1占比43%），近千亿AUM，战略转型方向清晰"],
    ["战略需要", "四大核心能力跃升：配得优 · 算得快 · 控得稳 · 连得通"],
    ["解决方案", "华锐POMS平台：全资产组合管理 + 量化策略 + 实时风控 + 指令闭环 + 精确成本核算"],
    ["生态协同", "与金证IBOR、金仕达行情中心共生——IBOR是数据引擎，POMS是业务智能"],
    ["实施路径", "三期18个月，Phase 1六个月见效"],
    ["保守价值", "年化直接价值1-1.8亿 + 极端事件保守避损1-2亿/次（系统赋能部分）"],
    ["投资保障", "Phase 1灯塔效应——6个月内CEO可看到全公司自营实时全貌"],
]
add_table_slide(slide, Inches(0.6), Inches(1.5), Inches(12.1), Emu(400000),
                headers, rows, font_size=11)

# Core judgment box
box = add_shape(slide, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.3),
                fill_color=RGBColor(0xF0, 0xF4, 0xF8), line_color=MEDIUM_BLUE)
tb = add_textbox(slide, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.1),
                 "核心判断：POMS不是一个IT系统，而是投资能力的基础设施。", font_size=13, bold=True, color=DARK_BLUE)
add_para(tb.text_frame, "价值的充分释放需要系统+人员+流程+管理四个支柱协同——系统是必要条件，没有它其他三个支柱无从发力。",
         font_size=11, color=MED_GRAY)
add_para(tb.text_frame, "POMS与国元已建技术生态是共生关系：金证IBOR是数据引擎，金仕达是信息高速公路，POMS是自动驾驶系统。",
         font_size=11, color=MED_GRAY)

set_notes(slide, "各位领导，一页纸看懂我们的方案。国元近千亿自营资金，战略目标非常清晰。要实现全资产精细化管理，需要四大能力跃升——配得优、算得快、控得稳、连得通。华锐POMS就是支撑这四个跃升的业务智能平台，它和国元已经在建的IBOR、行情中心形成完整的技术生态。")


# ============================================================
# ACT 1 DIVIDER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_act_divider(slide, "Act 1: 我们深刻理解国元",
                 "不是来卖系统的，是来解决战略问题的")


# ============================================================
# SLIDE 3 — Guoyuan Business: 4x Growth
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "国元自营业务：4年4倍增长，已成第一大支柱",
               "自营收入4年增长4倍至占比43%，'去方向性+FICC+资本中介'三位一体战略已定")

headers = ["年份", "自营收入", "占比", "财富信用", "关键转折"]
rows = [
    ["2022", "4.48亿", "8%", "19.29亿", "战略转型元年：确立'去方向性'"],
    ["2023", "13.07亿", "21%", "15.05亿", "+192%爆发，FICC资本中介快速发展"],
    ["2024", "23.11亿", "29%", "16.96亿", "首次超越财富业务，成为第一大收入"],
    ["2025H1", "14.60亿", "43%", "9.98亿", "半年超2023全年，领先财富1.46倍"],
]
add_table_slide(slide, Inches(0.6), Inches(1.5), Inches(12.1), Emu(400000),
                headers, rows, font_size=12)

# Org restructure section
tb = add_textbox(slide, Inches(0.6), Inches(3.8), Inches(12), Inches(0.4),
                 "2025年8月组织架构升级——三重买方信号", font_size=16, bold=True, color=DARK_BLUE)

headers2 = ["调整", "战略信号"]
rows2 = [
    ["自营业务委员会 → 金融市场业务委员会", "从投资本位到市场服务本位——需要统一组合管理平台"],
    ["固定收益部 → FICC业务总部", "从单一债券到综合FICC——需要跨资产组合管理"],
    ["权益投资部 → 证券投资部", "去方向性+多策略——需要实时风控+相对价值工具"],
    ["创新金融部 → 创新金融业务总部（升格）", "衍生品从边缘到核心——需要统一定价+成本核算"],
]
add_table_slide(slide, Inches(0.6), Inches(4.3), Inches(12.1), Emu(350000),
                headers2, rows2, font_size=11)

set_notes(slide, "我们对国元做了深入研究。自营业务4年增长4倍，2025年上半年占比43%。2025年8月的组织架构升级更说明问题——每一个调整都指向同一个方向：需要一个统一的组合管理平台。三大基础能力中明确提出的金融科技能力，正是POMS要提供的。")


# ============================================================
# SLIDE 4 — Business Requirements
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "我们理解国元的需求：资质已就位，缺的是统一平台",
               "5年积累完整FICC资质拼图——资质是通行证，POMS是驾驶能力")

# Four core requirements
tb = add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(2.5), "", font_size=14)
tf = tb.text_frame
tf.word_wrap = True

items = [
    ("1. 以稳定收益为目标", "低波动6%，从方向性投资转向精细化管理"),
    ("2. 全资产品种覆盖", "不是固收单品种，而是股票、黄金、债券、衍生品真正的多资产"),
    ("3. 丰富投资策略", "相对价值（利差/基差/新老券）+ 多策略多timeframe"),
    ("4. 以组合管理为核心", "不是以交易为核心，这是Aladdin级的思维高度"),
]
for i, (title, desc) in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"{title} — {desc}"
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Microsoft YaHei"
    p.space_before = Pt(8)

# Qualifications table
headers = ["已获资质", "解锁的业务能力", "统一平台能做什么"]
rows = [
    ["利率互换交易资格（2020）", "利率衍生品对冲", "统一定价引擎支撑利差策略"],
    ["场外期权二级交易商（2021）", "期权对冲+结构化产品", "跨资产组合风控+全包成本核算"],
    ["信用风险缓释工具交易商（2022）", "信用衍生品+债券借贷", "全资产相关性分析+多策略"],
    ["上市证券做市资格（2024）", "权益做市+流动性提供", "实时做市风控+限额管理"],
    ["碳排放权交易资格（2024）", "绿色金融+FICC延伸", "多资产统一组合管理"],
]
add_table_slide(slide, Inches(0.6), Inches(4.0), Inches(12.1), Emu(340000),
                headers, rows, font_size=10)

set_notes(slide, "国元的资质积累非常扎实——5年拿齐了FICC全套通行证。但这些能力分散在不同系统、不同团队里。POMS就是把这些资质的业务价值统一释放出来的平台。")


# ============================================================
# SLIDE 5 — Four Capability Leaps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "战略落地需要四大核心能力跃升",
               "配得优 · 算得快 · 控得稳 · 连得通")

# Four capability boxes
box_width = Inches(2.8)
box_height = Inches(3.8)
gap = Inches(0.3)
start_x = Inches(0.6)
start_y = Inches(1.6)

capabilities = [
    ("配得优", ACCENT_GREEN, [
        "千亿实现全资产最优配置",
        "全资产品种覆盖",
        "跨资产相关性分析",
        "数据驱动的试算工具",
        "自动再平衡",
        "IFRS9影响模拟"
    ]),
    ("算得快", LIGHT_BLUE, [
        "利差机会实时捕获",
        "实时风险计量",
        "自主定价能力",
        "利差/基差/新老券策略",
        "多策略回测验证",
        "精确全包成本核算"
    ]),
    ("控得稳", ACCENT_RED, [
        "回撤秒级预警干预",
        "极端行情压力测试",
        "流动性冲击可量化",
        "合规事前拦截",
        "实时限额管理",
        "现金流预测"
    ]),
    ("连得通", ACCENT_GOLD, [
        "投研到交易全链贯通",
        "指令自动传递",
        "multi-leg系统支持",
        "执行质量分析(TCA)",
        "实时绩效归因",
        "数据统一（对接IBOR）"
    ]),
]

for i, (title, color, items) in enumerate(capabilities):
    x = start_x + i * (box_width + gap)
    box = add_shape(slide, x, start_y, box_width, box_height,
                    fill_color=WHITE, line_color=color)

    # Title bar in box
    add_shape(slide, x, start_y, box_width, Inches(0.6), fill_color=color)
    add_textbox(slide, x, start_y + Inches(0.05), box_width, Inches(0.5),
                f"【{title}】", font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Items
    tb = add_textbox(slide, x + Inches(0.15), start_y + Inches(0.7),
                     box_width - Inches(0.3), box_height - Inches(0.8), "", font_size=11)
    tf = tb.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(6)

# Bottom tagline
add_textbox(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(0.5),
            "不是哪里不行，是战略要求更高——任何单点工具都解决不了，需要统一平台系统性支撑",
            font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

set_notes(slide, "这四个跃升目标——配得优、算得快、控得稳、连得通——不是孤立的需求，是战略落地的系统性能力要求。任何单点工具都解决不了，需要一个统一平台系统性支撑。")


# ============================================================
# ACT 2 DIVIDER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_act_divider(slide, "Act 2: 为什么现在必须行动",
                 "924真实教训 + 投资经理的一天 = 紧迫感")


# ============================================================
# SLIDE 6 — 924 Case Study
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "924事件证明：没有压力测试的自营部门在裸泳",
               "2024年924央行干预债市——有工具和没工具结果完全不同")

headers = ["时间", "没有POMS", "有POMS"]
rows = [
    ["09:30", "消息出来看新闻才知道", "CEP引擎自动捕获异常波动"],
    ["09:31", "—", "压测自动触发，实时损益-1.2亿"],
    ["09:35", "—", "系统生成减仓方案+合规预检"],
    ["09:38", "—", "一键执行，滑点仅3bp"],
    ["09:45", "开始手工算影响，数据是昨天的", "—"],
    ["10:15", "大致算出可能亏1.5亿但不确定", "—"],
    ["11:30", "执行减仓时价格又跌了", "—"],
    ["收盘", "亏2亿", "亏0.8亿 → 差异1.2亿"],
]
add_table_slide(slide, Inches(0.6), Inches(1.5), Inches(12.1), Emu(340000),
                headers, rows, font_size=11)

# Liquidity lesson box
box = add_shape(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.0),
                fill_color=RGBColor(0xFD, 0xF2, 0xF2), line_color=ACCENT_RED)
tb = add_textbox(slide, Inches(0.9), Inches(5.1), Inches(11.5), Inches(1.8),
                 "924的真正教训——流动性冲击比价格下跌更致命", font_size=14, bold=True, color=ACCENT_RED)
add_para(tb.text_frame,
         "真正的危险不是收益率飙升，而是不知道哪些券能卖、冲击成本有多大。",
         font_size=12, color=DARK_GRAY)
add_para(tb.text_frame,
         "POMS差异化能力：实时流动性冲击成本计算——逐一评估前20大持仓变现天数和冲击成本，优先选择冲击最小的减仓路径。",
         font_size=12, color=DARK_GRAY)

set_notes(slide, "这不是假设，是2024年9月24日真实发生的行业事件。同样的市场冲击，有系统的8分钟内完成评估和减仓，没有系统的2小时后还在算。差距是1.2亿，一天之内。特别要强调：924当天真正致命的是流动性冲击——不是价格跌了多少，而是你想卖的时候发现卖不掉或者要承受巨大折价。POMS实时计算每只券的流动性冲击成本，减仓时自动选择冲击最小的路径。")


# ============================================================
# SLIDE 7 — A Day in the Life
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "张总的一天：同一个人、同一个市场，工具不同结果完全不同",
               "POMS不是替代投资经理，是让优秀的人拥有更好的工具")

headers = ["时间", "没有POMS", "有POMS"]
rows = [
    ["08:30", "3个系统看昨日持仓，数据分散", "一屏看全组合，建议增配5%黄金"],
    ["09:15", "等昨日VaR邮件，过时12小时", "实时VaR 1.2亿，利率债集中度预警"],
    ["10:00", "Excel粗算后因不确定而放弃", "虚拟组合试算：波动率降32%"],
    ["11:30", "利差93%分位，缺工具不敢做", "推送预警+回测：胜率78%"],
    ["14:00", "手写指令→签字→价格已变", "一键→合规预检→秒级下达"],
    ["15:30", "不知道今天操作的贡献", "实时归因：利差+8bp"],
    ["17:00", "年底才知道落后同业", "年化5.8%，波动3.1%，同业前30%"],
]
add_table_slide(slide, Inches(0.6), Inches(1.5), Inches(12.1), Emu(350000),
                headers, rows, font_size=11)

# Knowledge retention box
box = add_shape(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.2),
                fill_color=RGBColor(0xF0, 0xF8, 0xF0), line_color=ACCENT_GREEN)
tb = add_textbox(slide, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.0),
                 "知识固化效应：如果张总离职，他的Excel模型和利差判断经验也将消失。",
                 font_size=13, bold=True, color=DARK_GRAY)
add_para(tb.text_frame,
         "POMS将策略逻辑、风险参数、历史决策系统化留存——从'依赖个人'升级为'机构化能力'，这正是国元'能力一体化'的关键一步。",
         font_size=12, color=MED_GRAY)

set_notes(slide, "请各位想象张总管理80亿组合。左边没有POMS——看不到全景，不敢做试算。右边有POMS——同一个人同样的市场，有了工具能看到机会、算清风险、秒级执行。还有一点很重要——如果张总离职了怎么办？有了POMS，知识被固化在系统中。这就是从依赖个人升级为机构化能力。")


# ============================================================
# SLIDE 8 — Why Now
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "为什么必须现在行动",
               "四大驱动力同时发力——不是要不要做，是能不能承受不做的代价")

headers = ["考核驱动", "监管驱动", "竞争驱动", "信创窗口"]
rows = [
    ["缺工具→错过行情→考核被动", "十五五明确数字化+信创", "华泰/平安已有完整平台", "依赖Calypso/Bloomberg"],
    ["IFRS9 2026年全面执行", "2024年券商罚单162次", "差距从可追赶变代差", "国产替代窗口期2-3年"],
    ["自营已占43%收入", "合规预检从建议变必须", "优秀PM被挖走", "对标Aladdin持续演进"],
]
add_table_slide(slide, Inches(0.6), Inches(1.5), Inches(12.1), Emu(450000),
                headers, rows, font_size=11)

# Benchmark box
tb = add_textbox(slide, Inches(0.6), Inches(3.7), Inches(12), Inches(0.4),
                 "标杆验证：统一组合管理平台是终局形态", font_size=16, bold=True, color=DARK_BLUE)

headers2 = ["贝莱德Aladdin", "华泰'大象'"]
rows2 = [
    ["管理11.6万亿美元AUM", "2020启动，3年建成"],
    ["核心是组合管理+风控能力", "2023年已对外开放赋能"],
    ["国元POMS愿景 = Mini-Aladdin", "国元18个月追赶3-5年"],
]
add_table_slide(slide, Inches(0.6), Inches(4.2), Inches(12.1), Emu(380000),
                headers2, rows2, font_size=12)

set_notes(slide, "四个驱动力同时发力。特别是竞争驱动——华泰大象2020年启动花了3年已经建成。一旦竞争对手建起完整平台而我们还没动，差距就不是慢一步，而是能力代差。")


# ============================================================
# ACT 3 DIVIDER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_act_divider(slide, "Act 3: POMS如何支撑四大能力跃升",
                 "架构 → 能力 → 成本 → 生态，逐层展开")


# ============================================================
# SLIDE 9 — POMS Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "华锐POMS平台：六大引擎构建投资决策的智慧大脑",
               "全资产组合管理 + 量化策略 + 实时风控 + 指令管理 + 成本核算 + 合规预检")

# Architecture diagram as labeled boxes
# Top: Investment Workstation
add_shape(slide, Inches(3.5), Inches(1.5), Inches(6.3), Inches(0.7),
          fill_color=LIGHT_BLUE, line_color=MEDIUM_BLUE)
add_textbox(slide, Inches(3.5), Inches(1.55), Inches(6.3), Inches(0.6),
            "▲ 投资工作台：PM台 | 交易台 | 风控台 | 策略台",
            font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Left: Research
add_shape(slide, Inches(0.5), Inches(2.8), Inches(2.5), Inches(2.5),
          fill_color=LIGHT_GRAY, line_color=MED_GRAY)
add_textbox(slide, Inches(0.5), Inches(2.9), Inches(2.5), Inches(2.3),
            "◄ 投研系统\n\n策略信号\n宏观因子\nAI模型",
            font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Center: POMS
add_shape(slide, Inches(3.5), Inches(2.5), Inches(6.3), Inches(3.2),
          fill_color=RGBColor(0xE8, 0xF0, 0xF8), line_color=DARK_BLUE)
add_textbox(slide, Inches(3.5), Inches(2.55), Inches(6.3), Inches(0.4),
            "华 锐  P O M S  平 台",
            font_size=16, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

engines = [
    "① 全资产组合管理引擎",
    "② 量化策略引擎（相对价值/回测）",
    "③ 实时风控 + 极端压力测试",
    "④ 指令与执行管理",
    "⑤ 绩效归因 + 精确成本核算",
    "⑥ 合规预检 + 实时限额管理",
]
for i, eng in enumerate(engines):
    y = Inches(3.1) + i * Inches(0.38)
    add_textbox(slide, Inches(4.2), y, Inches(5), Inches(0.35),
                eng, font_size=12, color=DARK_BLUE)

# Right: OEMS
add_shape(slide, Inches(10.3), Inches(2.8), Inches(2.5), Inches(2.5),
          fill_color=LIGHT_GRAY, line_color=MED_GRAY)
add_textbox(slide, Inches(10.3), Inches(2.9), Inches(2.5), Inches(2.3),
            "OEMS系统 ►\n\n交易执行\nmulti-leg\n算法交易",
            font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom: Data Infrastructure
add_shape(slide, Inches(3.5), Inches(6.0), Inches(6.3), Inches(0.7),
          fill_color=MED_GRAY)
add_textbox(slide, Inches(3.5), Inches(6.05), Inches(6.3), Inches(0.6),
            "▼ 已建数据基础设施：金证IBOR | 金仕达行情 | 数据总线",
            font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

set_notes(slide, "这是POMS的整体架构。核心是六大引擎。向下对接国元已经在建的金证IBOR和金仕达行情中心，向上支撑投资工作台，向左连投研，向右连交易执行。接下来看每个能力目标如何实现。")


# ============================================================
# SLIDE 10 — 配得优 + 算得快
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "配得优 + 算得快：让千亿找到最优配置，让每个bp都算清楚",
               "全资产组合管理 + 量化策略 + 精确成本核算 = 多赚的基础")

headers = ["能力需求", "POMS如何支撑", "关键能力"]
rows = [
    ["全资产配置", "全资产统一组合管理", "股/债/商品/衍生品/基金/黄金"],
    ["跨资产分析", "相关性矩阵+风险分散度", "找到低波动低相关组合"],
    ["数据驱动决策", "what-if虚拟组合试算", "100+并行模拟，秒级出结果"],
    ["实时风险计量", "分布式实时计算", "VaR/Greeks秒级输出"],
    ["自主定价", "全品种定价引擎", "利率曲线/波动率曲面"],
    ["全包成本", "精确成本核算引擎", "融资+佣金+CFETS+结算=真实P&L"],
]
add_table_slide(slide, Inches(0.6), Inches(1.5), Inches(12.1), Emu(360000),
                headers, rows, font_size=11)

# Cost example box
box = add_shape(slide, Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.2),
                fill_color=RGBColor(0xFD, 0xF8, 0xF0), line_color=ACCENT_GOLD)
tb = add_textbox(slide, Inches(0.9), Inches(4.7), Inches(11.5), Inches(2.0),
                 "成本核算示例：1亿杠杆交易的'隐藏真相'", font_size=14, bold=True, color=DARK_GRAY)
add_para(tb.text_frame, "看似收益：卖出价差 = +40万", font_size=12, color=ACCENT_GREEN)
add_para(tb.text_frame, "全包成本：融资44.4万 + 佣金1万 + CFETS 0.1万 + 结算0.05万 = 45.6万", font_size=12, color=ACCENT_RED)
add_para(tb.text_frame, "真实P&L：-5.6万 → 看似赚了40万，实际亏了5.6万！", font_size=13, bold=True, color=ACCENT_RED)

set_notes(slide, "配得优和算得快是多赚的基础。千亿如果只投利率债，当股票涨30%、黄金涨25%的时候只能看着。同时精确成本核算引擎算清每笔交易的真实P&L——500亿杠杆规模，隐性亏损可能比我们想象的多。")


# ============================================================
# SLIDE 11 — 控得稳
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "控得稳：实时风控 + 924级压力测试 + 合规预检",
               "从T+1人工到秒级自动预警，从无预案到三层压力测试")

headers = ["能力需求", "POMS如何支撑", "关键能力"]
rows = [
    ["实时回撤预警", "CEP实时回撤监控", "秒级检测+阈值自动预警"],
    ["极端行情预案", "三层压力测试引擎", "历史重演/假设模拟/反向压测"],
    ["流动性管理", "流动性压力测试", "持仓变现天数/冲击成本"],
    ["合规事前拦截", "投前合规预检", "事前拦截+事中阻止"],
    ["限额实时管理", "实时限额仪表盘", "下单前即知限额影响"],
    ["现金流预测", "T+N前瞻预测", "资金缺口提前预警"],
]
add_table_slide(slide, Inches(0.6), Inches(1.5), Inches(12.1), Emu(360000),
                headers, rows, font_size=11)

# Three-layer stress test
tb = add_textbox(slide, Inches(0.6), Inches(4.6), Inches(12), Inches(0.4),
                 "三层压测能力", font_size=16, bold=True, color=DARK_BLUE)

headers2 = ["层级", "方法", "典型场景"]
rows2 = [
    ["第1层", "历史重演", "924央行干预 | 包商银行 | 2013钱荒"],
    ["第2层", "假设模拟", "利率+100bp | 利差走阔50bp | 股市-10%"],
    ["第3层", "反向压测", "'我最多亏X亿'— 什么情景会导致？"],
]
add_table_slide(slide, Inches(0.6), Inches(5.1), Inches(12.1), Emu(380000),
                headers2, rows2, font_size=12)

set_notes(slide, "924事件的教训告诉我们常规风控不够。POMS提供三层压力测试。同时合规预检在下单前自动检查——2024年券商罚单162次，事前拦截是最有效的合规手段。")


# ============================================================
# SLIDE 12 — 连得通
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "连得通：从发现机会到交易完成，全链路贯通",
               "投研→决策→指令→执行→归因，全链路贯通")

headers = ["能力需求", "POMS如何支撑", "关键能力"]
rows = [
    ["投研到组合", "策略信号直接触发", "因子库/利差信号自动驱动"],
    ["决策到指令", "自动生成交易指令", "合规预检自动通过"],
    ["multi-leg执行", "多腿交易引擎", "利差交易同时买卖多只券"],
    ["执行质量", "执行质量分析+TCA", "滑点归因/TWAP/VWAP"],
    ["绩效评估", "实时绩效归因", "Campisi模型+成本归因"],
]
add_table_slide(slide, Inches(0.6), Inches(1.5), Inches(12.1), Emu(400000),
                headers, rows, font_size=12)

# Flow diagram
tb = add_textbox(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(1.5),
                 "", font_size=14)
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "当前链路（1-2天）：投研发现 → 传给PM → 手写指令 → 找领导签字 → 交易员执行 → 价格已变"
p.font.size = Pt(13)
p.font.color.rgb = ACCENT_RED
p.font.name = "Microsoft YaHei"

p2 = tf.add_paragraph()
p2.text = "POMS链路（分钟级）：策略信号 → 组合调整 → 自动生成指令 → 合规预检 → 审批 → 秒级下达"
p2.font.size = Pt(13)
p2.font.color.rgb = ACCENT_GREEN
p2.font.name = "Microsoft YaHei"
p2.space_before = Pt(12)

set_notes(slide, "连得通的价值是真金白银。当前链路1-2天，POMS让这个链路变成分钟级。")


# ============================================================
# SLIDE 13 — IFRS9 + Cost
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "IFRS9 + 全包成本：投资决策前就看清会计影响和真实成本",
               "2026年IFRS9全面执行——没有系统支撑利润表波动将不可控")

tb = add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(3), "", font_size=13)
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "IFRS9挑战"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.font.name = "Microsoft YaHei"

items = [
    "金融资产分FVTPL/FVOCI/AC三类",
    "分类直接影响利润表波动",
    "FVTPL：公允价值进P&L，波动大",
    "FVOCI：进OCI，利润稳但无法兑现价差",
    "投资前必须知道：放哪个账户，影响多大",
]
for item in items:
    p = tf.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Microsoft YaHei"
    p.space_before = Pt(4)

tb2 = add_textbox(slide, Inches(6.5), Inches(1.5), Inches(6.2), Inches(3), "", font_size=13)
tf2 = tb2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "全包成本管理"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.font.name = "Microsoft YaHei"

items2 = [
    "500亿杠杆，回购利率优化10bp = 2000-3000万/年",
    "佣金+CFETS+结算费年度近4000万，可优化",
    "每笔下单前自动显示：",
    "  全包成本 + IFRS9账户建议 + 真实净收益",
]
for item in items2:
    p = tf2.add_paragraph()
    p.text = f"• {item}' if not item.startswith('  ') else f'  {item.strip()}"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Microsoft YaHei"
    p.space_before = Pt(4)

set_notes(slide, "2026年新会计准则全面执行。一笔交易放FVTPL还是FVOCI直接决定利润表波动。POMS在每次下单前自动显示全包成本和会计影响。")


# ============================================================
# SLIDE 14 — Ecosystem Integration
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "POMS如何融入国元已建技术生态",
               "不是推倒重来——IBOR是发动机，POMS是自动驾驶系统")

# Car analogy box
box = add_shape(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(3.5),
                fill_color=RGBColor(0xF0, 0xF4, 0xF8), line_color=MEDIUM_BLUE)

tb = add_textbox(slide, Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.4),
                 "如果把国元的技术生态比作一辆智能汽车：", font_size=15, bold=True, color=DARK_BLUE)

analogy_items = [
    ("金证IBOR", "= 发动机（数据驱动力）"),
    ("金仕达行情中心", "= 仪表传感器（实时信息采集）"),
    ("金仕达数据总线", "= 车内总线（信息高速公路）"),
    ("恒生/衡泰/蜂虎", "= 各功能零部件（估值/交易/信评）"),
    ("华锐POMS", "= 自动驾驶系统（将所有数据转化为投资决策和行动）"),
]
for i, (name, desc) in enumerate(analogy_items):
    y = Inches(2.2) + i * Inches(0.45)
    color = ACCENT_GOLD if i == 4 else DARK_GRAY
    bold = i == 4
    size = 14 if i == 4 else 13
    add_textbox(slide, Inches(2), y, Inches(9), Inches(0.4),
                f"{name}  {desc}", font_size=size, bold=bold, color=color)

# Principles
tb = add_textbox(slide, Inches(0.6), Inches(5.3), Inches(12), Inches(0.4),
                 "设计原则", font_size=16, bold=True, color=DARK_BLUE)

headers = ["原则", "含义"]
rows = [
    ["协同而非替代", "消费各系统产出的数据，在其之上做组合优化、策略分析、风控预警"],
    ["增值而非重复", "不重建任何已有能力，专注于业务智能层"],
    ["渐进式上线", "Phase 1对接IBOR+行情中心，Phase 2扩展OEMS/投研"],
]
add_table_slide(slide, Inches(0.6), Inches(5.8), Inches(12.1), Emu(330000),
                headers, rows, font_size=11)

set_notes(slide, "这页是整个方案的关键。打一个比方——如果把国元的技术生态比作一辆智能汽车，金证IBOR是发动机，金仕达是仪表和传感器，恒生衡泰蜂虎是各个功能零部件。现在这辆车什么都有，唯独缺一个自动驾驶系统——把所有数据转化为投资决策和行动。华锐POMS就是这个自动驾驶系统。我们不替代任何已有零部件，我们让它们协同起来产生投资价值。")


# ============================================================
# ACT 4 DIVIDER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_act_divider(slide, "Act 4: 为什么选华锐",
                 "唯一能一体化交付的国产厂商")


# ============================================================
# SLIDE 15 — Why Huarui
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "华锐是唯一能一体化支撑四大能力跃升的国产厂商",
               "六大引擎全覆盖 + 毕马威金融科技50强六连冠 + 128件专利")

headers = ["#", "差异化优势", "对千亿自营的价值"]
rows = [
    ["1", "不只做一个点 — 六大引擎全覆盖", "一个厂商解决全部问题"],
    ["2", "不只是分析 — 覆盖交易执行、量化策略", "从策略到执行一条龙"],
    ["3", "不只是国内 — 香港国际业务中心", "跨境投资不用另找供应商"],
    ["4", "不只是替代 — 对标Aladdin/SecDB", "千亿规模值得世界级系统"],
]
add_table_slide(slide, Inches(0.6), Inches(1.5), Inches(12.1), Emu(400000),
                headers, rows, font_size=12)

# Credentials
tb = add_textbox(slide, Inches(0.6), Inches(3.8), Inches(12), Inches(2.5), "", font_size=13)
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "资质背书"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.font.name = "Microsoft YaHei"

creds = [
    "毕马威中国领先金融科技50强（2020-2025连续6年）",
    "国家级专精特新重点小巨人（2024年首批入选）",
    "IDC中国 FinTech50（2023-2025连续3年）",
    "9年 | 500人 | 10亿研发 | 128件专利 | 上海深圳+香港",
    "交付保障：核心团队来自头部券商和国际投行，Phase 1全程驻场",
]
for c in creds:
    p = tf.add_paragraph()
    p.text = f"• {c}"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Microsoft YaHei"
    p.space_before = Pt(6)

set_notes(slide, "市场上做组合管理的厂商不少但多数只覆盖一个点。华锐是唯一一个六大引擎全覆盖的。如CEO追问竞品：现有6家供应商各做垂直模块，华锐做的是横向的组合管理智能层，是互补不是竞争。")


# ============================================================
# SLIDE 16 — Roadmap
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "三期18个月，Phase 1六个月见效",
               "从底向上、分步见效——Phase 1六个月CEO即看全公司自营实时全貌")

# Three phase boxes
phases = [
    ("Phase 1: 筑基期\n0-6月", ACCENT_GREEN, [
        "★ 组合管理MVP",
        "★ 跨资产相关性",
        "   IFRS9基础",
        "   全包成本基础",
        "★ 基础风控",
        "★ 对接IBOR+行情",
    ], "CEO打开系统即看全貌"),
    ("Phase 2: 赋能期\n6-12月", LIGHT_BLUE, [
        "量化策略工具",
        "实时风控+压测",
        "指令管理",
        "绩效归因",
        "合规预检+限额",
        "信用风险",
    ], "实时风控上线\n指令全贯通"),
    ("Phase 3: 进化期\n12-18月", ACCENT_GOLD, [
        "策略回测框架",
        "multi-leg优化",
        "TCA持续优化",
        "AI引擎",
        "",
        "",
    ], "完整平台上线\n对标行业领先"),
]

for i, (title, color, items, milestone) in enumerate(phases):
    x = Inches(0.6) + i * Inches(4.2)
    # Phase box
    add_shape(slide, x, Inches(1.5), Inches(3.8), Inches(4.5),
              fill_color=WHITE, line_color=color)
    # Title
    add_shape(slide, x, Inches(1.5), Inches(3.8), Inches(0.8), fill_color=color)
    add_textbox(slide, x, Inches(1.55), Inches(3.8), Inches(0.7),
                title, font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Items
    tb = add_textbox(slide, x + Inches(0.15), Inches(2.4), Inches(3.5), Inches(2.2), "", font_size=11)
    tf = tb.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        if not item:
            continue
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(4)

    # Milestone
    add_shape(slide, x + Inches(0.2), Inches(4.7), Inches(3.4), Inches(0.7),
              fill_color=LIGHT_GRAY)
    add_textbox(slide, x + Inches(0.2), Inches(4.75), Inches(3.4), Inches(0.6),
                f"里程碑：{milestone}", font_size=10, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom commitment
add_textbox(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
            "底线承诺：Phase 1不达标则不启动Phase 2。每一期独立验收，国元完全掌握节奏。",
            font_size=13, bold=True, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)

set_notes(slide, "从底向上三期建设。Phase 1是关键——6个月内组合管理引擎上线，CEO打开系统就能看到全公司千亿自营的实时全貌。Phase 1不达标不启Phase 2，国元完全掌握节奏和主动权。")


# ============================================================
# ACT 5 DIVIDER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_act_divider(slide, "Act 5: 价值、风险、下一步",
                 "值多少 · 风险可控 · 不行动代价更大")


# ============================================================
# SLIDE 17 — ROI / Value
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "保守价值：年化直接价值1-1.8亿 + 极端事件避损1-2亿/次",
               "每个数字都是保守估计，且仅计算系统赋能部分")

# Four pillars
box = add_shape(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.8),
                fill_color=RGBColor(0xF0, 0xF4, 0xF8), line_color=MEDIUM_BLUE)
add_textbox(slide, Inches(0.9), Inches(1.45), Inches(11.5), Inches(0.7),
            "价值实现四大支柱：系统平台（POMS） + 人员能力 + 流程优化 + 管理机制  →  系统是必要条件，没有它其他三个支柱无从发力",
            font_size=12, color=DARK_BLUE)

# Value tables side by side
headers_l = ["多赚（Alpha增强）", "保守年化"]
rows_l = [
    ["配置优化（保守取行业1/10）", "2000-5000万"],
    ["量化策略", "1000-2000万"],
    ["再平衡效率", "2000-3000万"],
    ["小计", "0.5-1亿/年"],
]
add_table_slide(slide, Inches(0.6), Inches(2.5), Inches(5.8), Emu(350000),
                headers_l, rows_l, header_color=ACCENT_GREEN, font_size=11)

headers_r = ["少花（成本节省）", "保守年化"]
rows_r = [
    ["滑点优化", "1000-1500万"],
    ["融资优化", "2000-3000万"],
    ["TCA+成本核算+人效+信创+闲置", "~1500-3000万"],
    ["小计", "0.5-0.8亿/年"],
]
add_table_slide(slide, Inches(6.7), Inches(2.5), Inches(5.8), Emu(350000),
                headers_r, rows_r, header_color=MEDIUM_BLUE, font_size=11)

# Summary table
headers_s = ["指标", "数值"]
rows_s = [
    ["年化直接价值（多赚+少花）", "1-1.8亿/年"],
    ["极端事件保守避损", "1-2亿/次"],
    ["常规风控保守避损", "0.3-0.5亿/次"],
    ["5年累计直接价值", "5-9亿"],
]
add_table_slide(slide, Inches(2.5), Inches(4.8), Inches(8.3), Emu(330000),
                headers_s, rows_s, header_color=DARK_BLUE, font_size=12)

# Safety margin note
box = add_shape(slide, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.7),
                fill_color=RGBColor(0xFD, 0xF8, 0xF0), line_color=ACCENT_GOLD)
add_textbox(slide, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.6),
            "安全边际：配置Alpha假设0.02%-0.05%，仅为行业基准（0.3%-0.5%）的1/10 → 即使被挑战，还有9倍安全空间",
            font_size=12, bold=True, color=DARK_GRAY)

set_notes(slide, "我们的价值测算有两个特点：第一，全部是保守估计——配置优化用的是行业基准的十分之一。第二，只计算系统赋能的部分。我们用的配置Alpha假设是0.02%-0.05%，行业通常用0.3%-0.5%——我们只取了十分之一。即使被挑战说太乐观，我们还有9倍的安全边际。")


# ============================================================
# SLIDE 18 — Risk & Mitigation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "风险与保障：预见风险、提前化解",
               "每个风险都有对应的保障机制——最大的保障是Phase 1独立验收")

headers = ["风险", "保障措施"]
rows = [
    ["数据质量与对接", "Phase 1首月专项数据治理，逐字段校验"],
    ["业务适配度", "驻场需求调研+双周迭代，业务团队全程参与UAT"],
    ["系统集成复杂度", "标准API+适配层，Phase 1优先对接IBOR+行情中心"],
    ["用户接受度", "2-3位种子用户深度参与，先Excel替代再能力升级"],
    ["与现有供应商协调", "POMS是智能层，与金证/金仕达/恒生是互补关系"],
    ["项目周期", "三期解耦独立交付，Phase 1不达标不启Phase 2"],
]
add_table_slide(slide, Inches(0.6), Inches(1.5), Inches(12.1), Emu(420000),
                headers, rows, font_size=12)

set_notes(slide, "任何大型平台建设都有风险，我们不回避。特别要强调：POMS和国元现有的6家供应商是互补不是竞争。最重要的保障：三期解耦，Phase 1做完验收满意了再启Phase 2。")


# ============================================================
# SLIDE 19 — Cost of Inaction
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "不行动的代价远超行动的投入",
               "3年不行动保守代价3.3-6亿，还不含极端事件风险敞口")

headers = ["年份", "不行动的代价", "累计"]
rows = [
    ["Year 1", "配置机会（0.5-1亿）+ 隐性成本（0.2-0.5亿）+ 效率损失（0.15-0.3亿）", "~1-2亿"],
    ["Year 2", "累计Y1成本 + IFRS9波动 + 极端事件若发生（1-2亿）", "1.3-2.5亿+"],
    ["Year 3", "竞争对手平台上线→能力代差→PM被挖走→不可逆", "不可量化"],
]
add_table_slide(slide, Inches(0.6), Inches(1.5), Inches(12.1), Emu(500000),
                headers, rows, font_size=12)

# Summary box
box = add_shape(slide, Inches(2), Inches(3.8), Inches(9.3), Inches(2.5),
                fill_color=RGBColor(0xFD, 0xF2, 0xF2), line_color=ACCENT_RED)
tb = add_textbox(slide, Inches(2.3), Inches(3.9), Inches(8.7), Inches(2.3), "", font_size=14)
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "3年不行动累计代价"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT_RED
p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.CENTER

items = [
    "直接成本：保守 3.3-6亿",
    "+ 一次极端事件：1-2亿",
    "+ 人才流失+竞争力下滑 → 不可逆",
    "",
    "vs 平台投入 → 不行动的代价远超投入",
]
for item in items:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(4)

set_notes(slide, "换个角度看——不做的代价。Year 1约1-2亿。Year 3最关键——竞争对手建起平台后，优秀投资经理会去工具更好的平台，这是不可逆的。3年累计保守3.3-6亿，远超平台投入。投入金额在商务环节单独沟通。")


# ============================================================
# SLIDE 20 — Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
make_title_bar(slide, "建议的下一步",
               "建议立即启动Phase 1，6个月内见效")

steps = [
    ("1. 联合工作组组建（1周内）", "成立'业务+IT'联合工作组\n明确双方项目Owner、沟通机制、决策流程"),
    ("2. 需求确认研讨会（2-3周内）", "确认资产品种范围、与金证IBOR接口方案\n选定2-3位种子投资经理深度参与"),
    ("3. Phase 1详细方案（1个月内）", "技术架构详细设计（含6家现有供应商集成方案）\n数据对接+种子用户培训计划"),
    ("4. Phase 1启动 → 6个月交付", "CEO仪表盘上线，三大核心指标实时可见：\n① 全公司自营实时盈亏  ② 组合VaR  ③ 跨资产相关性热力图"),
]

for i, (title, desc) in enumerate(steps):
    y = Inches(1.6) + i * Inches(1.3)
    # Number circle
    colors = [ACCENT_GREEN, LIGHT_BLUE, ACCENT_GOLD, DARK_BLUE]
    add_shape(slide, Inches(0.6), y, Inches(0.05), Inches(1.0), fill_color=colors[i])

    add_textbox(slide, Inches(1.0), y, Inches(11), Inches(0.4),
                title, font_size=15, bold=True, color=DARK_BLUE)
    add_textbox(slide, Inches(1.0), y + Inches(0.4), Inches(11), Inches(0.8),
                desc, font_size=12, color=DARK_GRAY)

# Contact
add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
            "华锐技术  |  [联系人]  |  [电话]  |  [邮箱]",
            font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

set_notes(slide, "建议下一步分四步走。第一步最关键——一周内组建业务+IT联合工作组。第二步两周内做需求研讨会。第三步一个月内出详细方案。第四步启动建设——6个月后CEO打开手机就能看到三个数字：全公司自营实时P&L、组合VaR、跨资产相关性热力图。谢谢各位领导。")


# ============================================================
# SAVE
# ============================================================
output_path = "/mnt/d/Work/国元/CEO_Deck_v2.1.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
