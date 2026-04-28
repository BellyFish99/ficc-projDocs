"""
CEO Deck v2.2 — Built on Huarui's actual company PPT template.
Uses the M9 PPT as template source to inherit theme, backgrounds, logos, and fonts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# === Load company template ===
template_path = "MKT2026_AIM_M9.pptx"
# First, copy the template to preserve its theme
import shutil
shutil.copy("/mnt/d/Work/国元/MKT2026_AIM_华锐投资管理平台 M9介绍PPT(现券管理)_20260316.pptx", template_path)

prs = Presentation(template_path)

# Delete all existing slides (keep theme/master/layouts)
xml_slides = prs.slides._sldIdLst
for sldId in list(xml_slides):
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId:
        try:
            prs.part.drop_rel(rId)
        except (KeyError, Exception):
            pass
    xml_slides.remove(sldId)

# === Color constants from Huarui theme ===
DARK_NAVY = RGBColor(0x03, 0x0E, 0x42)      # Primary dark bg
MEDIUM_BLUE = RGBColor(0x15, 0x42, 0x8F)     # Secondary blue
BRAND_RED = RGBColor(0xF5, 0x4D, 0x61)       # Huarui accent red
BRAND_PINK = RGBColor(0xEC, 0x4C, 0x62)      # Text accent
BRAND_YELLOW = RGBColor(0xFF, 0xD0, 0x51)    # Highlight yellow
BRIGHT_YELLOW = RGBColor(0xFF, 0xFF, 0x00)   # Strong highlight
LIGHT_BLUE = RGBColor(0x37, 0x9D, 0xFF)      # Accent blue
BRIGHT_BLUE = RGBColor(0x32, 0x70, 0xFC)     # Links/accent
TEAL = RGBColor(0x56, 0xC4, 0xD0)            # Secondary accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x88, 0x88, 0x88)
TITLE_BLUE = RGBColor(0x28, 0x36, 0x63)      # Title text color

W = prs.slide_width   # 13.333"
H = prs.slide_height  # 7.5"

# === Layout references ===
LAYOUT_COVER = prs.slide_layouts[14]       # '封面'
LAYOUT_TOC = prs.slide_layouts[12]         # '目录'
LAYOUT_DARK = prs.slide_layouts[15]        # '3_内页深色渐变-Logo右上'
LAYOUT_LIGHT = prs.slide_layouts[11]       # '纯色内页（旧版兼容）-右上Logo'
LAYOUT_BLANK = prs.slide_layouts[16]       # '空白页'
LAYOUT_END = prs.slide_layouts[10]         # '末尾幻灯片'


def tb(slide, left, top, width, height, text="", size=14,
       bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    """Add textbox shorthand."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return txBox


def ap(text_frame, text, size=14, bold=False, color=DARK_GRAY,
       align=PP_ALIGN.LEFT, font="Microsoft YaHei", space=Pt(4)):
    """Add paragraph shorthand."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    p.space_before = space
    return p


def box(slide, left, top, width, height, fill=None, line=None):
    """Add rectangle shape."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def add_table(slide, left, top, width, rh, headers, rows, hdr_color=DARK_NAVY, fs=11):
    """Add formatted table."""
    nr = len(rows) + 1
    nc = len(headers)
    ts = slide.shapes.add_table(nr, nc, left, top, width, Emu(rh * nr))
    t = ts.table
    cw = width // nc
    for i in range(nc):
        t.columns[i].width = cw
    for i, h in enumerate(headers):
        c = t.cell(0, i)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = hdr_color
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(fs)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Microsoft YaHei"
    for r, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            c = t.cell(r + 1, c_idx)
            c.text = str(val)
            if r % 2 == 1:
                c.fill.solid()
                c.fill.fore_color.rgb = LIGHT_GRAY
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(fs)
                p.font.color.rgb = DARK_GRAY
                p.font.name = "Microsoft YaHei"
    return ts


def notes(slide, text):
    """Set speaker notes."""
    slide.notes_slide.notes_text_frame.text = text


def dark_slide(title_text, subtitle_text=""):
    """Create a dark-bg content slide with title."""
    slide = prs.slides.add_slide(LAYOUT_DARK)
    # Set title placeholder if available
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title_text
            for p in ph.text_frame.paragraphs:
                p.font.color.rgb = WHITE
                p.font.name = "Microsoft YaHei"
                p.font.bold = True
            break
    if subtitle_text:
        tb(slide, Inches(0.6), Inches(0.75), Inches(10), Inches(0.4),
           subtitle_text, size=13, color=BRAND_YELLOW, font="Microsoft YaHei")
    return slide


def light_slide(title_text, subtitle_text=""):
    """Create a light-bg content slide with title bar."""
    slide = prs.slides.add_slide(LAYOUT_LIGHT)
    # Title bar
    box(slide, Inches(0), Inches(0), W, Inches(1.1), fill=DARK_NAVY)
    tb(slide, Inches(0.6), Inches(0.12), Inches(11), Inches(0.5),
       title_text, size=22, bold=True, color=WHITE)
    if subtitle_text:
        tb(slide, Inches(0.6), Inches(0.6), Inches(11), Inches(0.35),
           subtitle_text, size=12, color=BRAND_YELLOW)
    return slide


def divider_slide(act_text, subtitle):
    """Create act divider on dark layout."""
    slide = prs.slides.add_slide(LAYOUT_DARK)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""
    # Red accent bar
    box(slide, Inches(5.5), Inches(2.8), Inches(2.3), Inches(0.05), fill=BRAND_RED)
    tb(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1),
       act_text, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Microsoft YaHei")
    tb(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.6),
       subtitle, size=16, color=BRAND_YELLOW, align=PP_ALIGN.CENTER, font="Microsoft YaHei")
    return slide


# ============================================================
# SLIDE 1 — COVER (using company cover layout)
# ============================================================
slide = prs.slides.add_slide(LAYOUT_COVER)
# Add our content on top of the company cover background
tb(slide, Inches(0.8), Inches(2.0), Inches(8), Inches(0.5),
   "/ / /", size=16, color=WHITE, font="Microsoft YaHei")
tb(slide, Inches(0.8), Inches(2.7), Inches(8), Inches(0.8),
   "POMS", size=48, bold=True, color=WHITE, font="Microsoft YaHei")
t = tb(slide, Inches(0.8), Inches(3.4), Inches(8), Inches(0.7),
       "Archforce Investment Management", size=18, color=WHITE, font="Microsoft YaHei")
# Red accent bar
box(slide, Inches(0.88), Inches(4.3), Inches(1.2), Inches(0.025), fill=BRAND_RED)
tb(slide, Inches(0.8), Inches(4.5), Inches(8), Inches(0.4),
   "2026.Q2", size=16, color=WHITE, font="Microsoft YaHei")
notes(slide, "封面页。开场前确认参会人员名单和预计时间（约55分钟+Q&A）。")


# ============================================================
# SLIDE 2 — Executive Summary
# ============================================================
slide = light_slide("执行摘要",
                    "一页看懂：千亿自营资金需要什么、我们提供什么、值多少")

headers = ["维度", "要点"]
rows = [
    ["国元现状", "自营已成第一大收入来源（2025H1占比43%），近千亿AUM，战略转型方向清晰"],
    ["战略需要", "四大核心能力跃升：配得优 · 算得快 · 控得稳 · 连得通"],
    ["解决方案", "华锐POMS平台：全资产组合管理 + 量化策略 + 实时风控 + 指令闭环 + 精确成本核算"],
    ["生态协同", "与金证IBOR、金仕达行情中心共生——IBOR是数据引擎，POMS是业务智能"],
    ["实施路径", "三期18个月，Phase 1六个月见效"],
    ["保守价值", "年化直接价值1-1.8亿 + 极端事件保守避损1-2亿/次（系统赋能部分）"],
    ["投资保障", "Phase 1灯塔效应——6个月内CEO可看到全公司自营实时全貌"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(380000),
          headers, rows, fs=10)

# Core judgment box
b = box(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.4),
        fill=RGBColor(0x05, 0x12, 0x4A), line=MEDIUM_BLUE)
t = tb(slide, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.35),
       "核心判断：POMS不是一个IT系统，而是投资能力的基础设施。",
       size=13, bold=True, color=WHITE)
ap(t.text_frame,
   "价值的充分释放需要系统+人员+流程+管理四个支柱协同——系统是必要条件，没有它其他三个支柱无从发力。",
   size=11, color=BRAND_YELLOW)
ap(t.text_frame,
   "POMS与国元已建技术生态是共生关系：金证IBOR是数据引擎，金仕达是信息高速公路，POMS是自动驾驶系统。",
   size=11, color=RGBColor(0xAA, 0xBB, 0xDD))

notes(slide, "各位领导，一页纸看懂我们的方案。国元近千亿自营资金，战略目标非常清晰。要实现全资产精细化管理，需要四大能力跃升。华锐POMS就是支撑这四个跃升的业务智能平台。")


# ============================================================
# ACT 1 DIVIDER
# ============================================================
divider_slide("Act 1: 我们深刻理解国元",
              "不是来卖系统的，是来解决战略问题的")


# ============================================================
# SLIDE 3 — 4x Growth
# ============================================================
slide = light_slide("国元自营业务：4年4倍增长，已成第一大支柱",
                    "自营收入4年增长4倍至占比43%，去方向性+FICC+资本中介三位一体战略已定")

headers = ["年份", "自营收入", "占比", "财富信用", "关键转折"]
rows = [
    ["2022", "4.48亿", "8%", "19.29亿", "战略转型元年：确立去方向性"],
    ["2023", "13.07亿", "21%", "15.05亿", "+192%爆发，FICC资本中介快速发展"],
    ["2024", "23.11亿", "29%", "16.96亿", "首次超越财富业务，成为第一大收入"],
    ["2025H1", "14.60亿", "43%", "9.98亿", "半年超2023全年，领先财富1.46倍"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(380000),
          headers, rows, fs=11)

tb(slide, Inches(0.6), Inches(3.6), Inches(12), Inches(0.35),
   "2025年8月组织架构升级——三重买方信号", size=15, bold=True, color=TITLE_BLUE)

headers2 = ["调整", "战略信号"]
rows2 = [
    ["自营业务委员会 → 金融市场业务委员会", "从投资本位到市场服务本位——需要统一组合管理平台"],
    ["固定收益部 → FICC业务总部", "从单一债券到综合FICC——需要跨资产组合管理"],
    ["权益投资部 → 证券投资部", "去方向性+多策略——需要实时风控+相对价值工具"],
    ["创新金融部 → 创新金融业务总部（升格）", "衍生品从边缘到核心——需要统一定价+成本核算"],
]
add_table(slide, Inches(0.6), Inches(4.1), Inches(12.1), Emu(330000),
          headers2, rows2, fs=10)

notes(slide, "我们对国元做了深入研究。自营业务4年增长4倍，2025年上半年占比43%。2025年8月的组织架构升级更说明问题——每一个调整都指向同一个方向：需要一个统一的组合管理平台。")


# ============================================================
# SLIDE 4 — Requirements
# ============================================================
slide = light_slide("我们理解国元的需求：资质已就位，缺的是统一平台",
                    "5年积累完整FICC资质拼图——资质是通行证，POMS是驾驶能力")

t = tb(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(1.8), "", size=13)
tf = t.text_frame
tf.word_wrap = True
for i, (title, desc) in enumerate([
    ("1. 以稳定收益为目标", "低波动6%，从方向性投资转向精细化管理"),
    ("2. 全资产品种覆盖", "不是固收单品种，而是股票、黄金、债券、衍生品真正的多资产"),
    ("3. 丰富投资策略", "相对价值（利差/基差/新老券）+ 多策略多timeframe"),
    ("4. 以组合管理为核心", "不是以交易为核心，这是Aladdin级的思维高度"),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"{title} -- {desc}"
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Microsoft YaHei"
    p.space_before = Pt(8)

headers = ["已获资质", "解锁的业务能力", "统一平台能做什么"]
rows = [
    ["利率互换交易资格（2020）", "利率衍生品对冲", "统一定价引擎支撑利差策略"],
    ["场外期权二级交易商（2021）", "期权对冲+结构化产品", "跨资产组合风控+全包成本核算"],
    ["信用风险缓释工具交易商（2022）", "信用衍生品+债券借贷", "全资产相关性分析+多策略"],
    ["上市证券做市资格（2024）", "权益做市+流动性提供", "实时做市风控+限额管理"],
    ["碳排放权交易资格（2024）", "绿色金融+FICC延伸", "多资产统一组合管理"],
]
add_table(slide, Inches(0.6), Inches(3.8), Inches(12.1), Emu(330000),
          headers, rows, fs=10)

notes(slide, "国元的资质积累非常扎实——5年拿齐了FICC全套通行证。但这些能力分散在不同系统、不同团队里。POMS就是把这些资质的业务价值统一释放出来的平台。")


# ============================================================
# SLIDE 5 — Four Capability Leaps
# ============================================================
slide = light_slide("战略落地需要四大核心能力跃升",
                    "配得优 · 算得快 · 控得稳 · 连得通")

# Four colored boxes
cap_data = [
    ("配得优", TEAL, ["全资产品种覆盖", "跨资产相关性分析", "数据驱动的试算工具", "自动再平衡", "IFRS9影响模拟"]),
    ("算得快", LIGHT_BLUE, ["实时风险计量", "自主定价能力", "利差/基差/新老券策略", "多策略回测验证", "精确全包成本核算"]),
    ("控得稳", BRAND_RED, ["秒级回撤预警", "极端压力测试", "流动性冲击量化", "投前合规预检", "实时限额管理"]),
    ("连得通", BRAND_YELLOW, ["投研到组合直连", "指令自动生成", "multi-leg执行", "TCA交易成本分析", "实时绩效归因"]),
]

for i, (title, color, items) in enumerate(cap_data):
    x = Inches(0.5) + i * Inches(3.15)
    y = Inches(1.5)
    bw = Inches(2.9)
    bh = Inches(4.5)
    box(slide, x, y, bw, bh, line=color)
    box(slide, x, y, bw, Inches(0.9), fill=color)
    tb(slide, x, y + Inches(0.05), bw, Inches(0.8),
       title, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    t = tb(slide, x + Inches(0.15), y + Inches(1.0), bw - Inches(0.3), bh - Inches(1.1), "", size=11)
    tf = t.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(8)

tb(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
   "不是哪里不行，是战略要求更高——任何单点工具都解决不了，需要统一平台系统性支撑",
   size=11, color=MED_GRAY, align=PP_ALIGN.CENTER)

notes(slide, "这四个跃升目标——配得优、算得快、控得稳、连得通——不是孤立的需求，是战略落地的系统性能力要求。任何单点工具都解决不了，需要一个统一平台系统性支撑。")


# ============================================================
# ACT 2 DIVIDER
# ============================================================
divider_slide("Act 2: 为什么现在必须行动",
              "924真实教训 + 投资经理的一天 = 紧迫感")


# ============================================================
# SLIDE 6 — 924 Case
# ============================================================
slide = light_slide("924事件证明：没有压力测试的自营部门在裸泳",
                    "2024年924央行干预债市——有工具和没工具结果完全不同")

headers = ["时间", "没有POMS", "有POMS"]
rows = [
    ["09:30", "消息出来看新闻才知道", "CEP引擎自动捕获异常波动"],
    ["09:31", "—", "压测自动触发，实时损益-1.2亿"],
    ["09:35", "—", "系统生成减仓方案+合规预检"],
    ["09:38", "—", "一键执行，滑点仅3bp"],
    ["09:45", "开始手工算影响，数据是昨天的", "—"],
    ["10:15", "大致算出可能亏1.5亿但不确定", "—"],
    ["11:30", "执行减仓时价格又跌了", "—"],
    ["Close", "亏2亿", "亏0.8亿"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(320000),
          headers, rows, fs=10)

# PROMINENT DELTA CALLOUT
tb(slide, Inches(2), Inches(4.3), Inches(9.3), Inches(0.6),
   "差异：1.2亿——一天之内", size=28, bold=True, color=BRAND_RED, align=PP_ALIGN.CENTER)

# Liquidity lesson
b = box(slide, Inches(0.6), Inches(5.1), Inches(12.1), Inches(1.8),
        fill=RGBColor(0x2A, 0x0A, 0x12), line=BRAND_RED)
t = tb(slide, Inches(0.9), Inches(4.9), Inches(11.5), Inches(0.35),
       "924的真正教训——流动性冲击比价格下跌更致命",
       size=14, bold=True, color=BRAND_RED)
ap(t.text_frame,
   "真正的危险不是收益率飙升，而是不知道哪些券能卖、冲击成本有多大。",
   size=12, color=WHITE)
ap(t.text_frame,
   "POMS差异化能力：实时流动性冲击成本计算——逐一评估前20大持仓变现天数和冲击成本，优先选择冲击最小的减仓路径。",
   size=12, color=BRAND_YELLOW)

notes(slide, "这不是假设，是2024年9月24日真实发生的行业事件。差距是1.2亿，一天之内。924当天真正致命的是流动性冲击——不是价格跌了多少，而是你想卖的时候发现卖不掉。POMS实时计算每只券的流动性冲击成本，减仓时自动选择冲击最小的路径。")


# ============================================================
# SLIDE 7 — Day in Life
# ============================================================
slide = light_slide("张总的一天：同一个人、同一个市场，工具不同结果完全不同",
                    "POMS不是替代投资经理，是让优秀的人拥有更好的工具")

headers = ["时间", "没有POMS", "有POMS"]
rows = [
    ["08:30", "3个系统看昨日持仓，数据分散", "一屏看全组合，建议增配5%黄金"],
    ["09:15", "等昨日VaR邮件，过时12小时", "实时VaR 1.2亿，利率债集中度预警"],
    ["10:00", "Excel粗算后因不确定而放弃", "虚拟组合试算：波动率降32%"],
    ["11:30", "利差93%分位，缺工具不敢做", "推送预警+回测：胜率78%"],
    ["14:00", "手写指令→签字→价格已变", "一键→合规预检→秒级下达"],
    ["15:30", "不知道今天操作的贡献", "实时归因：利差+8bp"],
    ["17:00", "年底才知道落后同业", "年化5.8%，波动3.1%，同业前30%"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(340000),
          headers, rows, fs=10)

# Knowledge retention
b = box(slide, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.2),
        fill=RGBColor(0x0A, 0x1A, 0x0A), line=TEAL)
t = tb(slide, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.35),
       "知识固化效应：如果张总离职，他的Excel模型和利差判断经验也将消失。",
       size=13, bold=True, color=TEAL)
ap(t.text_frame,
   "POMS将策略逻辑、风险参数、历史决策系统化留存——从依赖个人升级为机构化能力。",
   size=12, color=WHITE)

notes(slide, "请各位想象张总管理80亿组合。左边没有POMS——看不到全景，不敢做试算。右边有POMS——同一个人同样的市场，有了工具能看到机会、算清风险、秒级执行。还有一点很重要——如果张总离职了怎么办？有了POMS，知识被固化在系统中。这就是从依赖个人升级为机构化能力。")


# ============================================================
# SLIDE 8 — Why Now
# ============================================================
slide = light_slide("为什么必须现在行动",
                    "四大驱动力同时发力——不是要不要做，是能不能承受不做的代价")

headers = ["考核驱动", "监管驱动", "竞争驱动", "信创窗口"]
rows = [
    ["缺工具→错过行情→考核被动", "十五五明确数字化+信创", "华泰/平安已有完整平台", "依赖Calypso/Bloomberg"],
    ["IFRS9 2026年全面执行", "2024年券商罚单162次", "差距从可追赶变代差", "国产替代窗口期2-3年"],
    ["自营已占43%收入", "合规预检从建议变必须", "优秀PM被挖走", "对标Aladdin持续演进"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(430000),
          headers, rows, fs=10)

tb(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(0.35),
   "标杆验证：统一组合管理平台是终局形态", size=15, bold=True, color=TITLE_BLUE)

headers2 = ["贝莱德Aladdin", "华泰大象"]
rows2 = [
    ["管理11.6万亿美元AUM", "2020启动，3年建成"],
    ["核心是组合管理+风控能力", "2023年已对外开放赋能"],
    ["国元POMS愿景 = Mini-Aladdin", "国元18个月追赶3-5年"],
]
add_table(slide, Inches(0.6), Inches(4.0), Inches(12.1), Emu(360000),
          headers2, rows2, fs=11)

notes(slide, "四个驱动力同时发力。特别是竞争驱动——华泰大象2020年启动花了3年已经建成。一旦竞争对手建起完整平台而我们还没动，差距就不是慢一步，而是能力代差。")


# ============================================================
# ACT 3 DIVIDER
# ============================================================
divider_slide("Act 3: POMS如何支撑四大能力跃升",
              "架构 → 能力 → 成本 → 生态，逐层展开")


# ============================================================
# SLIDE 9 — Architecture
# ============================================================
slide = light_slide("华锐POMS平台：六大引擎构建投资决策的智慧大脑",
                    "全资产组合管理 + 量化策略 + 实时风控 + 指令管理 + 成本核算 + 合规预检")

# Architecture as labeled boxes
# Top layer
box(slide, Inches(3.5), Inches(1.4), Inches(6.3), Inches(0.65), fill=LIGHT_BLUE)
tb(slide, Inches(3.5), Inches(1.43), Inches(6.3), Inches(0.55),
   "投资工作台：PM台 | 交易台 | 风控台 | 策略台", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Left
box(slide, Inches(0.5), Inches(2.5), Inches(2.5), Inches(2.8), fill=RGBColor(0xE8, 0xEE, 0xF5))
tb(slide, Inches(0.5), Inches(2.6), Inches(2.5), Inches(2.6),
   "投研系统\n\n策略信号\n宏观因子\nAI模型", size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Center - POMS
box(slide, Inches(3.5), Inches(2.3), Inches(6.3), Inches(3.5), fill=RGBColor(0x05, 0x14, 0x4E))
tb(slide, Inches(3.5), Inches(2.35), Inches(6.3), Inches(0.35),
   "华 锐  P O M S  平 台", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

engines = [
    "① 全资产组合管理引擎",
    "② 量化策略引擎（相对价值/回测）",
    "③ 实时风控 + 极端压力测试",
    "④ 指令与执行管理",
    "⑤ 绩效归因 + 精确成本核算",
    "⑥ 合规预检 + 实时限额管理",
]
for i, eng in enumerate(engines):
    y = Inches(2.85) + i * Inches(0.42)
    tb(slide, Inches(4.2), y, Inches(5), Inches(0.35),
       eng, size=12, color=BRAND_YELLOW if i < 3 else WHITE)

# Right
box(slide, Inches(10.3), Inches(2.5), Inches(2.5), Inches(2.8), fill=RGBColor(0xE8, 0xEE, 0xF5))
tb(slide, Inches(10.3), Inches(2.6), Inches(2.5), Inches(2.6),
   "OEMS系统\n\n交易执行\nmulti-leg\n算法交易", size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Bottom
box(slide, Inches(3.5), Inches(6.1), Inches(6.3), Inches(0.65), fill=MED_GRAY)
tb(slide, Inches(3.5), Inches(6.13), Inches(6.3), Inches(0.55),
   "已建数据基础设施：金证IBOR | 金仕达行情 | 数据总线",
   size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

notes(slide, "这是POMS的整体架构。核心是六大引擎。向下对接国元已经在建的金证IBOR和金仕达行情中心，向上支撑投资工作台，向左连投研，向右连交易执行。")


# ============================================================
# SLIDES 10-12 — Capabilities (3 slides)
# ============================================================

# Slide 10: Optimal + Fast
slide = light_slide("配得优 + 算得快：让千亿找到最优配置，让每个bp都算清楚",
                    "全资产组合管理 + 量化策略 + 精确成本核算 = 多赚的基础")

headers = ["能力需求", "POMS如何支撑", "关键能力"]
rows = [
    ["全资产配置", "全资产统一组合管理", "股/债/商品/衍生品/基金/黄金"],
    ["跨资产分析", "相关性矩阵+风险分散度", "找到低波动低相关组合"],
    ["数据驱动决策", "what-if虚拟组合试算", "100+并行模拟，秒级出结果"],
    ["实时风险计量", "Distributed real-time computation", "VaR/Greeks output in seconds"],
    ["自主定价能力", "Full-spectrum pricing engine", "Yield curves / volatility surfaces"],
    ["精确全包成本核算", "Precise cost accounting engine", "Financing + commission + CFETS + settlement = true P&L"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(340000),
          headers, rows, fs=10)

# Cost example
b = box(slide, Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.5),
        fill=RGBColor(0x20, 0x15, 0x05), line=BRAND_YELLOW)
t = tb(slide, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.35),
       "成本核算示例：1亿杠杆交易的隐藏真相", size=14, bold=True, color=BRAND_YELLOW)
ap(t.text_frame, "看似收益：卖出价差 = +40万", size=12, color=TEAL)
ap(t.text_frame, "全包成本：融资44.4万 + 佣金1万 + CFETS 0.1万 + 结算0.05万 = 45.6万", size=12, color=BRAND_RED)
ap(t.text_frame, "真实P&L：-5.6万 → 看似赚了40万，实际亏了5.6万！", size=13, bold=True, color=BRAND_RED)

notes(slide, "配得优和算得快是多赚的基础。千亿如果只投利率债，当股票涨30%、黄金涨25%的时候只能看着。同时精确成本核算引擎算清每笔交易的真实P&L——500亿杠杆规模，隐性亏损可能比我们想象的多。")


# Slide 11: Stable Control
slide = light_slide("控得稳：实时风控 + 924级压力测试 + 合规预检",
                    "从T+1人工到秒级自动预警，从无预案到三层压力测试")

headers = ["能力需求", "POMS如何支撑", "关键能力"]
rows = [
    ["实时回撤预警", "CEP实时回撤监控", "秒级检测+阈值自动预警"],
    ["极端行情预案", "三层压力测试引擎", "历史重演/假设模拟/反向压测"],
    ["流动性管理", "流动性压力测试", "持仓变现天数/冲击成本"],
    ["投前合规预检", "Embedded compliance engine", "Pre-trade block + in-flight interception"],
    ["实时限额管理", "Real-time limit dashboard", "See limit impact before placing order"],
    ["现金流预测", "T+N前瞻预测", "资金缺口提前预警"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(340000),
          headers, rows, fs=10)

tb(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(0.35),
   "三层压测能力", size=15, bold=True, color=TITLE_BLUE)
headers2 = ["层级", "方法", "典型场景"]
rows2 = [
    ["第1层", "历史重演", "924央行干预 | 包商银行 | 2013钱荒"],
    ["第2层", "假设模拟", "利率+100bp | 利差走阔50bp | 股市-10%"],
    ["第3层", "反向压测", "我最多亏X亿 → 什么情景会导致？"],
]
add_table(slide, Inches(0.6), Inches(4.8), Inches(12.1), Emu(360000),
          headers2, rows2, fs=11)

notes(slide, "924事件的教训告诉我们常规风控不够。POMS提供三层压力测试。同时合规预检在下单前自动检查——2024年券商罚单162次，事前拦截是最有效的合规手段。")


# Slide 12: Seamless Flow
slide = light_slide("连得通：从发现机会到交易完成，全链路贯通",
                    "投研→决策→指令→执行→归因，全链路贯通")

headers = ["能力需求", "POMS如何支撑", "关键能力"]
rows = [
    ["投研到组合", "策略信号直接触发", "因子库/利差信号自动驱动"],
    ["Decision to order", "指令自动生成", "Compliance auto-cleared"],
    ["multi-leg执行", "Multi-leg trade engine", "Spread trade: buy + sell simultaneously"],
    ["执行质量", "执行质量分析+TCA", "滑点归因/TWAP/VWAP"],
    ["Performance evaluation", "实时绩效归因", "Campisi model + cost attribution"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(380000),
          headers, rows, fs=11)

# Before/After flow
b = box(slide, Inches(0.6), Inches(4.2), Inches(12.1), Inches(1.5), line=MEDIUM_BLUE)
t = tb(slide, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.5),
       "当前链路（1-2天）：投研发现 → 传给PM → 手写指令 → 找领导签字 → 交易员执行 → 价格已变",
       size=12, color=BRAND_RED)
ap(t.text_frame,
   "POMS链路（分钟级）：策略信号 → 组合调整 → 自动生成指令 → 合规预检 → 审批 → 秒级下达",
   size=12, color=TEAL)

notes(slide, "连得通的价值是真金白银。当前链路1-2天，POMS让这个链路变成分钟级。")


# ============================================================
# SLIDE 13 — IFRS9
# ============================================================
slide = light_slide("IFRS9 + 全包成本：投资决策前就看清会计影响和真实成本",
                    "2026年IFRS9全面执行——没有系统支撑利润表波动将不可控")

t = tb(slide, Inches(0.6), Inches(1.4), Inches(5.5), Inches(3.5), "", size=13)
tf = t.text_frame
tf.word_wrap = True
ap(tf, "IFRS9挑战", size=16, bold=True, color=TITLE_BLUE, space=Pt(0))
for item in [
    "金融资产分FVTPL/FVOCI/AC三类",
    "分类直接影响利润表波动",
    "FVTPL：公允价值变动进P&L，波动大",
    "FVOCI：进OCI，利润稳但无法兑现价差",
    "投资前必须知道：放哪个账户，影响多大",
]:
    ap(tf, f"  {item}", size=11, color=DARK_GRAY)

t2 = tb(slide, Inches(6.5), Inches(1.4), Inches(6), Inches(3.5), "", size=13)
tf2 = t2.text_frame
tf2.word_wrap = True
ap(tf2, "全包成本管理", size=16, bold=True, color=TITLE_BLUE, space=Pt(0))
for item in [
    "500亿杠杆，回购利率优化10bp = 2000-3000万/年",
    "佣金+CFETS+结算费年度近4000万，可见可优化",
    "每笔下单前自动显示：",
    "   全包成本 + IFRS9账户建议 + 真实净收益",
]:
    ap(tf2, f"  {item}", size=11, color=DARK_GRAY)

notes(slide, "2026年新会计准则全面执行。一笔交易放FVTPL还是FVOCI直接决定利润表波动。POMS在每次下单前自动显示全包成本和会计影响。")


# ============================================================
# SLIDE 14 — Ecosystem
# ============================================================
slide = light_slide("POMS如何融入国元已建技术生态",
                    "不是推倒重来——IBOR是发动机，POMS是自动驾驶系统")

# Car analogy
b = box(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.0),
        fill=RGBColor(0x05, 0x12, 0x4A))
t = tb(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.35),
       "如果把国元的技术生态比作一辆智能汽车：", size=14, bold=True, color=WHITE)

analogies = [
    ("金证IBOR", "= 发动机（数据驱动力）", WHITE),
    ("金仕达行情中心", "= 仪表传感器（实时信息采集）", WHITE),
    ("金仕达数据总线", "= 车内总线（信息高速公路）", WHITE),
    ("恒生/衡泰/蜂虎", "= 各功能零部件（估值/交易/信评）", WHITE),
    ("华锐POMS", "= 自动驾驶系统（将所有数据转化为投资决策和行动）", BRAND_YELLOW),
]
for i, (name, desc, color) in enumerate(analogies):
    y = Inches(2.0) + i * Inches(0.42)
    bold = i == 4
    sz = 13 if i == 4 else 12
    tb(slide, Inches(2), y, Inches(9.5), Inches(0.35),
       f"{name}  {desc}", size=sz, bold=bold, color=color)

# Design principles
tb(slide, Inches(0.6), Inches(4.7), Inches(12), Inches(0.35),
   "设计原则", size=15, bold=True, color=TITLE_BLUE)

headers = ["原则", "含义"]
rows = [
    ["协同而非替代", "消费各系统产出的数据，在其之上做组合优化、策略分析、风控预警"],
    ["增值而非重复", "不重建任何已有能力，专注于业务智能层"],
    ["渐进式上线", "Phase 1对接IBOR+行情中心，Phase 2扩展OEMS/投研"],
]
add_table(slide, Inches(0.6), Inches(5.2), Inches(12.1), Emu(340000),
          headers, rows, fs=11)

notes(slide, "这页是整个方案的关键。打一个比方——如果把国元的技术生态比作一辆智能汽车，金证IBOR是发动机，金仕达是仪表和传感器，恒生衡泰蜂虎是各个功能零部件。华锐POMS就是自动驾驶系统——把所有数据转化为投资决策和行动。我们不替代任何已有零部件，我们让它们协同起来产生投资价值。")


# ============================================================
# ACT 4 DIVIDER
# ============================================================
divider_slide("Act 4: 为什么选华锐",
              "唯一能一体化交付的国产厂商")


# ============================================================
# SLIDE 15 — Why Huarui
# ============================================================
slide = light_slide("华锐是唯一能一体化支撑四大能力跃升的国产厂商",
                    "六大引擎全覆盖 + 毕马威金融科技50强六连冠 + 128件专利")

headers = ["#", "差异化优势", "对千亿自营的价值"]
rows = [
    ["1", "不只做一个点——六大引擎全覆盖", "一个厂商解决全部问题"],
    ["2", "不只是分析——覆盖交易执行、量化策略", "从策略到执行一条龙"],
    ["3", "不只是国内——香港国际业务中心", "跨境投资不用另找供应商"],
    ["4", "不只是替代——对标Aladdin/SecDB", "千亿规模值得世界级系统"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(380000),
          headers, rows, fs=11)

t = tb(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(2.5), "", size=13)
tf = t.text_frame
tf.word_wrap = True
ap(tf, "资质背书", size=16, bold=True, color=TITLE_BLUE, space=Pt(0))
for c in [
    "毕马威中国领先金融科技50强（2020-2025连续6年）",
    "国家级专精特新重点小巨人（2024年首批入选）",
    "IDC中国 FinTech50（2023-2025连续3年）",
    "9年 | 500人 | 10亿研发 | 128件专利 | 上海深圳+香港",
    "交付保障：核心团队来自头部券商和国际投行，Phase 1全程驻场",
]:
    ap(tf, f"  {c}", size=11, color=DARK_GRAY)

notes(slide, "市场上做组合管理的厂商不少但多数只覆盖一个点。华锐是唯一一个六大引擎全覆盖的。如CEO追问竞品：现有6家供应商各做垂直模块，华锐做的是横向的组合管理智能层，是互补不是竞争。")


# ============================================================
# SLIDE 16 — Roadmap
# ============================================================
slide = light_slide("三期18个月，Phase 1六个月见效",
                    "从底向上、分步见效——Phase 1六个月CEO即看全公司自营实时全貌")

# Three phase boxes
phases_data = [
    ("Phase 1\n0-6月 筑基期", TEAL,
     ["★ 组合管理MVP", "★ 跨资产相关性", "  IFRS9基础", "  全包成本基础", "★ 基础风控", "★ 对接IBOR+行情"],
     "CEO打开系统\n即看全貌"),
    ("Phase 2\n6-12月 赋能期", LIGHT_BLUE,
     ["量化策略工具", "实时风控+压测", "指令管理", "绩效归因", "合规预检+限额", "信用风险"],
     "实时风控上线\n指令全贯通"),
    ("Phase 3\n12-18月 进化期", BRAND_YELLOW,
     ["策略回测框架", "multi-leg优化", "TCA持续优化", "AI引擎", "", ""],
     "完整平台上线\n对标行业领先"),
]

for i, (title, color, items, milestone) in enumerate(phases_data):
    x = Inches(0.5) + i * Inches(4.2)
    box(slide, x, Inches(1.4), Inches(3.8), Inches(4.8), line=color)
    box(slide, x, Inches(1.4), Inches(3.8), Inches(1.0), fill=color)
    tb(slide, x, Inches(1.45), Inches(3.8), Inches(0.9),
       title, size=12, bold=True, color=WHITE if color != BRAND_YELLOW else DARK_NAVY, align=PP_ALIGN.CENTER)

    t = tb(slide, x + Inches(0.15), Inches(2.5), Inches(3.5), Inches(2.5), "", size=11)
    tf = t.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        if not item: continue
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(4)

    box(slide, x + Inches(0.15), Inches(5.0), Inches(3.5), Inches(0.6), fill=LIGHT_GRAY)
    tb(slide, x + Inches(0.15), Inches(5.05), Inches(3.5), Inches(0.5),
       f"Milestone: {milestone}", size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Bottom commitment
tb(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
   "底线承诺：Phase 1不达标则不启动Phase 2。每一期独立验收，国元完全掌握节奏。",
   size=12, bold=True, color=BRAND_RED, align=PP_ALIGN.CENTER)

notes(slide, "从底向上三期建设。Phase 1是关键——6个月内组合管理引擎上线，CEO打开系统就能看到全公司千亿自营的实时全貌。Phase 1不达标不启Phase 2，国元完全掌握节奏和主动权。")


# ============================================================
# ACT 5 DIVIDER
# ============================================================
divider_slide("Act 5: 价值、风险、下一步",
              "值多少 · 风险可控 · 不行动代价更大")


# ============================================================
# SLIDE 17 — ROI
# ============================================================
slide = light_slide("保守价值：年化直接价值1-1.8亿 + 极端事件避损1-2亿/次",
                    "每个数字都是保守估计，仅计算系统赋能部分——9倍安全边际")

# Four pillars bar
b = box(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.7),
        fill=RGBColor(0x05, 0x12, 0x4A))
tb(slide, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.6),
   "四大支柱：系统平台(POMS) + 人员能力 + 流程优化 + 管理机制 → 系统是必要条件，没有它其他支柱无从发力",
   size=11, color=BRAND_YELLOW)

# Two value tables
headers_l = ["多赚（Alpha增强）", "保守年化"]
rows_l = [
    ["配置优化（保守取行业基准1/10）", "2000-5000万"],
    ["量化策略", "1000-2000万"],
    ["再平衡效率", "2000-3000万"],
    ["小计", "0.5-1亿/年"],
]
add_table(slide, Inches(0.6), Inches(2.3), Inches(5.8), Emu(320000),
          headers_l, rows_l, hdr_color=TEAL, fs=10)

headers_r = ["少花（成本节省）", "保守年化"]
rows_r = [
    ["滑点优化", "1000-1500万"],
    ["融资优化", "2000-3000万"],
    ["TCA+成本核算+人效+信创+闲置", "~1500-3000万"],
    ["小计", "0.5-0.8亿/年"],
]
add_table(slide, Inches(6.7), Inches(2.3), Inches(5.8), Emu(320000),
          headers_r, rows_r, hdr_color=MEDIUM_BLUE, fs=10)

# Summary
headers_s = ["指标", "数值"]
rows_s = [
    ["年化直接价值（多赚+少花）", "1-1.8亿/年"],
    ["极端事件保守避损", "1-2亿/次"],
    ["常规风控保守避损", "0.3-0.5亿/次"],
    ["5年累计直接价值", "5-9亿"],
]
add_table(slide, Inches(2.5), Inches(4.5), Inches(8.3), Emu(320000),
          headers_s, rows_s, hdr_color=DARK_NAVY, fs=11)

# Safety margin
b = box(slide, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.65),
        fill=RGBColor(0x20, 0x15, 0x05), line=BRAND_YELLOW)
tb(slide, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.55),
   "安全边际：配置Alpha假设0.02%-0.05%，仅为行业基准（0.3%-0.5%）的1/10 → 即使被挑战，还有9倍安全空间\n注：行业基准通常为0.3%-0.5%，我们仅取约1/10以确保财务可靠性",
   size=12, bold=True, color=BRAND_YELLOW)

notes(slide, "我们的价值测算有两个特点：第一，全部是保守估计——配置优化用的是行业基准的十分之一。第二，只计算系统赋能的部分。我们用的配置Alpha假设是0.02%-0.05%，行业通常用0.3%-0.5%——我们只取了十分之一。即使被挑战说太乐观，我们还有9倍的安全边际。")


# ============================================================
# SLIDE 18 — Risk
# ============================================================
slide = light_slide("风险与保障：预见风险、提前化解",
                    "每个风险都有对应的保障机制——最大的保障是Phase 1独立验收")

headers = ["风险", "保障措施"]
rows = [
    ["数据质量与对接", "Phase 1首月专项数据治理，逐字段校验"],
    ["业务适配度", "驻场需求调研+双周迭代，业务团队全程参与UAT"],
    ["系统集成复杂度", "标准API+适配层，Phase 1优先对接IBOR+行情中心"],
    ["用户接受度", "2-3位种子PM深度参与设计，先Excel替代再能力升级"],
    ["与现有供应商协调", "POMS是智能层，与金证/金仕达/恒生是互补关系"],
    ["项目周期", "三期解耦独立交付，Phase 1不达标不启Phase 2"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(400000),
          headers, rows, fs=11)

notes(slide, "任何大型平台建设都有风险，我们不回避。特别要强调：POMS和国元现有的6家供应商是互补不是竞争。最重要的保障：三期解耦，Phase 1做完验收满意了再启Phase 2。")


# ============================================================
# SLIDE 19 — Cost of Inaction
# ============================================================
slide = light_slide("不行动的代价远超行动的投入",
                    "3年不行动保守代价3.3-6亿，还不含极端事件风险敞口")

headers = ["年份", "不行动的代价", "累计"]
rows = [
    ["Year 1", "配置机会（0.5-1亿）+ 隐性成本（0.2-0.5亿）+ 效率损失（0.15-0.3亿）", "~1-2亿"],
    ["Year 2", "累计Y1成本 + IFRS9波动 + 极端事件若发生（1-2亿）", "1.3-2.5亿+"],
    ["Year 3", "竞争对手平台上线→能力代差→PM被挖走→不可逆", "不可量化"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(450000),
          headers, rows, fs=11)

# Summary box
b = box(slide, Inches(2), Inches(3.7), Inches(9.3), Inches(2.8),
        fill=RGBColor(0x2A, 0x0A, 0x0A), line=BRAND_RED)
t = tb(slide, Inches(2.3), Inches(3.8), Inches(8.7), Inches(0.4),
       "3年不行动累计代价", size=20, bold=True, color=BRAND_RED, align=PP_ALIGN.CENTER)
for line in [
    "直接成本：保守 3.3-6亿",
    "+ 一次极端事件：1-2亿",
    "+ 人才流失+竞争力下滑 → 不可逆",
    "",
    "vs 平台投入 → 不行动的代价远超投入",
]:
    ap(t.text_frame, line, size=14, color=WHITE, align=PP_ALIGN.CENTER)

notes(slide, "换个角度看——不做的代价。Year 1约1-2亿。Year 3最关键——竞争对手建起平台后，优秀投资经理会去工具更好的平台，这是不可逆的。3年累计保守3.3-6亿，远超平台投入。")


# ============================================================
# SLIDE 20 — Next Steps
# ============================================================
slide = light_slide("建议的下一步",
                    "建议立即启动Phase 1，6个月内见效")

steps = [
    ("1. 联合工作组组建（1周内）", "成立业务+IT联合工作组\n明确双方项目Owner、沟通机制、决策流程", TEAL),
    ("2. 需求确认研讨会（2-3周内）", "确认资产品种范围、与金证IBOR接口方案、优先级排序\n选定2-3位种子PM深度参与Phase 1设计", LIGHT_BLUE),
    ("3. Phase 1详细方案（1个月内）", "技术架构详细设计（含与6家现有供应商系统的集成方案）\n数据对接方案 + 种子用户培训计划", BRIGHT_BLUE),
    ("4. Phase 1启动 → 6个月交付", "CEO仪表盘上线，三大核心指标实时可见：\n① 全公司自营实时盈亏  ② 组合VaR  ③ 跨资产相关性热力图", BRAND_YELLOW),
]

for i, (title, desc, color) in enumerate(steps):
    y = Inches(1.4) + i * Inches(1.35)
    box(slide, Inches(0.6), y, Inches(0.08), Inches(1.1), fill=color)
    tb(slide, Inches(1.0), y, Inches(11), Inches(0.35),
       title, size=15, bold=True, color=TITLE_BLUE)
    tb(slide, Inches(1.0), y + Inches(0.4), Inches(11), Inches(0.8),
       desc, size=12, color=DARK_GRAY)

# Contact
tb(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
   "华锐技术  |  [联系人]  |  [电话]  |  [邮箱]",
   size=14, color=MED_GRAY, align=PP_ALIGN.CENTER)

notes(slide, "建议下一步分四步走。第一步最关键——一周内组建业务+IT联合工作组。第二步两周内做需求研讨会。第三步一个月内出详细方案。第四步启动建设——6个月后CEO打开手机就能看到三个数字：全公司自营实时P&L、组合VaR、跨资产相关性热力图。谢谢各位领导。")


# ============================================================
# END SLIDE
# ============================================================
slide = prs.slides.add_slide(LAYOUT_END)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "谢谢"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Microsoft YaHei"
    elif ph.placeholder_format.idx == 13:
        ph.text = "华锐技术 | Archforce Investment Management"
        for p in ph.text_frame.paragraphs:
            p.font.color.rgb = BRAND_YELLOW
            p.font.name = "Microsoft YaHei"


# ============================================================
# SAVE
# ============================================================
import os
os.remove(template_path)  # Clean up temp copy

output = "CEO_Deck_v2.3.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
