"""
Embed drawio-rendered PNGs into gy_ppt.pptx.
Each diagram is placed maintaining aspect ratio, centred horizontally,
starting at the top of the content area and capped at the available height.
"""

import os
import shutil
import sys
sys.path.insert(0, "gy")

import drawio_renderer
from pptx import Presentation
from pptx.util import Inches

shutil.copy("gy/gy_ppt.pptx", "gy/backup/gy_ppt_pre_diagrams.pptx")
prs = Presentation("gy/gy_ppt.pptx")
SLIDE_W = 13.33  # inches

os.makedirs("gy/tmp_diagrams", exist_ok=True)

# (slide_index, drawio_filename_stem, y_start_in, avail_h_in)
PLACEMENTS = [
    (9,  "diagram_4_竞争定位",      1.10, 5.26),  # S10 competitor
    (13, "diagram_H_1+16架构",      1.73, 4.62),  # S14 product arch
    (29, "diagram_7_CEP风控引擎",   1.10, 4.45),  # S30 CEP engine
    (30, "diagram_8_指令执行泳道",  1.10, 4.55),  # S31 order exec
    (32, "diagram_6_实施路线图",    1.10, 5.00),  # S33 roadmap
    (38, "diagram_3_924行情推演",   1.10, 4.51),  # S39 924 case
    (39, "diagram_5_ROI决策框架",   1.10, 4.28),  # S40 ROI
]

for slide_idx, stem, y_in, avail_h_in in PLACEMENTS:
    drawio_path = f"gy/drawio/{stem}.drawio"
    png_path = f"gy/tmp_diagrams/{stem}.png"

    w_px, h_px = drawio_renderer.convert(
        drawio_path, png_path, target_w_px=1300
    )
    ar = w_px / h_px

    # Fit within (SLIDE_W × avail_h_in), maintain AR, centre horizontally
    h_at_full_w = SLIDE_W / ar
    if h_at_full_w <= avail_h_in:
        img_w = SLIDE_W
        img_h = h_at_full_w
    else:
        img_h = avail_h_in
        img_w = avail_h_in * ar

    img_x = (SLIDE_W - img_w) / 2

    slide = prs.slides[slide_idx]
    slide.shapes.add_picture(
        png_path,
        Inches(img_x), Inches(y_in),
        Inches(img_w), Inches(img_h),
    )
    print(f"S{slide_idx+1} ← {stem}: ({img_x:.2f}\", {y_in:.2f}\") {img_w:.2f}\"×{img_h:.2f}\"")

prs.save("gy/gy_ppt.pptx")
print("\n✓ gy_ppt.pptx saved with 7 embedded diagrams")
