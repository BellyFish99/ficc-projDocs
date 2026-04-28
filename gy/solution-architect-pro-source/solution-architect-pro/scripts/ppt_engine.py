import os
import re
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# === Color Palette ===
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_GOLD = RGBColor(0xC4, 0x9A, 0x2A)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=14, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = alignment
    return txBox

def make_title_bar(slide, title_text, subtitle_text=""):
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=DARK_BLUE)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6), title_text, font_size=24, bold=True, color=WHITE)
    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.4), subtitle_text, font_size=13, color=ACCENT_GOLD)

def parse_manifest(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Parse YAML frontmatter
    yaml_match = re.search(r'^---\s*(.*?)\s*---', content, re.DOTALL)
    variables = yaml.safe_load(yaml_match.group(1)) if yaml_match else {}
    
    # Replace variables in content
    def replace_var(match):
        keys = match.group(1).split('.')
        val = variables
        for k in keys:
            val = val.get(k, match.group(0))
        return str(val)
    
    content = re.sub(r'\{\{(.*?)\}\}', replace_var, content)
    
    # Extract slides
    slides = re.findall(r'<!-- @slide:start -->\s*(.*?)\s*<!-- @slide:end -->', content, re.DOTALL)
    return variables, slides

def create_ppt(manifest_path, output_path):
    variables, slide_contents = parse_manifest(manifest_path)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Cover Slide (Special logic)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BLUE)
    add_textbox(slide, Inches(1.5), Inches(2), Inches(10), Inches(1), variables['client'] + "全资产POMS平台解决方案", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.5), "让" + variables['metrics']['aum'] + "自营资金拥有世界级的组合管理能力", font_size=22, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)
    
    # Process extracted slides
    for content in slide_contents:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, WHITE)
        
        # Simple parser for slide content
        title_match = re.search(r'## (.*)', content)
        action_title_match = re.search(r'\*\*Action Title\*\*: (.*)', content)
        
        title = title_match.group(1) if title_match else "Slide"
        action_title = action_title_match.group(1) if action_title_match else ""
        
        make_title_bar(slide, title, action_title)
        
        # Body content (simplified: just a big textbox for now)
        body_text = re.sub(r'## .*| \*\*Action Title\*\*: .*| \*\*Speaker Notes\*\*: .*', '', content, flags=re.MULTILINE).strip()
        add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5), body_text, font_size=14)
        
        # Speaker Notes
        notes_match = re.search(r'\*\*Speaker Notes\*\*: (.*)', content, re.DOTALL)
        if notes_match:
            slide.notes_slide.notes_text_frame.text = notes_match.group(1).strip()
            
    prs.save(output_path)
    print(f"PPT saved to {output_path}")

if __name__ == "__main__":
    create_ppt("/mnt/d/Work/gy/solution_master_v8.md", "/mnt/d/Work/gy/CEO_Deck_v8_from_Master.pptx")
