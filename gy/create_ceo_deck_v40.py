"""
CEO Deck v7 — 38 slides, built on Huarui M9 template.
Content sourced from solution_master_v40.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import shutil, copy

TEMPLATE = "/mnt/d/Work/gy/MKT2026_AIM_华锐投资管理平台 M9介绍PPT(现券管理)_20260316.pptx"
OUTPUT = "/mnt/d/Work/gy/CEO_Deck_v7.pptx"

shutil.copy(TEMPLATE, OUTPUT)
prs = Presentation(OUTPUT)

# Delete all existing slides
xml_slides = prs.slides._sldIdLst
for sldId in list(xml_slides):
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId:
        try: prs.part.drop_rel(rId)
        except: pass
    xml_slides.remove(sldId)

# Colors
DARK_NAVY   = RGBColor(0x03, 0x0E, 0x42)
MEDIUM_BLUE = RGBColor(0x15, 0x42, 0x8F)
BRAND_RED   = RGBColor(0xF5, 0x4D, 0x61)
BRAND_YELLOW= RGBColor(0xFF, 0xD0, 0x51)
LIGHT_BLUE  = RGBColor(0x37, 0x9D, 0xFF)
TEAL        = RGBColor(0x56, 0xC4, 0xD0)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MED_GRAY    = RGBColor(0x88, 0x88, 0x88)
TITLE_BLUE  = RGBColor(0x28, 0x36, 0x63)
ACCENT_ORANGE = RGBColor(0xEE, 0x82, 0x2F)
DEEP_RED    = RGBColor(0x2A, 0x0A, 0x12)

W = prs.slide_width
H = prs.slide_height

LAYOUT_COVER = prs.slide_layouts[14]
LAYOUT_TOC   = prs.slide_layouts[12]
LAYOUT_DARK  = prs.slide_layouts[15]
LAYOUT_LIGHT = prs.slide_layouts[11]
LAYOUT_BLANK = prs.slide_layouts[16]
LAYOUT_END   = prs.slide_layouts[10]

FONT = "MiSans Normal"

def tb(slide, left, top, width, height, text="", size=14,
       bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT, font=None, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font or FONT
    p.alignment = align
    return txBox

def ap(text_frame, text, size=13, bold=False, color=DARK_GRAY,
       align=PP_ALIGN.LEFT, font=None, space=Pt(4)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font or FONT
    p.alignment = align
    p.space_before = space
    return p

def box(slide, left, top, width, height, fill=None, line=None, line_width=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(line_width)
    else:
        s.line.fill.background()
    return s

def add_table(slide, left, top, width, rh, headers, rows, hdr_color=DARK_NAVY, fs=10, col_widths=None):
    nr = len(rows) + 1
    nc = len(headers)
    ts = slide.shapes.add_table(nr, nc, left, top, width, Emu(rh * nr))
    t = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            t.columns[i].width = cw
    else:
        cw = width // nc
        for i in range(nc): t.columns[i].width = cw
    for i, h in enumerate(headers):
        c = t.cell(0, i)
        c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = hdr_color
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(fs); p.font.bold = True
            p.font.color.rgb = WHITE; p.font.name = FONT
    for r, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = t.cell(r+1, ci)
            c.text = str(val)
            if r % 2 == 1:
                c.fill.solid(); c.fill.fore_color.rgb = LIGHT_GRAY
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(fs); p.font.color.rgb = DARK_GRAY; p.font.name = FONT
    return ts

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def dark_slide(title_text, subtitle_text=""):
    slide = prs.slides.add_slide(LAYOUT_DARK)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title_text
            for p in ph.text_frame.paragraphs:
                p.font.color.rgb = WHITE; p.font.name = FONT; p.font.bold = True
            break
    if subtitle_text:
        tb(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(0.4),
           subtitle_text, size=13, color=BRAND_YELLOW)
    return slide

def light_slide(title_text, subtitle_text=""):
    slide = prs.slides.add_slide(LAYOUT_LIGHT)
    box(slide, Inches(0), Inches(0), W, Inches(1.05), fill=DARK_NAVY)
    tb(slide, Inches(0.5), Inches(0.1), Inches(11.5), Inches(0.55),
       title_text, size=22, bold=True, color=WHITE)
    if subtitle_text:
        tb(slide, Inches(0.5), Inches(0.62), Inches(11.5), Inches(0.35),
           subtitle_text, size=12, color=BRAND_YELLOW)
    return slide

def divider_slide(section_num, act_text, subtitle):
    slide = prs.slides.add_slide(LAYOUT_DARK)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0: ph.text = ""
    box(slide, Inches(0), Inches(0), Inches(0.5), H, fill=BRAND_RED)
    tb(slide, Inches(0.8), Inches(2.2), Inches(12), Inches(0.5),
       section_num, size=14, color=BRAND_RED, bold=True)
    box(slide, Inches(0.8), Inches(2.8), Inches(5), Inches(0.05), fill=BRAND_RED)
    tb(slide, Inches(0.8), Inches(3.0), Inches(12), Inches(1.1),
       act_text, size=34, bold=True, color=WHITE)
    tb(slide, Inches(0.8), Inches(4.2), Inches(12), Inches(0.6),
       subtitle, size=16, color=BRAND_YELLOW)
    return slide

def callout_box(slide, left, top, width, height, title, body_lines,
                bg=DEEP_RED, border=BRAND_RED, title_color=BRAND_RED, body_color=WHITE):
    box(slide, left, top, width, height, fill=bg, line=border)
    t = tb(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.35),
           title, size=13, bold=True, color=title_color)
    tf = t.text_frame
    for line in body_lines:
        ap(tf, line, size=11, color=body_color)
    return t

# ─── SLIDE 01: COVER ────────────────────────────────────────────────────────
slide = prs.slides.add_slide(LAYOUT_COVER)
tb(slide, Inches(0.8), Inches(1.8), Inches(8), Inches(0.45),
   "国元证券全资产POMS平台解决方案", size=15, color=WHITE)
tb(slide, Inches(0.8), Inches(2.4), Inches(9), Inches(1.1),
   "华锐POMS：千亿资金的智慧管理引擎", size=38, bold=True, color=WHITE)
tb(slide, Inches(0.8), Inches(3.55), Inches(8), Inches(0.45),
   "中国版Mini-Aladdin  ·  配得优·算得快·控得稳·连得通·领得先", size=14, color=BRAND_YELLOW)
box(slide, Inches(0.85), Inches(4.2), Inches(1.5), Inches(0.04), fill=BRAND_RED)
tb(slide, Inches(0.8), Inches(4.35), Inches(8), Inches(0.35),
   "华锐科技  ·  2026年Q2", size=13, color=WHITE)
notes(slide, "Cover. 38-slide proposal for Guoyuan Securities. POMS = Portfolio & Order Management System. ~60 min + Q&A.")

# ─── SLIDE 02: SITUATION ────────────────────────────────────────────────────
slide = light_slide("国元自营：四年四倍增长，正站在能力跃升的关键窗口",
                    "Situation — 近千亿规模 · 战略机遇 · 必须行动的三重驱动")
headers = ["维度", "现状与数据", "关键含义"]
rows = [
    ["自营规模", "近千亿AUM，2025H1自营收入14.6亿（+43%营收占比）", "千亿规模下，每0.1%收益/风险改善 = 1亿价值"],
    ["收益目标", "低波动6-7%稳定收益战略目标", "需要全资产精细化管理，而非单一利率债方向判断"],
    ["杠杆规模", "约500亿回购融资杠杆", "融资成本管理、流动性风险是核心风控议题"],
    ["年交易量", "约3000亿年度交易量（3倍周转）", "执行质量每提升1bp = 3000万/年"],
    ["组织变革", "2025年8月重组：设立金融市场委员会+FICC总部+证券投资部", "战略转向全FICC一体化，需统一平台支撑"],
    ["时间窗口", "十五五数字化转型+信创替代+IFRS9 2026年全面执行", "政策驱动 + 竞争驱动 + 合规驱动三重叠加"],
]
add_table(slide, Inches(0.5), Inches(1.15), Inches(12.3), Emu(360000), headers, rows,
          col_widths=[Inches(1.8), Inches(5.5), Inches(5.0)])
callout_box(slide, Inches(0.5), Inches(4.95), Inches(12.3), Inches(1.3),
    "核心命题：POMS不是IT项目，是投资能力的投资",
    ["千亿自营资金，每年交易3000亿，管理500亿杠杆——这个规模，需要世界级的管理工具。",
     "华泰/平安平台化已进入深水区。窗口期有限——现在是追赶的最佳时机，延迟的代价是差距指数级拉大。"],
    bg=RGBColor(0x05, 0x12, 0x4A), border=LIGHT_BLUE, title_color=LIGHT_BLUE, body_color=WHITE)
notes(slide, "Situation. Guoyuan prop trading grew 4x in 4 years. Three forcing functions: regulatory (15th 5-year plan), competitive (Huatai/PA already deep), compliance (IFRS9 2026).")

# ─── DIVIDER: SECTION 1 ─────────────────────────────────────────────────────
divider_slide("SECTION 1 — COMPLICATION",
              "五大能力缺口：千亿资金的五个战略制约",
              "配得优 · 算得快 · 控得稳 · 连得通 · 领得先")

# ─── SLIDE 03: FIVE GAPS OVERVIEW ───────────────────────────────────────────
slide = light_slide("五大能力跃升目标：从L1起点到L4行业领先",
                    "当前能力成熟度 L1，目标 L4——每个缺口都是千亿资金的价值漏损点")
headers = ["能力域", "目标", "当前", "目标", "缺口", "核心问题"]
rows = [
    ["① 全资产组合管理", "配得优", "L1 ▓░░░░", "L4", "████", "只有利率债，看不到全资产机会；组合靠拍脑袋调仓"],
    ["② 高性能计算", "算得快", "L1 ▓░░░░", "L4", "████", "风险T+1，无实时估值；定价依赖Bloomberg"],
    ["③ 量化策略", "算得快", "L0 ░░░░░", "L3", "███", "无利差/基差/新老券分析工具；策略靠经验"],
    ["④ 实时风控", "控得稳", "L1 ▓░░░░", "L4", "████", "回撤发生才发现；924类极端事件无预案"],
    ["⑤ 决策→执行闭环", "连得通", "L2 ▓▓░░░", "L4", "██", "研究与投资脱节；多腿交易无系统支持"],
    ["⑥ 信创领先", "领得先", "L1 ▓░░░░", "L3", "███", "依赖Calypso/Bloomberg，信创压力大"],
]
add_table(slide, Inches(0.5), Inches(1.15), Inches(12.3), Emu(350000), headers, rows,
          col_widths=[Inches(2.5), Inches(1.2), Inches(1.5), Inches(0.8), Inches(0.8), Inches(5.5)])
callout_box(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(1.3),
    "McKinsey问题树：将'近千亿自营资金实现低波动6%稳定收益'MECE分解为5个核心能力跃升目标",
    ["每个L1→L4的缺口，都是每年数千万到数亿的价值漏损——能力差距 = 收益差距。",
     "五大缺口相互依赖：没有组合管理引擎，计算能力无处落脚；没有实时风控，量化策略无法放量。"],
    bg=RGBColor(0x05, 0x12, 0x4A), border=BRAND_YELLOW, title_color=BRAND_YELLOW)
notes(slide, "Five capability gaps. Each represents a strategic constraint on the 100B AUM. L1→L4 maturity model.")

# ─── SLIDE 04: GAP 1 配得优 ──────────────────────────────────────────────────
slide = light_slide("缺口①【配得优】：千亿资金看不到全局，配置靠经验而非数据",
                    "当前L1 — 3个系统才能拼出全貌；跨资产机会系统性错失")
# Two-column layout
box(slide, Inches(0.5), Inches(1.12), Inches(5.9), Inches(4.7),
    fill=RGBColor(0xF8, 0xF9, 0xFF), line=RGBColor(0xCC, 0xD0, 0xE0))
t = tb(slide, Inches(0.6), Inches(1.2), Inches(5.7), Inches(0.35),
       "当前痛点（As-Is）", size=14, bold=True, color=DARK_NAVY)
tf = t.text_frame
for line in [
    "▪ 打开3个系统才能拼出昨日完整持仓",
    "▪ 利率债、权益、黄金分散在不同账户，无统一视图",
    "▪ 跨资产相关性分析靠Excel，每次需2-3天",
    "▪ 组合调仓靠投资经理经验，无数量化依据",
    "▪ 公司→部门→交易台→个人贡献无法穿透评估",
    "▪ 虚拟组合/what-if试算没有工具",
    "▪ SAA战略指令无法量化传导到TAA执行层",
]:
    ap(tf, line, size=12, color=DARK_GRAY, space=Pt(5))

box(slide, Inches(6.55), Inches(1.12), Inches(6.25), Inches(4.7),
    fill=RGBColor(0xF0, 0xF8, 0xF4), line=RGBColor(0x56, 0xC4, 0xD0))
t2 = tb(slide, Inches(6.65), Inches(1.2), Inches(6.05), Inches(0.35),
        "POMS解决方案（To-Be）", size=14, bold=True, color=TEAL)
tf2 = t2.text_frame
for line in [
    "✓ 全资产统一持仓模型：8大类资产一屏管理",
    "✓ 多层组合树：公司→部门→策略→个人，穿透实时",
    "✓ 因子驱动：2,200+风险因子重构配置决策",
    "✓ QP优化引擎：500+约束下自动求最优调仓路径",
    "✓ 跨资产相关性矩阵：实时识别'假分散'陷阱",
    "✓ Shadow Portfolio：100+虚拟组合并行试算",
    "✓ SAA→TAA：投委会战略指令→量化执行，零偏差",
]:
    ap(tf2, line, size=12, color=DARK_GRAY, space=Pt(5))

callout_box(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.85),
    "价值量化（千亿AUM）",
    ["配置优化收益+0.02%-0.05% → 每年增加2000-5000万  |  再平衡效率提升 → 减少跟踪误差5-10bps = 2000-3000万/年"],
    bg=RGBColor(0x05, 0x12, 0x4A), border=TEAL, title_color=TEAL)
notes(slide, "Gap 1: Portfolio management. Currently need 3 systems to see full picture. POMS unifies 8 asset classes.")

# ─── SLIDE 05: GAP 2 算得快 ──────────────────────────────────────────────────
slide = light_slide("缺口②③【算得快】：风险T+1、定价依赖外部、利差策略无工具",
                    "当前L1 — 实时计算能力缺失使相对价值策略无从落地")
headers = ["痛点", "当前状态", "竞争影响", "POMS方案"]
rows = [
    ["风险计量", "T+1批量跑批，今天看昨天数据", "极端行情下无法实时响应，等于盲飞", "分布式实时计算，VaR/CVaR秒级全量"],
    ["定价能力", "依赖Bloomberg/Calypso，成本高", "许可费+数据费每年数百万，且受制于人", "自研全品种定价引擎，利率曲线/波动率曲面/MC"],
    ["利差分析", "无系统化利差/基差工具", "看到机会无法量化验证，错失Alpha", "利差策略/基差策略/新老券分析工具集"],
    ["多策略回测", "无多timeframe信号回测框架", "策略只能靠经验，无法系统验证", "多timeframe信号引擎+历史回测框架"],
    ["Shadow试算", "100+虚拟组合无法并行", "调仓决策需1-2天，机会可能已过", "100+虚拟组合毫秒级并行计算"],
]
add_table(slide, Inches(0.5), Inches(1.15), Inches(12.3), Emu(370000), headers, rows,
          col_widths=[Inches(1.5), Inches(3.0), Inches(3.5), Inches(4.3)])
# Stats row
for i, (num, label, color) in enumerate([
    ("2,000+", "持仓规模\n(2000次定价/秒)", LIGHT_BLUE),
    ("毫秒级", "定价与估值\n响应延迟", TEAL),
    ("10万+", "CEP事件吞吐\n(events/sec)", BRAND_YELLOW),
    ("0", "Bloomberg依赖\n(完全自研替代)", BRAND_RED),
]):
    x = Inches(0.5) + i * Inches(3.1)
    box(slide, x, Inches(5.05), Inches(2.8), Inches(0.95), fill=DARK_NAVY, line=color)
    tb(slide, x + Inches(0.1), Inches(5.1), Inches(2.6), Inches(0.45),
       num, size=24, bold=True, color=color, align=PP_ALIGN.CENTER)
    tb(slide, x + Inches(0.1), Inches(5.55), Inches(2.6), Inches(0.4),
       label, size=10, color=WHITE, align=PP_ALIGN.CENTER)
notes(slide, "Gap 2+3: Computing speed. T+1 risk is a fundamental blocker for relative value strategies.")

# ─── SLIDE 06: GAP 3 控得稳 + 924 ──────────────────────────────────────────
slide = light_slide("缺口④【控得稳】：实时风控缺失——924事件是最真实的教训",
                    "当前L1 — 回撤发生才发现；极端市场无预案；2024年924是行业警醒")
# 924 timeline
headers = ["时间", "Without POMS（当前）", "With POMS（目标）"]
rows = [
    ["09:30", "从新闻得知央行干预债市", "CEP引擎捕获利率异常波动，实时触发警报"],
    ["09:31", "—", "压力测试自动启动'924情景' → 组合实时损益: -1.2亿"],
    ["09:35", "—", "系统生成减仓建议：优先前10只高流动性债券，合规预检通过"],
    ["09:38", "—", "一键multi-leg执行，滑点仅3bp"],
    ["09:45", "开始手工计算，数据是昨天的", "—"],
    ["10:15", "大致估出亏损~1.5亿，但不确定", "—"],
    ["11:30", "手工指令找领导签字，价格又跌了", "—"],
    ["收盘", "实际亏损2亿，事后才看清损失", "实际亏损0.8亿，全程实时可见 → 差异: 1.2亿"],
]
add_table(slide, Inches(0.5), Inches(1.15), Inches(12.3), Emu(320000), headers, rows,
          col_widths=[Inches(1.0), Inches(5.65), Inches(5.65)])
callout_box(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(1.0),
    "924教训的本质：流动性冲击 > 价格冲击",
    ["真正的危险不是债券跌价，而是不知道哪些债券'卖不出去'——POMS实时计算每只持仓的流动性冲击成本",
     "极端压力测试三层：历史情景(924/包商银行/2013钱荒) + 假设情景 + 反向压力测试（倒推临界点）"])
callout_box(slide, Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.7),
    "POMS风控能力",
    ["事件驱动CEP引擎（10万+ events/sec）· 秒级回撤监控 · 流动性冲击模型 · 可配置5000+合规规则 · 全公司统一风险视图"],
    bg=RGBColor(0x05, 0x12, 0x4A), border=LIGHT_BLUE, title_color=LIGHT_BLUE)
notes(slide, "Gap 4: Risk control. 924 event Sep 24 2024 is the real industry proof point. 1.2B difference in one day.")

# ─── SLIDE 07: GAP 4+5 连得通+领得先 ───────────────────────────────────────
slide = light_slide("缺口⑤⑥【连得通·领得先】：研究→交易断层 + 信创窗口倒计时",
                    "决策→执行全链路漏损 | Bloomberg/Calypso依赖 | 头部券商差距拉大")
# Two panels
box(slide, Inches(0.5), Inches(1.12), Inches(6.1), Inches(4.3),
    fill=RGBColor(0xF8, 0xF9, 0xFF), line=MEDIUM_BLUE)
t = tb(slide, Inches(0.6), Inches(1.2), Inches(5.9), Inches(0.35),
       "连得通：研究→决策→指令→执行 断层", size=14, bold=True, color=DARK_NAVY)
tf = t.text_frame
for line in [
    "▪ 投研发现利差机会，结论停在PPT里",
    "▪ 组合调整→指令生成 全靠手工，延迟1-2天",
    "▪ multi-leg交易无系统支持（如利差交易同时买卖两只券）",
    "▪ 多柜台（场内/场外/期货）各自独立，统一视图缺失",
    "▪ 绩效归因结果不驱动下一次决策（断环）",
    "POMS方案：",
    "✓ 投研信号→因子库→组合再平衡→指令→OEMS 全链路",
    "✓ multi-leg交易引擎：利差交易一键执行",
    "✓ Campisi 2.0 + Brinson归因闭环驱动下次决策",
]:
    ap(tf, line, size=11, color=DARK_GRAY if not line.startswith("✓") else TEAL, space=Pt(4))

box(slide, Inches(6.75), Inches(1.12), Inches(6.1), Inches(4.3),
    fill=RGBColor(0xF8, 0xF4, 0xFF), line=BRAND_RED)
t2 = tb(slide, Inches(6.85), Inches(1.2), Inches(5.9), Inches(0.35),
        "领得先：信创窗口期 + 竞争代差", size=14, bold=True, color=BRAND_RED)
tf2 = t2.text_frame
for line in [
    "▪ 当前依赖Bloomberg/Calypso（数百万/年许可费）",
    "▪ '十五五'明确数字化转型+信创替代硬要求",
    "▪ 华泰'大象'平台+平安FITS已进入深水区",
    "▪ 差距每延迟一年，追赶成本翻倍",
    "▪ IFRS9 2026年全面执行——无系统P&L波动不可控",
    "POMS方案：",
    "✓ 全栈信创适配，国产OS/DB/中间件，消除外部依赖",
    "✓ 18个月追赶行业头部3-5年建设成果",
    "✓ 对标Aladdin/SecDB持续演进路径",
]:
    ap(tf2, line, size=11, color=DARK_GRAY if not line.startswith("✓") else TEAL, space=Pt(4))

callout_box(slide, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.75),
    "五大缺口 × 五大引擎 = 精确映射",
    ["华锐POMS五大引擎逐一对应五大跃升目标——不是功能堆砌，是需求→方案的精确映射。No problem, no solution."],
    bg=RGBColor(0x05, 0x12, 0x4A), border=BRAND_YELLOW, title_color=BRAND_YELLOW)
notes(slide, "Gap 5+6: Connectivity and leadership. Research-to-trade gap is a daily value leakage. Xinchuang (信创) window is closing.")

# ─── DIVIDER: SECTION 2 ─────────────────────────────────────────────────────
divider_slide("SECTION 2 — WHY ARCHFORCE",
              "为什么是华锐：唯一的一体化国产选择",
              "厂商矩阵 · 四大不可替代优势 · 为什么是现在")

# ─── SLIDE 08: VENDOR MATRIX ────────────────────────────────────────────────
slide = light_slide("厂商对比：华锐是唯一同时覆盖FICC+买方组合管理的国产方案",
                    "市场上没有第二家国产厂商能一体化满足国元五大跃升目标")
headers = ["能力域", "国际厂商\n(Calypso/Murex)", "数据厂商\n(Wind/同花顺)", "交易系统厂商\n(金证/恒生)", "华锐POMS"]
rows = [
    ["全资产组合管理", "✓ 强（高成本）", "✗ 无", "△ 弱（交易为主）", "✓✓ 专项设计"],
    ["实时风险计量", "✓ 强", "△ 分析为主", "△ 弱", "✓✓ 千亿级验证"],
    ["量化策略引擎", "△ 部分", "△ 数据分析", "✗ 无", "✓ 相对价值+多策略"],
    ["事件驱动风控", "✓ 强", "✗ 无", "△ 弱", "✓✓ CEP+极端压测"],
    ["指令执行全链路", "✓ 强", "✗ 无", "✓ 强（交易核心）", "✓ 全生命周期"],
    ["绩效归因", "✓ 强", "△ 报告为主", "✗ 无", "✓✓ Campisi 2.0+Brinson"],
    ["信创适配", "✗ 不支持", "△ 部分", "✓ 部分", "✓✓ 全栈国产化"],
    ["对接已建IBOR", "需定制", "✗ 无", "需定制", "✓✓ 专项对接方案"],
    ["价格与服务连续性", "高成本+受制于人", "按数据订阅", "系统级依赖", "✓ 国内团队+快速响应"],
]
add_table(slide, Inches(0.5), Inches(1.12), Inches(12.3), Emu(340000), headers, rows,
          col_widths=[Inches(2.2), Inches(2.0), Inches(2.1), Inches(2.2), Inches(3.8)])
callout_box(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.75),
    "华锐独特定位：FICC卖方基因 + 买方组合管理能力 = 双重核心能力",
    ["没有第二家国产厂商同时具备'高性能FICC系统'和'Aladdin级组合管理平台'的双重基因——这是华锐的护城河"],
    bg=RGBColor(0x05, 0x12, 0x4A), border=BRAND_RED, title_color=BRAND_RED)
notes(slide, "Vendor matrix. Archforce is the only domestic vendor covering both FICC and buy-side portfolio management.")

# ─── SLIDE 09: 4 IRREPLACEABLE ADVANTAGES ───────────────────────────────────
slide = light_slide("华锐四大不可替代优势：为什么一体化 > 拼凑集成",
                    "6年毕马威领先金融科技50强 · 128件发明专利 · 千亿级客户验证")
for i, (icon, title, points, color) in enumerate([
    ("①", "不只是数据对接\n完整FICC业务平台", [
        "从数据集成→业务应用的全栈平台",
        "覆盖组合管理/定价/风控/归因/执行",
        "一个厂商解决全部问题，无需拼凑集成",
        "避免多厂商接口灾难（数据不一致/责任推诿）",
    ], LIGHT_BLUE),
    ("②", "不只是分析工具\n覆盖交易执行全链路", [
        "相对价值策略到multi-leg执行一条龙",
        "QP优化引擎→指令→OEMS→TCA闭环",
        "投研信号直接驱动组合调整，消除死亡谷",
        "Campisi 2.0+Brinson归因持续改进Alpha",
    ], TEAL),
    ("③", "不只是国内\n香港国际业务中心", [
        "支持跨境投资，港股/外债/外汇统一管理",
        "国际化资产不用另找供应商",
        "HKMA/SEC合规框架支持",
        "跨境业务扩张无缝衔接",
    ], BRAND_YELLOW),
    ("④", "不只是替代\n对标Aladdin定义天花板", [
        "愿景：国元的Mini-Aladdin",
        "千亿规模值得世界级系统",
        "10亿研发投入+500人团队持续演进",
        "3年路线图，而非项目交付即止步",
    ], BRAND_RED),
]):
    x = Inches(0.5) + i * Inches(3.1)
    box(slide, x, Inches(1.12), Inches(2.9), Inches(4.9), fill=DARK_NAVY, line=color)
    tb(slide, x+Inches(0.1), Inches(1.18), Inches(2.7), Inches(0.6),
       icon, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
    tb(slide, x+Inches(0.1), Inches(1.75), Inches(2.7), Inches(0.55),
       title, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(slide, x+Inches(0.3), Inches(2.38), Inches(2.3), Inches(0.03), fill=color)
    t = tb(slide, x+Inches(0.15), Inches(2.5), Inches(2.7), Inches(2.5), "", size=11)
    tf = t.text_frame
    tf.word_wrap = True
    for j, pt in enumerate(points):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = "▸ " + pt; p.font.size = Pt(11); p.font.color.rgb = WHITE
        p.font.name = FONT; p.space_before = Pt(5)

# Credentials bar
box(slide, Inches(0.5), Inches(6.18), Inches(12.3), Inches(0.65), fill=DARK_NAVY)
tb(slide, Inches(0.6), Inches(6.23), Inches(12.1), Inches(0.5),
   "毕马威领先金融科技50强（连续6年）  ·  国家级专精特新重点小巨人  ·  IDC中国FinTech50（连续3年）  ·  128件发明专利  ·  500人团队  ·  10亿研发投入",
   size=10, color=BRAND_YELLOW, align=PP_ALIGN.CENTER)
notes(slide, "4 irreplaceable advantages. The key differentiator: only domestic vendor with both FICC sell-side AND buy-side portfolio management DNA.")

# ─── SLIDE 10: WHY NOW ──────────────────────────────────────────────────────
slide = light_slide("为什么是现在：三大时间窗口同时打开，延迟的代价指数级上升",
                    "监管驱动 + 竞争驱动 + 合规驱动 = 18个月黄金建设窗口")
headers = ["驱动力", "具体压力", "时间节点", "延迟一年的代价"]
rows = [
    ["监管/信创", "'十五五'规划明确数字化转型+信创替代硬指标", "2026-2030", "不达标的合规风险；窗口期关闭后成本翻倍"],
    ["IFRS9合规", "中国版IFRS9 2026年全面执行，无系统→P&L波动不可控", "2026年1月", "1年内必须上线基础功能；否则利润表大幅波动"],
    ["竞争差距", "华泰'大象'+平安FITS平台化进入深水区，持续投入", "持续扩大", "差距从'可追赶'变为'代差'，人才和客户双流失"],
    ["机会成本", "千亿规模下，每延迟1年 = 错失1-2亿直接价值", "每年", "3年不行动累计代价保守估算3.3-6亿"],
    ["极端事件", "924类事件平均5年一次，下次无预案就是纯粹风险暴露", "随机", "一次极端事件避损价值1-2亿，有无POMS差别显著"],
]
add_table(slide, Inches(0.5), Inches(1.15), Inches(12.3), Emu(380000), headers, rows,
          col_widths=[Inches(1.8), Inches(4.5), Inches(1.8), Inches(4.2)])
callout_box(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.3),
    "华锐独有：18个月追赶头部3-5年",
    ["华泰花了3年建设'大象'平台，平安花了4年建设FITS——华锐的方案让国元在18个月内达到同等水平。",
     "后发优势：站在华泰/平安踩过的坑上，直接跳到最佳实践；避免重复犯错，起点更高、速度更快。",
     "Phase 1在6个月内即可让CEO看到全公司自营实时全貌——这是最快的价值验证。"])
notes(slide, "Why now: 3 forcing functions converging. IFRS9 is a hard deadline. 18-month catch-up vs 3-5 years competitors spent.")

# ─── DIVIDER: SECTION 3 ─────────────────────────────────────────────────────
divider_slide("SECTION 3 — POMS PLATFORM",
              "华锐POMS：五大引擎 · 十七大核心模块",
              "全资产组合管理 · 高性能计算 · 量化策略 · 实时风控 · 指令执行")

# ─── SLIDE 11: FIVE ENGINES ARCHITECTURE ────────────────────────────────────
slide = light_slide("POMS平台：五大引擎·四向集成·十七大核心模块",
                    "对接已建IBOR+行情中心+数据总线；构建投资业务智能层——中国版Mini-Aladdin")
# Architecture diagram using shapes
# Top: Workbench
box(slide, Inches(3.0), Inches(1.12), Inches(7.3), Inches(0.55), fill=DARK_NAVY, line=LIGHT_BLUE)
tb(slide, Inches(3.0), Inches(1.15), Inches(7.3), Inches(0.45),
   "  ↑ 投资工作台（PM台 · 交易台 · 风控台 · 策略台 · 运营台）", size=12, color=WHITE, bold=True)

# Left: Research
box(slide, Inches(0.3), Inches(2.0), Inches(2.5), Inches(3.2), fill=RGBColor(0x1C, 0x2C, 0x5A), line=TEAL)
t = tb(slide, Inches(0.4), Inches(2.05), Inches(2.3), Inches(0.35), "← 投研系统", size=11, bold=True, color=TEAL)
tf = t.text_frame
for line in ["策略信号", "研究数据", "宏观因子", "利差/基差分析"]:
    ap(tf, "· " + line, size=10, color=WHITE, space=Pt(4))

# Center: Five engines
box(slide, Inches(3.0), Inches(1.75), Inches(7.3), Inches(3.7), fill=RGBColor(0x03, 0x1A, 0x52), line=LIGHT_BLUE)
for i, (eng, color2) in enumerate([
    ("① 全资产组合管理引擎（配得优）— 本次深度介绍", BRAND_YELLOW),
    ("② 高性能计算引擎（算得快）— 定价/估值/风险计量", LIGHT_BLUE),
    ("③ 量化策略引擎（算得快+配得优）— 相对价值/多策略/回测", TEAL),
    ("④ 事件驱动风控引擎（控得稳）— CEP/回撤/压力测试/合规", BRAND_RED),
    ("⑤ 指令与执行管理引擎（连得通）— 指令→审批→执行→归因", BRAND_YELLOW),
]):
    y = Inches(1.85) + i * Inches(0.68)
    box(slide, Inches(3.1), y, Inches(7.1), Inches(0.6),
        fill=RGBColor(0x05, 0x20, 0x60) if i != 0 else RGBColor(0x1A, 0x3A, 0x00),
        line=color2)
    tb(slide, Inches(3.2), y + Inches(0.1), Inches(6.9), Inches(0.45),
       eng, size=11, bold=(i == 0), color=color2 if i == 0 else WHITE)

# Right: OEMS
box(slide, Inches(10.5), Inches(2.0), Inches(2.5), Inches(3.2), fill=RGBColor(0x1C, 0x2C, 0x5A), line=TEAL)
t2 = tb(slide, Inches(10.55), Inches(2.05), Inches(2.35), Inches(0.35), "OEMS系统 →", size=11, bold=True, color=TEAL)
tf2 = t2.text_frame
for line in ["交易执行", "算法交易", "multi-leg", "场内/场外/期货"]:
    ap(tf2, "· " + line, size=10, color=WHITE, space=Pt(4))

# Bottom: IBOR
box(slide, Inches(3.0), Inches(5.62), Inches(7.3), Inches(0.55), fill=DARK_NAVY, line=BRAND_YELLOW)
tb(slide, Inches(3.0), Inches(5.65), Inches(7.3), Inches(0.45),
   "  ↓ IBOR 3.0 + 行情中心 + 数据总线（在建）  |  统一簿记 · 实时头寸 · 行情 · 数据总线",
   size=11, color=BRAND_YELLOW)

callout_box(slide, Inches(0.3), Inches(6.35), Inches(13.0), Inches(0.5),
    "架构定位：POMS是业务智能层，不替代IBOR——IBOR=数据引擎，POMS=智慧驾驶层（Mini-Aladdin）",
    [], bg=RGBColor(0x05, 0x12, 0x4A), border=BRAND_YELLOW, title_color=BRAND_YELLOW, body_color=WHITE)
notes(slide, "5 engines architecture. POMS = business intelligence layer on top of IBOR data foundation. 4-directional integration.")

# ─── SLIDE 12: 4.4.1 OPENER ─────────────────────────────────────────────────
slide = light_slide("引擎①：全资产组合管理引擎——十大核心能力构成国元自营驾驶舱",
                    "配得优 — 这是国元最核心的需求，也是本次深度介绍的重点（13张详细展开）")
# 10 capabilities grid: 2 rows × 5 cols
caps = [
    ("①", "因子驱动\n组合构建", "2,200+风险因子\nSAA→TAA精准联动", LIGHT_BLUE),
    ("②", "全资产统一\n持仓模型", "8大类资产\n一屏管理", TEAL),
    ("③", "QP优化\n引擎", "500+约束\n毫秒级求解", BRAND_YELLOW),
    ("④", "跨资产相关性\n+Shadow组合", "100+虚拟组合\n并行试算", LIGHT_BLUE),
    ("⑤", "Campisi 2.0\n固收绩效归因", "五大效应\n精确拆解", TEAL),
    ("⑥", "Brinson模型\n权益/基金归因", "配置+选券\n效应分解", BRAND_YELLOW),
    ("⑦", "因子风险分解\n+MCTR", "风险预算\n精细分配", BRAND_RED),
    ("⑧", "Component VaR\n主动风险归因", "四层穿透\n逐仓风险溯源", LIGHT_BLUE),
    ("⑨", "敏感度分析\n(DV01/CS01)", "实时风险\n地图", TEAL),
    ("⑩", "情景/压力测试\n流动性+衍生品", "三层体系\n924类预案", BRAND_RED),
]
for i, (num, title, desc, color) in enumerate(caps):
    row = i // 5
    col = i % 5
    x = Inches(0.5) + col * Inches(2.55)
    y = Inches(1.15) + row * Inches(2.45)
    box(slide, x, y, Inches(2.4), Inches(2.25), fill=DARK_NAVY, line=color)
    tb(slide, x+Inches(0.05), y+Inches(0.05), Inches(2.3), Inches(0.45),
       num, size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
    tb(slide, x+Inches(0.05), y+Inches(0.5), Inches(2.3), Inches(0.7),
       title, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(slide, x+Inches(0.4), y+Inches(1.25), Inches(1.6), Inches(0.03), fill=color)
    tb(slide, x+Inches(0.05), y+Inches(1.35), Inches(2.3), Inches(0.8),
       desc, size=10, color=MED_GRAY, align=PP_ALIGN.CENTER)

callout_box(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.65),
    "Action Title：国元近千亿自营资金的核心挑战不是'该不该买债'，而是'如何在全资产中找到低波动、高确定性的最优配置'",
    [], bg=RGBColor(0x05, 0x12, 0x4A), border=BRAND_YELLOW, title_color=BRAND_YELLOW)
notes(slide, "4.4.1 opener: 10 core capabilities. Next 12 slides deep-dive each.")

# ─── SLIDE 13: 4.4.1.1 因子驱动 ────────────────────────────────────────────
slide = light_slide("4.4.1.1 因子驱动：用2,200+风险因子重构决策体系",
                    "传统'品种权重配置'无法解释组合真实风险来源——因子语言才能连接SAA战略与TAA执行")
headers = ["盲区", "典型场景", "后果"]
rows = [
    ["看似分散，实则集中", "债券60%/股票20%/黄金10%，但债券与黄金同向受利率驱动", "以为分散，利率风险敞口实际超过80%"],
    ["SAA战略无法落地", "投委会定了'降低利率敏感度'，交易台不知如何量化执行", "战略停留在PPT，执行靠经验判断"],
    ["风险超限才发现", "发现组合回撤超标时，无法快速定位是哪类风险因子导致", "事后归因，无法事前防控"],
]
tb(slide, Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.35),
   "当前痛点 — 权重管理的三大盲区", size=13, bold=True, color=DARK_NAVY)
add_table(slide, Inches(0.5), Inches(1.52), Inches(12.3), Emu(340000), ["盲区", "典型场景", "后果"], rows,
          col_widths=[Inches(2.5), Inches(5.5), Inches(4.3)])

tb(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.35),
   "POMS方案 — 因子驱动的三层转变", size=13, bold=True, color=TEAL)
for i, (layer, title, content) in enumerate([
    ("第一层", "统一风险语言",
     "全持仓分解为2,200+风险因子：利率(关键期限/久期/凸性/DV01) + 信用(DTS/评级/行业利差) + 权益(Beta/行业/风格) + 商品(黄金/原油) · 跨资产敞口一目了然"),
    ("第二层", "SAA→TAA精准联动",
     "投委会指令'将利率风险预算从40%降至25%' → 直接转化为因子约束输入QP优化 → 系统自动计算调仓路径 · 战略决策→量化执行，消除理解偏差"),
    ("第三层", "持续监控与预警",
     "因子敞口实时监控，超限自动预警 · 每日因子贡献报告 · 动态风险预算：已用/剩余/预警线三档显示"),
]):
    x = Inches(0.5) + i * Inches(4.1)
    box(slide, x, Inches(3.45), Inches(3.9), Inches(2.45), fill=DARK_NAVY, line=LIGHT_BLUE)
    tb(slide, x+Inches(0.1), Inches(3.52), Inches(3.7), Inches(0.3),
       f"{layer}：{title}", size=11, bold=True, color=BRAND_YELLOW)
    t = tb(slide, x+Inches(0.1), Inches(3.88), Inches(3.7), Inches(2.0), "", size=10)
    tf = t.text_frame; tf.word_wrap = True
    for sent in content.split("·"):
        if sent.strip():
            ap(tf, "▸ " + sent.strip(), size=10, color=WHITE, space=Pt(4))

callout_box(slide, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.72),
    "核心价值：决策变革——从传统'品种配置'转向'因子风险预算'",
    ["用2,200+因子语言统一SAA战略与TAA执行，确保低波动6%稳定收益目标的可持续性"],
    bg=RGBColor(0x05, 0x12, 0x4A), border=LIGHT_BLUE, title_color=LIGHT_BLUE)
notes(slide, "4.4.1.1: Factor-driven portfolio construction. 2200+ risk factors. SAA→TAA linkage.")

# ─── SLIDE 14: 4.4.1.2 统一持仓模型 ───────────────────────────────────────
slide = light_slide("4.4.1.2 全资产统一持仓模型——八类资产一屏管理",
                    "投资经理第一次在一个界面看到完整组合，不再需要打开3个系统才能拼出全貌")
# Asset tree
box(slide, Inches(0.5), Inches(1.12), Inches(6.1), Inches(3.8), fill=RGBColor(0x03, 0x1A, 0x52), line=LIGHT_BLUE)
t = tb(slide, Inches(0.6), Inches(1.18), Inches(5.9), Inches(0.35), "8大资产类别全覆盖", size=13, bold=True, color=LIGHT_BLUE)
tf = t.text_frame
assets = ["利率债/信用债/可转债（含回购/质押）", "A股/港股/ETF/场内基金",
          "公募基金/私募基金/FOF", "黄金（现货/期货/黄金ETF）",
          "商品期货（能化/农产品/有色）", "衍生品（国债期货/股指/期权/IRS）",
          "非标资产（信托/ABS/永续债/优先股）", "外汇（即期/远期/掉期）"]
for a in assets:
    ap(tf, "▸ " + a, size=11, color=WHITE, space=Pt(3))

# Portfolio tree
box(slide, Inches(0.5), Inches(5.05), Inches(6.1), Inches(1.72), fill=RGBColor(0x03, 0x1A, 0x52), line=TEAL)
t2 = tb(slide, Inches(0.6), Inches(5.1), Inches(5.9), Inches(0.35), "多层组合树（穿透式管理）", size=13, bold=True, color=TEAL)
tf2 = t2.text_frame
for line in ["公司层（全公司自营总览）→ 部门层（固收/权益/量化/综合）",
             "→ 策略层（利率债/信用/相对价值策略）",
             "→ 子组合层 → 投资经理层（个人持仓贡献）"]:
    ap(tf2, line, size=11, color=WHITE, space=Pt(3))

# Analysis dimensions table
headers2 = ["分析维度", "展示内容", "更新频率"]
rows2 = [
    ["品种分布", "各资产类别市值占比、盈亏贡献", "实时Tick级"],
    ["久期分布", "组合加权久期、DV01、各期限桶敞口", "实时"],
    ["评级分布", "AAA/AA+/AA/AA-各评级市值比例", "日内刷新"],
    ["行业分布", "金融/城投/产业各行业集中度", "日内刷新"],
    ["集中度", "单一发行人/单一券种/对手方敞口", "实时"],
]
add_table(slide, Inches(6.75), Inches(1.12), Inches(6.1), Emu(340000),
          ["分析维度", "展示内容", "频率"], rows2,
          col_widths=[Inches(1.5), Inches(3.4), Inches(1.2)])

callout_box(slide, Inches(6.75), Inches(4.15), Inches(6.1), Inches(1.0),
    "真实组合 + Shadow Portfolio并行",
    ["真实组合：T+0实时，对接IBOR统一簿记",
     "Shadow Portfolio：100+虚拟组合并行，下单前先试算后决定"],
    bg=RGBColor(0x05, 0x12, 0x4A), border=TEAL, title_color=TEAL)
callout_box(slide, Inches(6.75), Inches(5.3), Inches(6.1), Inches(1.47),
    "核心价值",
    ["投资经理第一次在一个界面看到完整组合",
     "不再有'看不见的风险'，不再有数据孤岛"],
    bg=RGBColor(0x05, 0x12, 0x4A), border=BRAND_YELLOW, title_color=BRAND_YELLOW)
notes(slide, "4.4.1.2: Unified position model. 8 asset classes. Multi-level portfolio tree. T+0 real-time.")

# ─── SLIDE 15: 4.4.1.3 QP优化引擎 ─────────────────────────────────────────
slide = light_slide("4.4.1.3 QP优化引擎——再平衡从'拍脑袋'到'数学求解'",
                    "传统再平衡靠经验判断——POMS内置二次规划求解器，500+约束下毫秒级求解最优调仓路径")
# Tech specs
box(slide, Inches(0.5), Inches(1.12), Inches(7.8), Inches(3.45), fill=RGBColor(0x03, 0x1A, 0x52), line=LIGHT_BLUE)
t = tb(slide, Inches(0.6), Inches(1.18), Inches(7.6), Inches(0.35),
       "优化引擎技术规格", size=13, bold=True, color=LIGHT_BLUE)
tf = t.text_frame
for line in [
    "目标函数：Minimize(w₁×跟踪误差 + w₂×交易成本 + w₃×融资成本 + w₄×IFRS9波动)",
    "① 合规约束：单券集中度上限 / 行业限额 / 评级准入 / 净资本占用",
    "② 风险预算：因子敞口上下限 / VaR上限 / DV01 / CS01",
    "③ 流动性约束：单日可变现量 / 市场冲击成本上限",
    "④ IFRS9约束：FVTPL账户波动限制 / OCI账户再分类触发",
    "求解能力：2,000+持仓规模 · 毫秒级求解 · 支持整数规划（面值取整）",
    "多期优化：考虑未来现金流和到期券 · 每日滚动自动更新",
]:
    ap(tf, ("  " if i > 0 and i < 5 else "") + line, size=11,
       color=BRAND_YELLOW if "目标函数" in line else WHITE, space=Pt(4))

# Workflow
box(slide, Inches(8.45), Inches(1.12), Inches(4.4), Inches(3.45), fill=RGBColor(0x03, 0x1A, 0x52), line=BRAND_YELLOW)
t2 = tb(slide, Inches(8.55), Inches(1.18), Inches(4.2), Inches(0.35),
        "再平衡工作流", size=13, bold=True, color=BRAND_YELLOW)
tf2 = t2.text_frame
for line in [
    "触发方式：",
    "① 定时：每日收盘/每周一/月初",
    "② 事件：敞口超限/行情剧变/SAA更新",
    "③ 手动：投资经理主动发起",
    "输出结果：",
    "• 推荐调仓清单（优先级排序）",
    "• 预期效果对比（调仓前后）",
    "• Top5/Top10笔效果分析",
    "• 一键生成交易指令",
]:
    ap(tf2, line, size=11, color=BRAND_YELLOW if line.endswith("：") else WHITE, space=Pt(4))

# Value metrics
for i, (val, label, color) in enumerate([
    ("2-3h → 15min", "再平衡时间（人工→系统）", LIGHT_BLUE),
    ("5-10bp/次", "交易摩擦节省（最优路径）", TEAL),
    ("0次", "合规违规（事前约束）", BRAND_RED),
    ("500+", "同时生效约束条件", BRAND_YELLOW),
]):
    x = Inches(0.5) + i * Inches(3.1)
    box(slide, x, Inches(4.7), Inches(2.9), Inches(0.85), fill=DARK_NAVY, line=color)
    tb(slide, x+Inches(0.1), Inches(4.76), Inches(2.7), Inches(0.4),
       val, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    tb(slide, x+Inches(0.1), Inches(5.15), Inches(2.7), Inches(0.35),
       label, size=10, color=WHITE, align=PP_ALIGN.CENTER)

callout_box(slide, Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.72),
    "核心价值：量化价值明确——再平衡效率提升让跟踪误差减少5-10bps，千亿AUM下每年节省2000-3000万",
    ["合规违规率从事后发现→事前约束归零；交易员决策时间从2-3小时压缩到15分钟确认"])
notes(slide, "4.4.1.3: QP optimization engine. 500+ constraints, millisecond solve. 5-10bp per rebalancing saved.")

# ─── SLIDE 16: 4.4.1.4 跨资产相关性 + Shadow Portfolio ──────────────────────
slide = light_slide("4.4.1.4 跨资产相关性分析 + Shadow Portfolio虚拟试算",
                    "实时相关性矩阵消除'假分散'陷阱 · Shadow Portfolio让调仓'看见结果再行动'")
box(slide, Inches(0.5), Inches(1.12), Inches(6.1), Inches(3.9), fill=RGBColor(0x03, 0x1A, 0x52), line=LIGHT_BLUE)
t = tb(slide, Inches(0.6), Inches(1.18), Inches(5.9), Inches(0.35),
       "跨资产相关性分析", size=13, bold=True, color=LIGHT_BLUE)
tf = t.text_frame
for line in [
    "相关性矩阵：覆盖全持仓（股/债/商品/衍生品/基金）",
    "多时间窗口：20日/60日/252日滚动相关性",
    "历史相关性 vs 压力情景相关性（极端市场中相关性骤升）",
    "动态相关性漂移监控：超过阈值自动预警",
    "相关性热力图：识别'假分散'资产（名义不同但高相关）",
    "分散化比率(Diversification Ratio)：量化分散化程度",
    "低相关资产推荐：约束条件下推荐候选资产",
]:
    ap(tf, "▸ " + line, size=11, color=WHITE, space=Pt(4))

box(slide, Inches(6.75), Inches(1.12), Inches(6.1), Inches(3.9), fill=RGBColor(0x03, 0x1A, 0x52), line=BRAND_YELLOW)
t2 = tb(slide, Inches(6.85), Inches(1.18), Inches(5.9), Inches(0.35),
        "Shadow Portfolio虚拟试算", size=13, bold=True, color=BRAND_YELLOW)
tf2 = t2.text_frame
for line in [
    "100+虚拟组合并行运行，与真实组合同步计算",
    "试算维度：增减仓/换仓/加杠杆/新资产类别",
    "即时显示：风险指标/归因/相关性变化/预期收益",
    "对比视图：试算 vs 当前 所有指标差异",
    "典型问题（秒级给答案）：",
    "  '把5亿利率债换成黄金ETF → Sharpe多少？'",
    "  '叠加100亿国债期货空头 → 久期如何变化？'",
    "  '信用债从20%降到10% → 极端损失减少多少？'",
]:
    ap(tf2, line, size=11, color=BRAND_YELLOW if "典型问题" in line else WHITE, space=Pt(4))

callout_box(slide, Inches(0.5), Inches(5.15), Inches(12.3), Inches(1.12),
    "核心价值：'假分散'陷阱 + 决策可视化",
    ["实时相关性监控：防止'看似分散、实则集中'的配置错误",
     "Shadow Portfolio：调仓决策从'大概猜测'变成'看见结果再行动'——每个调仓决策都有数据依据",
     "投资价值：优化配置潜在收益2000-5000万/年（千亿AUM，Sharpe Ratio提升0.2-0.4）"])
notes(slide, "4.4.1.4: Cross-asset correlation + Shadow Portfolio. Real-time correlation matrix, 100+ virtual portfolios.")

# ─── SLIDE 17: 4.4.1.5 Campisi 2.0 ────────────────────────────────────────
slide = light_slide("4.4.1.5 Campisi 2.0——固定收益绩效归因深度拆解",
                    "债券组合收益'从哪里来'不再是黑箱——五大效应精确拆解让Alpha有据可查")
headers = ["归因效应", "计算公式/逻辑", "实战示例(本月bp)"]
rows = [
    ["① 收入效应", "持有期间票息收入 + 摊销 = '不做任何决策的基础收益'", "+18.5bp（符合预期）"],
    ["② 国债效应-平移", "-Modified Duration × ΔYield(parallel shift)", "-12.3bp（利率上行拖累）"],
    ["② 国债效应-斜率", "各期限桶Key Rate Duration × 斜率变化", "+3.2bp（短端下行收益）"],
    ["② 国债效应-曲率", "凸性调整 + 曲率Key Rate Duration × 曲率变化", "含于斜率效应"],
    ["③ 利差效应-DTS", "Duration Times Spread变化 × 持仓DTS", "+6.8bp（信用利差收窄）"],
    ["③ 行业配置效应", "超/低配某行业相对基准的贡献", "+2.1bp（超配城投正确）"],
    ["③ 个券选择效应", "个券利差变动超出行业平均的部分（真实Alpha）", "+4.7bp（选券Alpha）"],
    ["④ 交易效应", "实际交易时机带来的超额收益或损耗", "-1.2bp（换手时机略有损耗）"],
    ["⑤ 残差", "实际收益 - 上述各项（通常<5bp，过大表明模型需检查）", "+0.0bp（完美匹配）"],
]
add_table(slide, Inches(0.5), Inches(1.15), Inches(12.3), Emu(320000), headers, rows,
          col_widths=[Inches(2.5), Inches(6.5), Inches(3.3)])
callout_box(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.32),
    "核心价值：绩效考核从主观印象变成数据说话",
    ["投资经理第一次清楚知道：'这个月的超额收益，有多少是运气（市场给的收入效应），有多少是真本事（选券Alpha）'",
     "超额收益分解：利率判断能力（平移+斜率效应）+ 信用配置能力（DTS+行业）+ 选券能力（个券Alpha）+ 交易能力",
     "投委会可以用数据决定资源分配——哪个策略/PM值得更多资金？"])
notes(slide, "4.4.1.5: Campisi 2.0 fixed income attribution. 5 effects. Real example with 21.8bp total.")

# ─── SLIDE 18: 4.4.1.6 Brinson ────────────────────────────────────────────
slide = light_slide("4.4.1.6 Brinson模型——权益与基金绩效归因",
                    "超额收益来自'配了什么'还是'选了什么'——Brinson精确拆分配置效应与选择效应")
# Three attribution components
for i, (title, formula, example, color) in enumerate([
    ("第一层：大类资产配置效应",
     "(组合权重 - 基准权重) × (基准资产收益 - 基准总收益)",
     "超配黄金10%，黄金当月+5% vs 基准+2%\n→ 配置贡献：+0.3%",
     LIGHT_BLUE),
    ("第二层：个券/个基选择效应",
     "基准权重 × (组合资产收益 - 基准资产收益)",
     "权益仓位内，选的股票+8% vs 沪深300 +5%\n→ 选择贡献：+0.15%",
     TEAL),
    ("第三层：交互效应",
     "(组合权重-基准权重) × (组合资产收益-基准资产收益)",
     "超配了且选得好的综合效应\n（或超配了但选得差的叠加损耗）",
     BRAND_YELLOW),
]):
    x = Inches(0.5) + i * Inches(4.1)
    box(slide, x, Inches(1.12), Inches(3.9), Inches(2.7), fill=DARK_NAVY, line=color)
    tb(slide, x+Inches(0.1), Inches(1.18), Inches(3.7), Inches(0.4),
       title, size=12, bold=True, color=color)
    tb(slide, x+Inches(0.1), Inches(1.65), Inches(3.7), Inches(0.5),
       "公式：" + formula, size=10, color=MED_GRAY)
    box(slide, x+Inches(0.3), Inches(2.22), Inches(3.3), Inches(0.03), fill=color)
    tb(slide, x+Inches(0.1), Inches(2.32), Inches(3.7), Inches(0.45),
       example, size=10, color=WHITE)

# Penetration capabilities
box(slide, Inches(0.5), Inches(3.95), Inches(12.3), Inches(1.3), fill=RGBColor(0x03, 0x1A, 0x52), line=TEAL)
t = tb(slide, Inches(0.6), Inches(4.02), Inches(12.1), Inches(0.35),
       "穿透分析能力 + 投资能力量化评估", size=13, bold=True, color=TEAL)
tf = t.text_frame
for line in [
    "多层穿透：大类→行业→子行业→个股/个基 | 基金FOF穿透：穿透底层持仓看真实行业/风格暴露",
    "季度报告：配置能力 +Xbp | 选券能力 +Ybp | 择时能力 +Zbp | 连续12月滚动识别稳定性vs运气",
    "投资经理排名（同部门横向对比，激励机制数据支撑）",
]:
    ap(tf, "▸ " + line, size=11, color=WHITE, space=Pt(3))

callout_box(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.88),
    "核心价值：基金经理/投资经理的真实投资能力第一次可以被量化",
    ["绩效考核不再靠'感觉'——投委会用数据决定资源分配，哪个PM值得更多风险预算"])
notes(slide, "4.4.1.6: Brinson attribution for equity/fund. 3-layer decomposition. Fund-of-funds penetration.")

# ─── SLIDE 19: 4.4.1.7 MCTR + 因子风险分解 ────────────────────────────────
slide = light_slide("4.4.1.7 因子风险分解 + 边际风险贡献MCTR",
                    "从'知道组合VaR是多少'升级到'知道每一笔持仓消耗了多少风险预算'")
box(slide, Inches(0.5), Inches(1.12), Inches(6.1), Inches(4.0), fill=RGBColor(0x03, 0x1A, 0x52), line=LIGHT_BLUE)
t = tb(slide, Inches(0.6), Inches(1.18), Inches(5.9), Inches(0.35),
       "因子风险分解（全组合维度）", size=13, bold=True, color=LIGHT_BLUE)
tf = t.text_frame
for line in [
    "① 利率风险：DV01敞口 × 利率波动率 → 贡献XX%",
    "   细分：短端(1Y内)/中端(1-5Y)/长端(5Y+)/超长端(10Y+)",
    "② 信用风险：CS01敞口 × 信用利差波动 → 贡献XX%",
    "   细分：AAA/AA+/AA城投/产业/金融各评级",
    "③ 权益风险：Beta敞口 × 股市波动 → 贡献XX%",
    "   细分：大盘/小盘/价值/成长风格因子",
    "④ 商品风险：黄金/原油/有色价格敏感度 → 贡献XX%",
    "特异性风险：个券/个股集中度风险（过高→需分散）",
    "输出：风险饼图 · 预算仪表盘（已用/剩余/预警线）",
]:
    ap(tf, ("   " if line.startswith("   ") else "▸ ") + line.strip(), size=10,
       color=MED_GRAY if line.startswith("   ") else WHITE, space=Pt(3))

box(slide, Inches(6.75), Inches(1.12), Inches(6.1), Inches(4.0), fill=RGBColor(0x03, 0x1A, 0x52), line=BRAND_YELLOW)
t2 = tb(slide, Inches(6.85), Inches(1.18), Inches(5.9), Inches(0.35),
        "边际风险贡献 MCTR", size=13, bold=True, color=BRAND_YELLOW)
tf2 = t2.text_frame
for line in [
    "MCTR_i = ∂(组合波动率) / ∂(持仓_i的权重)",
    "即：增加1%该持仓，组合总风险增加多少",
    "应用价值：",
    "① 风险消耗排名：找出哪10个持仓消耗了60%风险预算",
    "② 减仓优先级：需降风险时，先减MCTR最高的持仓",
    "③ 加仓引导：优先选MCTR为负的品种（降低组合风险）",
    "④ QP优化输入：作为约束条件实现风险预算最优分配",
    "实时仪表盘示例：",
    "  #1 XX城投2030：MCTR 18.3bp → 12.1% ⚠️",
    "  #2 XX转债：MCTR 12.7bp → 8.4% 正常",
    "  #3 XX股票：MCTR -3.2bp → -2.1% ✅分散化",
]:
    ap(tf2, line, size=10,
       color=BRAND_YELLOW if "MCTR_i" in line or "应用价值" in line or "实时仪表盘" in line else WHITE,
       space=Pt(3))

callout_box(slide, Inches(0.5), Inches(5.25), Inches(12.3), Inches(1.02),
    "核心价值：风险预算分配从粗放走向精细",
    ["风险管理升级：'知道组合VaR' → '知道每笔持仓消耗了多少风险预算'",
     "操作价值：需要降风险时，精确知道减哪些持仓效率最大；加仓时知道哪些品种不消耗额外风险预算"])
notes(slide, "4.4.1.7: Factor risk decomposition + MCTR. From knowing total VaR to knowing each position's risk budget consumption.")

# ─── SLIDE 20: 4.4.1.8 Component VaR ──────────────────────────────────────
slide = light_slide("4.4.1.8 Component VaR + 主动风险归因",
                    "从'组合整体有多大风险'到'每一个持仓/行业/发行人对VaR的贡献是多少'")
headers2 = ["层级", "穿透维度", "分析内容", "管理价值"]
rows2 = [
    ["第一层", "资产类别", "债券VaR/权益VaR/商品VaR + 各类之间分散化节省", "了解风险来自哪个大类"],
    ["第二层", "行业/子市场", "城投/金融债/产业债/TMT/消费各行业VaR", "识别行业集中度风险"],
    ["第三层", "评级/信用质量", "AAA/AA+/AA各评级VaR占比", "信用风险精细管控"],
    ["第四层", "单一发行人", "发行人A: Component VaR → 超集中度预警", "单一发行人风险阻止"],
]
add_table(slide, Inches(0.5), Inches(1.15), Inches(12.3), Emu(350000), headers2, rows2,
          col_widths=[Inches(1.2), Inches(1.8), Inches(5.3), Inches(4.0)])

# Three VaR methods
for i, (method, desc, color) in enumerate([
    ("历史模拟法", "500+日历史数据\n99%置信度\n线性风险精确", LIGHT_BLUE),
    ("Monte Carlo", "10,000次模拟\n捕获非线性风险\n期权/结构化产品", BRAND_YELLOW),
    ("参数法", "解析解\n毫秒级输出\n实时监控适用", TEAL),
]):
    x = Inches(0.5) + i * Inches(4.1)
    box(slide, x, Inches(4.0), Inches(3.9), Inches(1.35), fill=DARK_NAVY, line=color)
    tb(slide, x+Inches(0.1), Inches(4.07), Inches(3.7), Inches(0.35),
       method, size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    tb(slide, x+Inches(0.1), Inches(4.5), Inches(3.7), Inches(0.8),
       desc, size=11, color=WHITE, align=PP_ALIGN.CENTER)

tb(slide, Inches(0.5), Inches(5.42), Inches(12.3), Inches(0.3),
   "三种方法互为验证，差异超过阈值自动预警；跟踪误差分解：配置效应TE + 选券效应TE，信息比率(IR)持续监控",
   size=10, color=MED_GRAY)
callout_box(slide, Inches(0.5), Inches(5.78), Inches(12.3), Inches(0.57),
    "核心价值：风险管理精细化到可操作层面——每个发行人的Component VaR超限即预警，立刻知道减哪个",
    [])
notes(slide, "4.4.1.8: Component VaR. 4-layer penetration. 3 VaR methods cross-validate.")

# ─── SLIDE 21: 4.4.1.9 敏感度分析 ─────────────────────────────────────────
slide = light_slide("4.4.1.9 敏感度分析——实时风险地图，预判市场波动影响",
                    "市场利率明天上行10bp，组合会亏多少？——敏感度分析实时回答，让投资经理提前看清风险地图")
# Left: sensitivity metrics
box(slide, Inches(0.5), Inches(1.12), Inches(6.1), Inches(4.1), fill=RGBColor(0x03, 0x1A, 0x52), line=LIGHT_BLUE)
t = tb(slide, Inches(0.6), Inches(1.18), Inches(5.9), Inches(0.35),
       "多维度敏感度实时仪表盘", size=13, bold=True, color=LIGHT_BLUE)
tf = t.text_frame
for line in [
    "利率敏感度：",
    "  DV01（利率每变动1bp → P&L变化）",
    "  Key Rate Duration（各期限点敏感度）",
    "  凸性（大幅变动时的非线性保护）",
    "信用敏感度：",
    "  CS01（信用利差每变动1bp → P&L变化）",
    "  DTS（Duration Times Spread，更精确）",
    "权益敏感度：Delta · Beta · 行业Beta",
    "商品/汇率：黄金/原油/外汇每单位变动影响",
]:
    is_header = line.endswith("：")
    ap(tf, line, size=11, color=BRAND_YELLOW if is_header else WHITE, space=Pt(3))

# Right: shock grid
box(slide, Inches(6.75), Inches(1.12), Inches(6.1), Inches(4.1), fill=RGBColor(0x03, 0x1A, 0x52), line=BRAND_RED)
t2 = tb(slide, Inches(6.85), Inches(1.18), Inches(5.9), Inches(0.35),
        "冲击分析矩阵（利率 × 信用利差）", size=13, bold=True, color=BRAND_RED)
tf2 = t2.text_frame
shock_grid = [
    "              信用 -50bp  -20bp  0bp  +20bp  +50bp",
    "利率 -50bp    +8500  +6200 +4100  +1800   -800",
    "利率 -20bp    +5200  +3100 +1200   -900  -3200",
    "利率   0bp    +4100  +1800     0  -2100  -5400",
    "利率 +20bp    +2800   +400 -1800  -4100  -7200",
    "利率 +50bp     +800  -1800 -4200  -6800  -9800",
    "（万元P&L，红色区域为当前情景最可能损失区间）",
]
for line in shock_grid:
    ap(tf2, line, size=9, color=BRAND_RED if "红色" in line else WHITE,
       font="Courier New", space=Pt(3))

callout_box(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.0),
    "核心价值：从'被动应对'到'主动预判'",
    ["提前知道每种市场情景的P&L影响，决策从'市场动了才反应'变为'看懂矩阵，提前布防'",
     "冲击矩阵可视化风险地图：投资经理一眼看出'最怕利率大幅上行+信用利差走阔的组合情景'"])
notes(slide, "4.4.1.9: Sensitivity analysis. Real-time DV01/CS01/Delta dashboard. Shock grid for scenario planning.")

# ─── SLIDE 22: 4.4.1.10 情景分析+压力测试 ─────────────────────────────────
slide = light_slide("4.4.1.10 情景分析与压力测试——三层体系让924类事件从'灾难'变为'有预案'",
                    "2024年924事件证明：没有极端情景预案 = 被动承受——POMS三层压力测试体系")
for i, (layer, title, content, color) in enumerate([
    ("第一层", "历史情景重演",
     ["924央行货币政策转向（2024.9）：短端-30bp/中端-20bp/权益+9.6%",
      "包商银行事件（2019.5）：城投AA+走阔+80bp/AA+走阔+150bp",
      "2013钱荒：隔夜Shibor突破30%，长端利率+100bp",
      "→ 自动计算当前组合在上述事件下的损益"],
     BRAND_RED),
    ("第二层", "假设情景模拟",
     ["自定义任意冲击：'利率+100bp + 信用+50bp + 股市-10%'",
      "'某省城投违约，行业利差走阔200bp'",
      "即时输出：总损益/最受损前10持仓/流动性危机能力",
      "→ 建议应急对冲方案（哪些衍生品可以对冲）"],
     BRAND_YELLOW),
    ("第三层", "反向压力测试",
     ["从'最大可承受损失'倒推'什么情景触及止损线'",
      "算法：搜索最可能导致-2亿损失的利率/信用组合",
      "输出：'临界情景是利率+65bp+信用+30bp，历史每3年1次'",
      "→ 提前建立应急预案和止损机制"],
     LIGHT_BLUE),
]):
    x = Inches(0.5) + i * Inches(4.1)
    box(slide, x, Inches(1.12), Inches(3.9), Inches(4.05), fill=DARK_NAVY, line=color)
    tb(slide, x+Inches(0.1), Inches(1.18), Inches(3.7), Inches(0.3),
       layer, size=11, bold=True, color=color)
    tb(slide, x+Inches(0.1), Inches(1.52), Inches(3.7), Inches(0.38),
       title, size=14, bold=True, color=WHITE)
    box(slide, x+Inches(0.3), Inches(1.96), Inches(3.3), Inches(0.03), fill=color)
    tf = tb(slide, x+Inches(0.1), Inches(2.07), Inches(3.7), Inches(3.0), "").text_frame
    tf.word_wrap = True
    for line in content:
        ap(tf, "▸ " + line, size=10, color=WHITE, space=Pt(5))

callout_box(slide, Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.95),
    "924教训：提前压测 = 有预案 · 没有压测 = 被动承受",
    ["POMS将极端情景从'事后复盘'变成'事前演练'：同样的924事件，有预案亏0.8亿，无预案亏2亿",
     "差异1.2亿——这是压力测试系统在一天内创造的价值"])
notes(slide, "4.4.1.10: Scenario analysis + stress testing. 3 layers. 924 event: 1.2B difference with/without POMS.")

# ─── SLIDE 23: 4.4.1.11 流动性风险管理 ────────────────────────────────────
slide = light_slide("4.4.1.11 流动性风险管理——变现确定性量化",
                    "大额持仓'卖不出去'的流动性风险比价格风险更致命——POMS量化每个持仓的变现时间和冲击成本")
headers3 = ["持仓", "规模", "日均成交量", "估算变现天数", "变现难度"]
rows3 = [
    ["国债2034", "10亿", "50亿/日", "1天", "✅ 极高流动性"],
    ["AA城投2028", "5亿", "3亿/日", "8天", "⚠️ 中等流动性"],
    ["AA-产业债2030", "2亿", "0.5亿/日", "20天+", "🔴 低流动性"],
    ["非标ABS", "3亿", "无场内交易", "协议转让", "🔴 极低流动性"],
]
tb(slide, Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.32),
   "持仓变现能力评估（基于Almgren-Chriss模型 — 市场冲击成本 ∝ √变现规模/日均成交量）",
   size=12, bold=True, color=DARK_NAVY)
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Emu(360000), headers3, rows3,
          col_widths=[Inches(2.5), Inches(1.5), Inches(2.3), Inches(2.3), Inches(3.7)])

box(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(1.55), fill=RGBColor(0x03, 0x1A, 0x52), line=TEAL)
t = tb(slide, Inches(0.6), Inches(3.66), Inches(12.1), Inches(0.35),
       "流动性压力测试 + 流动性监控仪表盘", size=13, bold=True, color=TEAL)
tf = t.text_frame
for line in [
    "场景：需要在X天内变现Y亿资金（满足回购到期/赎回）→ 系统分析：可变现排序/变现成本累计/最优路径建议",
    "流动性调整VaR(LVaR) = 传统VaR + 流动性溢价风险  → 比传统VaR更保守，更接近真实风险",
    "仪表盘：流动性覆盖率 · 1/3/5/10天可变现缓冲 · 回购到期日历（未来30日每日到期规模）",
]:
    ap(tf, "▸ " + line, size=11, color=WHITE, space=Pt(4))

callout_box(slide, Inches(0.5), Inches(5.32), Inches(12.3), Inches(1.0),
    "核心价值：在极端市场中，'能卖多少·能卖多快·要付多少代价'有确定性答案",
    ["924事件教训：价格跌是确定的，卖不出去才是真正的灾难",
     "潜在价值：避免流动性紧急变现的额外损失，每次极端事件可减少数千万到数亿的冲击成本"])
notes(slide, "4.4.1.11: Liquidity risk management. Almgren-Chriss model. LVaR. Liquidation timeline dashboard.")

# ─── SLIDE 24: 4.4.1.12 衍生品工具矩阵 ────────────────────────────────────
slide = light_slide("4.4.1.12 衍生品工具矩阵——快速/灵活/稳健的组合对冲体系",
                    "系统告诉投资经理：用哪个工具、用多少、怎么做——将对冲从艺术变成科学")
headers4 = ["风险类型", "对冲工具", "用途", "效率"]
rows4 = [
    ["利率风险", "国债期货（2Y/5Y/10Y/30Y）", "快速调整组合久期", "秒级，高流动性"],
    ["利率风险", "利率互换IRS（1Y-10Y）", "精确管理浮动/固定利率敞口", "灵活，OTC"],
    ["信用风险", "信用违约互换CDS", "对冲特定发行人违约风险", "精准，OTC"],
    ["权益风险", "股指期货（IF/IH/IC/IM）", "快速对冲Beta敞口", "秒级，高流动性"],
    ["权益风险", "股指期权", "非线性保护（尾部风险对冲）", "精准，成本可控"],
    ["商品风险", "黄金期货/黄金期权", "对冲黄金价格风险", "日内调整"],
    ["汇率风险", "外汇远期/掉期", "锁定外汇敞口", "OTC，灵活"],
]
add_table(slide, Inches(0.5), Inches(1.15), Inches(12.3), Emu(330000), headers4, rows4,
          col_widths=[Inches(2.0), Inches(3.3), Inches(4.5), Inches(2.5)])

box(slide, Inches(0.5), Inches(4.05), Inches(12.3), Inches(1.55), fill=RGBColor(0x03, 0x1A, 0x52), line=BRAND_YELLOW)
t = tb(slide, Inches(0.6), Inches(4.12), Inches(12.1), Inches(0.35),
       "对冲决策引擎工作流", size=13, bold=True, color=BRAND_YELLOW)
tf = t.text_frame
for line in [
    "① 输入：当前组合风险敞口（DV01/CS01/Beta实时仪表盘数据）",
    "② 目标：设定对冲后目标风险水平（如'将DV01从500万降到200万'）",
    "③ 系统推荐：卖出XX手10年期国债期货 + 做多XX笔5年IRS → 预期DV01: 500万→198万(-60.4%) · 对冲成本: XX万",
    "④ 一键生成对冲指令，流转执行管理引擎  |  多维场景：快速对冲/精细对冲/尾部对冲/跨资产联合对冲",
]:
    ap(tf, "▸ " + line, size=11, color=WHITE, space=Pt(4))

callout_box(slide, Inches(0.5), Inches(5.78), Inches(12.3), Inches(0.57),
    "核心价值：衍生品对冲从'知道要用但不知道用多少'变成'系统给出最优方案，一键执行'",
    [])
notes(slide, "4.4.1.12: Derivatives hedging matrix. Decision engine. 7 tool types. One-click hedge generation.")

# ─── SLIDE 25: ENGINE 2 高性能计算 ─────────────────────────────────────────
slide = light_slide("引擎②：高性能计算引擎——算得快 [算得快]",
                    "实时估值/定价/风险计量 · 自研替代Bloomberg · 千亿规模2000万次定价/日")
for i, (title, items, color) in enumerate([
    ("定价引擎", ["覆盖权益/利率/商品/外汇/衍生品/结构化产品", "利率曲线/波动率曲面/MC多方法支撑", "Full Repricing：非线性风险精确捕获", "自研替代Bloomberg，消除许可费"], LIGHT_BLUE),
    ("实时估值引擎", ["行情驱动流式估值，Tick级刷新", "2,000持仓 × 每秒 = 2000次定价/秒", "VaR/CVaR/Greeks实时计算", "组合估值精度到Greeks级别"], TEAL),
    ("分布式计算架构", ["K8s编排弹性扩缩容", "8-16计算节点，VaR=2000万次定价/日", "CEP引擎：10万+ events/sec", "二次开发SDK（Python/Java/C++）"], BRAND_YELLOW),
]):
    x = Inches(0.5) + i * Inches(4.1)
    box(slide, x, Inches(1.12), Inches(3.9), Inches(4.3), fill=DARK_NAVY, line=color)
    tb(slide, x+Inches(0.1), Inches(1.18), Inches(3.7), Inches(0.42),
       title, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    box(slide, x+Inches(0.5), Inches(1.65), Inches(2.9), Inches(0.03), fill=color)
    tf = tb(slide, x+Inches(0.1), Inches(1.78), Inches(3.7), Inches(3.5), "").text_frame
    tf.word_wrap = True
    for item in items:
        ap(tf, "▸ " + item, size=12, color=WHITE, space=Pt(8))

callout_box(slide, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.82),
    "规模验证（千亿AUM估算）",
    ["存储：~10TB（3年滚动 含行情+持仓+交易+风控）· 并发：40用户 × 5请求/秒 = 200 req/sec",
     "Bloomberg依赖归零：许可费+数据费每年数百万，完全替代 → 300-500万/年成本节省"])
notes(slide, "Engine 2: High-performance computing. Self-developed pricing to replace Bloomberg. 20M pricings/day.")

# ─── SLIDE 26: ENGINE 3 量化策略 ────────────────────────────────────────────
slide = light_slide("引擎③：量化策略引擎——算得快+配得优 [相对价值·多策略·回测]",
                    "中国FICC市场系统性相对价值机会丰富——华锐量化工具让机会从'看到但抓不住'变为'系统化执行'")
for i, (title, items, color) in enumerate([
    ("相对价值分析工具集", [
        "利差策略：新老券/品种内/跨品种利差实时监控",
        "基差策略：国债期货基差/CTD券跟踪/期现套利",
        "曲线策略：收益率曲线形态/斜率/蝴蝶策略",
        "历史分位数：当前利差在历史X%分位，告警",
        "胜率回测：类似历史位置的胜率和平均收益"], LIGHT_BLUE),
    ("多timeframe信号引擎", [
        "日内信号：利差偏离>2σ → 自动推送机会提示",
        "日间信号：趋势跟踪+均值回归策略",
        "周/月级：宏观因子驱动的资产配置信号",
        "多策略并行：不同周期信号同时运行",
        "实时回测：新信号上线前历史数据验证"], TEAL),
    ("策略回测与验证框架", [
        "历史数据：5年+债券/权益/商品数据库",
        "快速回测：假设→验证结果<1天（vs 原1-2周）",
        "绩效指标：年化收益/夏普/最大回撤/胜率",
        "模拟盘验证：回测通过后模拟盘实盘验证",
        "AI引擎接口：机器学习策略预留接口"], BRAND_YELLOW),
]):
    x = Inches(0.5) + i * Inches(4.1)
    box(slide, x, Inches(1.12), Inches(3.9), Inches(4.3), fill=DARK_NAVY, line=color)
    tb(slide, x+Inches(0.1), Inches(1.18), Inches(3.7), Inches(0.42),
       title, size=13, bold=True, color=color)
    box(slide, x+Inches(0.3), Inches(1.65), Inches(3.3), Inches(0.03), fill=color)
    tf = tb(slide, x+Inches(0.1), Inches(1.78), Inches(3.7), Inches(3.5), "").text_frame
    tf.word_wrap = True
    for item in items:
        ap(tf, "▸ " + item, size=11, color=WHITE, space=Pt(6))

callout_box(slide, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.82),
    "量化价值（千亿AUM）：相对价值Alpha年化+20-50bps → 新增1000-2000万/年",
    ["中国FICC利差市场系统性机会多，但需要工具才能系统化捕捉——有工具 vs 无工具 = 1000-2000万/年的Alpha差距"])
notes(slide, "Engine 3: Quant strategy. Relative value tools, multi-timeframe signals, backtesting. 1000-2000M/yr alpha.")

# ─── SLIDE 27: ENGINE 4 风控 ────────────────────────────────────────────────
slide = light_slide("引擎④：事件驱动风控引擎——控得稳 [CEP·回撤·极端压测·合规]",
                    "从'被动等报告'到'主动推送预警'——CEP引擎实现秒级风险响应")
# Key capabilities
cols_data = [
    ("CEP复杂事件处理", ["行情/交易/越限/流动性事件统一处理", "10万+ events/sec事件吞吐", "毫秒级规则匹配和预警触发", "多级预警：黄色预警/红色告警/自动止损"], BRAND_RED),
    ("实时回撤监控", ["秒级检测组合回撤（非T+1）", "分层阈值：预警线/止损线/强制减仓", "自动触发减仓建议（结合流动性排序）", "回撤归因：是利率/信用/权益哪个因子"], BRAND_YELLOW),
    ("合规规则引擎", ["5000+规则库（净资本/集中度/评级/关联交易）", "投前预检：下单前自动合规校验", "投中拦截：违规交易实时阻止", "2024券商受处罚162次 → 预检归零"], LIGHT_BLUE),
    ("全公司风险视图", ["全公司VaR/DV01/CS01/集中度一屏", "跨组合/跨账簿/跨部门统一视图", "实时限额利用率（绿/黄/红交通灯）", "压力测试结果实时集成"], TEAL),
]
for i, (title, items, color) in enumerate(cols_data):
    x = Inches(0.5) + i * Inches(3.1)
    box(slide, x, Inches(1.12), Inches(2.9), Inches(4.3), fill=DARK_NAVY, line=color)
    tb(slide, x+Inches(0.1), Inches(1.18), Inches(2.7), Inches(0.42),
       title, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
    box(slide, x+Inches(0.3), Inches(1.65), Inches(2.3), Inches(0.03), fill=color)
    tf = tb(slide, x+Inches(0.1), Inches(1.78), Inches(2.7), Inches(3.5), "").text_frame
    tf.word_wrap = True
    for item in items:
        ap(tf, "▸ " + item, size=10, color=WHITE, space=Pt(5))

callout_box(slide, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.82),
    "风险控制价值（千亿AUM）",
    ["常规回撤控制：最大回撤从-3%控制到-1.5% → 减少3000-5000万/次",
     "极端事件预案（924类）：有预案 vs 无预案 → 单次差异1-2亿 · 合规拦截：每年避免2-3次违规处罚"])
notes(slide, "Engine 4: Event-driven risk control. CEP 100K+ events/sec. 5000+ compliance rules. Pre-trade checking.")

# ─── SLIDE 28: ENGINE 5 指令执行 ────────────────────────────────────────────
slide = light_slide("引擎⑤：指令与执行管理引擎——连得通 [决策→指令→执行→归因]",
                    "投研发现机会→组合决策→指令生成→OEMS执行→TCA分析→归因驱动下次决策——全链路贯通")
# Flow diagram
steps = [
    ("投研信号", "利差机会推送\n宏观因子信号", TEAL),
    ("组合决策", "再平衡优化\nShadow验证", LIGHT_BLUE),
    ("合规预检", "5000+规则\n事前拦截", BRAND_YELLOW),
    ("指令生成", "multi-leg指令\n批量审批", BRAND_RED),
    ("OEMS执行", "场内/场外/期货\n算法交易", TEAL),
    ("TCA归因", "滑点分析\n绩效归因", LIGHT_BLUE),
]
for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.5) + i * Inches(2.1)
    box(slide, x, Inches(1.12), Inches(1.9), Inches(2.0), fill=DARK_NAVY, line=color)
    tb(slide, x+Inches(0.05), Inches(1.18), Inches(1.8), Inches(0.45),
       title, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
    tb(slide, x+Inches(0.05), Inches(1.68), Inches(1.8), Inches(0.8),
       desc, size=10, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 5:
        tb(slide, x+Inches(1.9), Inches(1.55), Inches(0.3), Inches(0.4),
           "→", size=16, bold=True, color=color, align=PP_ALIGN.CENTER)

# Value table
headers5 = ["功能", "当前状态", "POMS能力", "量化价值"]
rows5 = [
    ["指令生成", "手工填写指令单，1-2小时", "再平衡自动生成，15分钟确认", "效率提升8-10x"],
    ["multi-leg执行", "手工分腿，各自执行，价差损失", "系统同时多腿，原子执行", "滑点减少50%+"],
    ["合规预检", "事后发现违规，处罚风险", "投前自动校验，实时拦截", "违规归零"],
    ["执行质量TCA", "无从评估滑点来源", "vs TWAP/VWAP持续优化", "500-1000万/年"],
    ["绩效归因闭环", "年底才看结果，不驱动改进", "实时归因→下次决策输入", "Alpha持续迭代"],
]
add_table(slide, Inches(0.5), Inches(3.35), Inches(12.3), Emu(330000), headers5, rows5,
          col_widths=[Inches(2.0), Inches(3.0), Inches(3.3), Inches(4.0)])

callout_box(slide, Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.7),
    "年度交易量3000亿：执行质量每提升1bp = 3000万/年 · 滑点优化0.3-0.5bp → 节省1000-1500万/年",
    [])
notes(slide, "Engine 5: Order and execution management. Full lifecycle from research signal to execution to attribution.")

# ─── DIVIDER: SECTION 4 ─────────────────────────────────────────────────────
divider_slide("SECTION 4 — IMPLEMENTATION",
              "三期18个月：Phase 1六个月CEO看到全貌",
              "从底向上·分步见效·灯塔效应")

# ─── SLIDE 29: ROADMAP ──────────────────────────────────────────────────────
slide = light_slide("三期18个月实施路径——从数据筑基到全面进化",
                    "先数据后应用 · 先看见后计算 · 先自营后做市 · Phase 1六个月灯塔效应")
for i, (phase, months, tag, delivs, milestone, value_kpi, color) in enumerate([
    ("Phase 1", "0-6月", "筑基期：看得清+配得优",
     ["POMS数据集成（对接IBOR+衡泰xIR+柜台+行情）",
      "全资产组合管理引擎MVP（全品种一屏）",
      "实时持仓敞口分析（公司→交易台→个人穿透）",
      "跨资产相关性分析（低波动最优组合）",
      "定价平台基础功能 · 基础风控+回撤预警"],
     "千亿自营资金在一个平台实时可见",
     "CEO打开系统即看全公司自营全貌", TEAL),
    ("Phase 2", "6-12月", "赋能期：算得快+控得稳",
     ["高性能计算引擎（实时风险计量+估值全量）",
      "量化策略工具（利差/基差/新老券分析）",
      "事件驱动风控引擎（CEP+极端压力测试）",
      "虚拟组合what-if试算（100+并行）",
      "绩效归因（Campisi+Brinson）· 指令管理"],
     "从T+1到实时风控 · 策略驱动投资决策",
     "利差机会秒级捕捉 · 回撤实时预警", LIGHT_BLUE),
    ("Phase 3", "12-18月", "进化期：连得通+领得先",
     ["多策略多timeframe回测框架",
      "multi-leg交易执行优化+滑点归因",
      "四向集成全面打通（投研←→POMS←→OEMS）",
      "现券交易系统升级+AI引擎",
      "FOF管理 · 做市能力储备评估"],
     "完整POMS平台上线，对标行业领先",
     "全链路打通 · 向做市/客需延伸基础", BRAND_YELLOW),
]):
    x = Inches(0.5) + i * Inches(4.1)
    box(slide, x, Inches(1.1), Inches(3.9), Inches(5.0), fill=DARK_NAVY, line=color)
    box(slide, x, Inches(1.1), Inches(3.9), Inches(0.55), fill=color)
    tb(slide, x+Inches(0.1), Inches(1.12), Inches(3.7), Inches(0.45),
       f"{phase}（{months}）", size=14, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
    tb(slide, x+Inches(0.1), Inches(1.72), Inches(3.7), Inches(0.38),
       tag, size=11, bold=True, color=color)
    tf = tb(slide, x+Inches(0.1), Inches(2.15), Inches(3.7), Inches(2.5), "").text_frame
    tf.word_wrap = True
    for d in delivs:
        ap(tf, "▸ " + d, size=10, color=WHITE, space=Pt(4))
    box(slide, x+Inches(0.1), Inches(4.75), Inches(3.7), Inches(0.3), fill=RGBColor(0x05, 0x12, 0x4A))
    tb(slide, x+Inches(0.15), Inches(4.78), Inches(3.6), Inches(0.25),
       "里程碑：" + milestone, size=9, color=color)
    tb(slide, x+Inches(0.15), Inches(5.08), Inches(3.6), Inches(0.25),
       "验证：" + value_kpi, size=9, color=WHITE)
notes(slide, "3-phase 18-month roadmap. Phase 1: data integration + portfolio visibility in 6 months.")

# ─── SLIDE 30: PHASE 1 DAY IN THE LIFE ─────────────────────────────────────
slide = light_slide("Phase 1灯塔效应：投资经理的一天——从'看不到'到'全资产精准配置'",
                    "6个月内实现：千亿自营资金一个平台实时可见——这是最直观的变化")
headers6 = ["时间", "Without POMS（今天）", "With POMS（Phase 1上线后）"]
rows6 = [
    ["08:30", "打开3个系统查看昨日持仓，看不到全资产机会", "登录POMS，全组合实时全景一屏 · 相关性矩阵：黄金与利率债相关性仅0.1，建议增配5%"],
    ["09:15", "等风控发邮件——昨日VaR报告（数据过时12小时）", "实时仪表盘：VaR 1.2亿 · DV01 850万 · CS01 420万 · 黄色预警：利率债集中度85%"],
    ["10:00", "想试增配黄金ETF，Excel粗算，不确定就放弃", "虚拟组合一键试算：增配5%黄金ETF → 波动率-32% · 预期收益仅-3% · Sharpe 0.8→1.1"],
    ["11:30", "看到10Y-5Y利差走阔，缺工具，无法量化决策", "利差监控推送：'10Y-5Y利差52bp，历史93%分位' · 回测：类似位置胜率78%，均收益12bp"],
    ["14:00", "手写指令单→签字→交易员执行，价格已变", "一键生成指令：合规预检通过→领导线上审批→秒级下达→滑点仅0.3bp"],
    ["15:30", "无法知道今天操作对组合的真实贡献", "实时归因：利差策略+8bp · 久期管理+3bp · 交易滑点-0.3bp"],
    ["17:00", "若全年仅做利率债，收益取决于单一方向判断", "系统实时：本月多资产配置 年化收益5.8% · 波动率3.1% · Sharpe 1.45"],
]
add_table(slide, Inches(0.5), Inches(1.15), Inches(12.3), Emu(315000), headers6, rows6,
          col_widths=[Inches(0.85), Inches(4.55), Inches(6.9)])
callout_box(slide, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.75),
    "Phase 1交付承诺：6个月内，CEO打开系统即看全公司自营实时全貌——这是最快的价值验证",
    ["Phase 1本身就是对华锐交付能力的验证：6个月先看见、先感受，后续深度功能基于信任分步建设"])
notes(slide, "Phase 1 lighthouse effect. Day-in-the-life comparison shows concrete daily value to investment managers.")

# ─── SLIDE 31: RISK REGISTER ────────────────────────────────────────────────
slide = light_slide("风险识别与应对：主要风险及缓解措施",
                    "提前识别、主动应对——分步建设策略降低整体风险")
headers7 = ["风险", "概率", "影响", "应对措施", "责任方"]
rows7 = [
    ["衡泰xIR数据对接复杂", "高", "高", "Phase 1提前启动接口调研；已有IBOR对接经验；渐进替换策略", "双方"],
    ["千亿数据量性能瓶颈", "中", "高", "Phase 1含性能压测专项；分布式架构预留弹性扩展", "华锐"],
    ["业务需求变更频繁", "中", "中", "敏捷迭代（双周）+需求冻结期+变更委员会机制", "双方"],
    ["关键投资经理参与不足", "中", "高", "项目发起人保障业务代表参与；专人对接关键用户", "国元"],
    ["数据质量问题影响准确性", "高", "高", "Phase 1含数据治理专项；入库自动校验规则", "双方"],
    ["信创环境适配挑战", "中", "中", "华锐已有信创全栈验证经验；提前环境验证专项", "华锐"],
    ["IFRS9合规时间压力", "中", "高", "Phase 1优先覆盖IFRS9基础功能；按监管时间表倒排", "双方"],
]
add_table(slide, Inches(0.5), Inches(1.15), Inches(12.3), Emu(345000), headers7, rows7,
          col_widths=[Inches(2.8), Inches(0.8), Inches(0.8), Inches(6.2), Inches(1.7)])
callout_box(slide, Inches(0.5), Inches(5.62), Inches(12.3), Inches(0.75),
    "风险管理原则：先建数据根基，分步见效，Phase 1即验证华锐交付能力",
    ["分步建设策略：每个Phase单独验证，结果满意后继续——最大化投资保护"])
notes(slide, "Risk register. 7 key risks with mitigation. Data integration and performance are the highest priority.")

# ─── SLIDE 32: GOVERNANCE ───────────────────────────────────────────────────
slide = light_slide("项目治理与关键里程碑",
                    "双周迭代 · 月度管理层汇报 · 季度Steering Committee · 结果驱动")
# Governance table
headers8 = ["角色", "国元侧", "华锐侧", "职责"]
rows8 = [
    ["项目发起人", "分管副总裁", "华锐VP", "战略决策·资源保障·阶段验收"],
    ["项目经理", "IT部项目经理", "华锐项目总监", "日常协调·进度追踪·风险上报"],
    ["业务代表", "自营部投资经理+交易员", "华锐业务顾问", "需求确认·验收测试·用户培训"],
    ["技术代表", "IT部架构师", "华锐技术架构师", "技术方案·集成设计·性能验证"],
    ["风控代表", "风控合规部", "华锐风控顾问", "合规规则·风控指标·监管对接"],
]
add_table(slide, Inches(0.5), Inches(1.15), Inches(12.3), Emu(370000), headers8, rows8,
          col_widths=[Inches(2.0), Inches(3.0), Inches(2.5), Inches(4.8)])

# Milestone timeline
tb(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.35),
   "关键里程碑时间线", size=13, bold=True, color=DARK_NAVY)
milestones = [
    ("M1", "项目启动\n数据调研"),
    ("M3", "IBOR对接\n数据集成"),
    ("M6", "Phase 1\n上线验收"),
    ("M9", "实时风控\n策略工具"),
    ("M12", "Phase 2\n上线验收"),
    ("M15", "全链路\n打通"),
    ("M18", "Phase 3\n完整上线"),
]
box(slide, Inches(0.5), Inches(4.65), Inches(12.3), Inches(0.04), fill=LIGHT_BLUE)
for i, (m, label) in enumerate(milestones):
    x = Inches(0.5) + i * Inches(1.87)
    box(slide, x+Inches(0.7), Inches(4.52), Inches(0.25), Inches(0.25),
        fill=BRAND_RED if m in ("M6","M12","M18") else LIGHT_BLUE)
    tb(slide, x+Inches(0.3), Inches(4.85), Inches(1.6), Inches(0.35),
       m, size=12, bold=True, color=BRAND_RED if m in ("M6","M12","M18") else LIGHT_BLUE,
       align=PP_ALIGN.CENTER)
    tb(slide, x+Inches(0.1), Inches(5.22), Inches(1.8), Inches(0.45),
       label, size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

callout_box(slide, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.5),
    "治理节奏：双周Sprint迭代 · 月度管理层汇报 · 季度Steering Committee · 每Phase独立验收",
    [])
notes(slide, "Governance structure. 7-milestone timeline. Biweekly sprints. Quarterly steering committee.")

# ─── DIVIDER: SECTION 5 ─────────────────────────────────────────────────────
divider_slide("SECTION 5 — VALUE & DECISION",
              "价值量化与决策：多赚·少亏·少花",
              "年化1-1.8亿直接价值 · 极端事件1-2亿/次避损 · 3年不行动代价3.3-6亿")

# ─── SLIDE 33: VALUE DASHBOARD ──────────────────────────────────────────────
slide = light_slide("POMS平台价值总览：多赚·少亏·少花——千亿AUM保守估计",
                    "三维价值：年化多赚0.5-1亿 + 年化少花0.5-0.8亿 + 极端事件保守避损1-2亿/次")
# Three value pillars
for i, (icon, title, val1, items, color) in enumerate([
    ("💰", "多赚\n(Alpha增强)",
     "0.5-1亿/年",
     ["全资产配置优化：+0.02%-0.05% → 2000-5000万",
      "相对价值/量化策略：新Alpha源 → 1000-2000万",
      "再平衡效率：跟踪误差-5-10bps → 2000-3000万"],
     TEAL),
    ("🛡", "少亏\n(损失规避)",
     "常规3000-5000万/次\n极端1-2亿/次",
     ["实时回撤控制：从-3%控制到-1.5%",
      "924类极端事件预案：差异高达1.2亿/事件",
      "合规事前拦截：每年避免2-3次违规处罚",
      "现金流预测：避免结算失败（数百万/次）"],
     BRAND_RED),
    ("✂", "少花\n(成本节省)",
     "0.5-0.8亿/年",
     ["交易滑点优化：1000-1500万/年",
      "融资成本精细管理：2000-3000万/年",
      "精确成本核算：发现隐性损失300-500万",
      "信创替代Bloomberg：300-500万/年"],
     BRAND_YELLOW),
]):
    x = Inches(0.5) + i * Inches(4.1)
    box(slide, x, Inches(1.12), Inches(3.9), Inches(4.25), fill=DARK_NAVY, line=color)
    tb(slide, x+Inches(0.1), Inches(1.18), Inches(3.7), Inches(0.55),
       title, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    box(slide, x+Inches(0.3), Inches(1.78), Inches(3.3), Inches(0.04), fill=color)
    tb(slide, x+Inches(0.1), Inches(1.9), Inches(3.7), Inches(0.55),
       val1, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tf = tb(slide, x+Inches(0.1), Inches(2.52), Inches(3.7), Inches(2.7), "").text_frame
    tf.word_wrap = True
    for item in items:
        ap(tf, "▸ " + item, size=10, color=WHITE, space=Pt(5))

callout_box(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.88),
    "年化直接价值：1-1.8亿 · 5年累计：5-9亿（不含风险避损）",
    ["重要前提：价值实现需四个支柱协同——系统平台(POMS) + 人员能力 + 流程优化 + 管理机制",
     "系统是必要条件，不是充分条件——但没有系统，其他三个支柱无从发力"])
notes(slide, "Value dashboard. Three pillars: earn more / lose less / spend less. 1-1.8B annual direct value.")

# ─── SLIDE 34: 924 WAR ROOM ─────────────────────────────────────────────────
slide = light_slide("924 War Room案例：POMS在一天内创造1.2亿价值",
                    "2024年9月24日 央行货币政策转向——真实行业事件，有无POMS差异1.2亿")
headers9 = ["时间", "Without POMS", "With POMS"]
rows9 = [
    ["09:30", "从新闻得知央行干预消息", "CEP引擎捕获利率异常波动，实时预警触发"],
    ["09:31", "—", "压力测试自动启动'924情景' → 组合实时损益: -1.2亿 · 最脆弱因子: 5Y+利率债久期"],
    ["09:35", "—", "系统生成减仓建议：优先前10只高流动性债券 · 合规预检自动通过 · 限额检查OK"],
    ["09:38", "—", "PM确认 → 一键multi-leg执行 → OEMS同步 · 开始成交，滑点仅3bp"],
    ["09:45", "开始手工计算影响，数据是昨天的", "—"],
    ["10:00", "—", "全包成本引擎：减仓交易成本 = 佣金12万+CFETS 1.5万+结算0.8万 = 14.3万"],
    ["10:15", "大致估出可能亏1.5亿，但不确定", "—"],
    ["11:30", "想减仓，不知道哪些券流动性最好", "—"],
    ["11:30", "手写指令单，找领导签字，价格已跌了50bp", "—"],
    ["收盘", "实际亏损2亿 · 事后才看到完整损失 · 不知道全包成本", "实际亏损0.8亿 · 全程实时可见 · 全包成本精确到分"],
]
add_table(slide, Inches(0.5), Inches(1.15), Inches(12.3), Emu(295000), headers9, rows9,
          col_widths=[Inches(0.9), Inches(5.2), Inches(6.2)])
callout_box(slide, Inches(0.5), Inches(5.42), Inches(12.3), Inches(0.92),
    "差异1.2亿——这就是POMS在一天内创造的价值",
    ["核心：流动性冲击比价格下跌更致命——不知道'能卖多快、要付多少代价'才是真正的风险",
     "POMS的三个关键优势：① 事前预案（压力测试） ② 实时流动性分析 ③ 一键执行（减少操作延迟损失）"])
notes(slide, "924 War Room case study. 1.2B difference in one day. Real industry event Sep 24 2024.")

# ─── SLIDE 35: COST OF INACTION ─────────────────────────────────────────────
slide = light_slide("不行动的代价：3年累计潜在代价3.3-8亿",
                    "不是'要不要做'的问题，而是'能不能承受不做的代价'的问题")
for i, (year, items, total, color) in enumerate([
    ("Year 1", [
        "错过跨资产配置机会 → 潜在机会成本0.5-1亿",
        "继续承受隐性交易成本 → 2000-5000万",
        "人工风控合规效率低下 → 1500-3000万",
    ], "Year 1小计：~1-2亿", TEAL),
    ("Year 2", [
        "924类极端事件若发生无预案 → 单次1-2亿",
        "IFRS9 2026全面执行无系统 → P&L波动不可控",
        "继续累计Year 1各项成本 → 1-2亿",
        "新增IFRS9无系统P&L波动 → 0.3-0.5亿",
    ], "Year 2小计：1.3-2.5亿+极端风险敞口", BRAND_YELLOW),
    ("Year 3", [
        "竞争对手平台化深水区 → 差距从'可追赶'变'代差'",
        "头部券商能力拉开 → 优秀投资经理被挖走",
        "客户和资金方注意到差距 → 自营规模可能缩减",
        "持续累计各项直接成本 → 1-1.5亿",
    ], "Year 3小计：1-1.5亿+能力代差（人才利润双损）", BRAND_RED),
]):
    x = Inches(0.5) + i * Inches(4.1)
    box(slide, x, Inches(1.12), Inches(3.9), Inches(4.05), fill=DARK_NAVY, line=color)
    tb(slide, x+Inches(0.1), Inches(1.18), Inches(3.7), Inches(0.42),
       year, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    box(slide, x+Inches(0.3), Inches(1.65), Inches(3.3), Inches(0.03), fill=color)
    tf = tb(slide, x+Inches(0.1), Inches(1.78), Inches(3.7), Inches(2.5), "").text_frame
    tf.word_wrap = True
    for item in items:
        ap(tf, "▸ " + item, size=10, color=WHITE, space=Pt(5))
    box(slide, x+Inches(0.1), Inches(4.35), Inches(3.7), Inches(0.6), fill=RGBColor(0x05, 0x12, 0x4A), line=color)
    tb(slide, x+Inches(0.15), Inches(4.42), Inches(3.6), Inches(0.45),
       total, size=10, bold=True, color=color)

callout_box(slide, Inches(0.5), Inches(5.38), Inches(12.3), Inches(0.97),
    "3年不行动累计代价（保守估算）",
    ["直接成本：3.3-6亿 · 含1次极端事件：4.3-8亿",
     "战略代价：人才流失+竞争力下滑+规模缩减——不可逆的能力代差，量化不了但影响更大"])
notes(slide, "Cost of inaction. 3-year waterfall. 3.3-8B including extreme event. Irreversible strategic gap.")

# ─── SLIDE 36: KPI SCORECARD ────────────────────────────────────────────────
slide = light_slide("KPI度量体系：可衡量才可管理——每个功能都有明确KPI",
                    "投入产出透明化，价值实现可追踪可验收")
headers10 = ["功能模块", "KPI指标", "基线(无POMS)", "目标(有POMS)", "量化价值"]
rows10 = [
    ["全资产组合管理", "Sharpe Ratio", "当前值", "+0.2~0.4", "2000-5000万/年"],
    ["虚拟组合试算", "试算响应时间", "1-2天(Excel)", "<10秒", "决策速度100x↑"],
    ["自动再平衡", "CIO决策→全组合指令", "1-2天", "<5分钟", "2000-3000万/年"],
    ["实时风险监控", "风控延迟", "T+1（24小时）", "<1秒", "924类事件1-2亿↓"],
    ["回撤控制", "年度最大回撤", "无控制", "<-2%", "3000-5000万/次"],
    ["极端压力测试", "情景覆盖率", "0", "日频10+情景", "提前预案"],
    ["合规事前预检", "违规拦截率", "0%（事后发现）", ">95%（事前）", "2-3次处罚/年↓"],
    ["融资成本管理", "加权融资利率", "市场均价", "-5~10bps", "2000-3000万/年"],
    ["执行质量TCA", "vs TWAP基准", "不可见", "持续优化", "500-1000万/年"],
    ["IFRS9管理", "P&L波动率", "当前值", "降低30-50%", "ROE稳定性↑"],
    ["审计追踪", "监管响应时间", "2-4周", "<4小时", "合规风险↓"],
    ["现金流预测", "结算失败率", "X次/年", "0", "数百万/次↓"],
]
add_table(slide, Inches(0.5), Inches(1.15), Inches(12.3), Emu(310000), headers10, rows10,
          col_widths=[Inches(2.5), Inches(2.2), Inches(2.0), Inches(2.0), Inches(3.6)])
callout_box(slide, Inches(0.5), Inches(5.82), Inches(12.3), Inches(0.52),
    "治理承诺：每个Phase验收时，按上述KPI逐项对标——结果说话，不靠PPT",
    [])
notes(slide, "KPI scorecard. 12 measurable metrics. Phase-by-phase validation against each KPI.")

# ─── SLIDE 37: NEXT STEPS ───────────────────────────────────────────────────
slide = light_slide("下一步：三个决策动作——今天讨论什么",
                    "Phase 1六个月启动 · 最小风险验证华锐交付能力 · 奠定千亿自营数字化管理基础")
for i, (num, title, actions, color) in enumerate([
    ("决策①", "立项批准Phase 1", [
        "批准Phase 1项目立项（0-6个月）",
        "确认项目发起人和核心团队名单",
        "明确首批对接IBOR数据范围",
        "华锐提交Phase 1详细SOW和报价",
    ], TEAL),
    ("决策②", "技术调研与验证", [
        "安排华锐与IT架构师技术深研会",
        "IBOR对接接口调研（提前启动最关键）",
        "信创环境测试计划",
        "Phase 1性能基准和验收标准确认",
    ], LIGHT_BLUE),
    ("决策③", "业务代表参与", [
        "自营部投资经理+交易员参与需求工作坊",
        "确认Phase 1功能优先级（哪些先上）",
        "Shadow Portfolio试算场景确认",
        "绩效归因报告格式和内容确认",
    ], BRAND_YELLOW),
]):
    x = Inches(0.5) + i * Inches(4.1)
    box(slide, x, Inches(1.12), Inches(3.9), Inches(4.25), fill=DARK_NAVY, line=color)
    box(slide, x, Inches(1.12), Inches(3.9), Inches(0.55), fill=color)
    tb(slide, x+Inches(0.1), Inches(1.15), Inches(3.7), Inches(0.48),
       num, size=18, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
    tb(slide, x+Inches(0.1), Inches(1.72), Inches(3.7), Inches(0.42),
       title, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    box(slide, x+Inches(0.3), Inches(2.18), Inches(3.3), Inches(0.03), fill=color)
    tf = tb(slide, x+Inches(0.1), Inches(2.3), Inches(3.7), Inches(2.9), "").text_frame
    tf.word_wrap = True
    for a in actions:
        ap(tf, "▸ " + a, size=12, color=WHITE, space=Pt(8))

callout_box(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.88),
    "建议：Phase 1六个月灯塔项目，最小风险验证交付能力",
    ["Phase 1投入最小化，产出最大化——6个月内CEO看到全公司自营实时全貌",
     "基于Phase 1验收结果决定Phase 2——分步建设、结果说话、风险可控"])
notes(slide, "Next steps: 3 decisions. Phase 1 kickoff, technical deep-dive, business rep workshop.")

# ─── SLIDE 38: CLOSING ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(LAYOUT_DARK)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0: ph.text = ""
box(slide, Inches(0), Inches(0), Inches(0.5), H, fill=BRAND_RED)
box(slide, Inches(0.8), Inches(2.85), Inches(8), Inches(0.04), fill=BRAND_RED)
tb(slide, Inches(0.8), Inches(1.5), Inches(12), Inches(0.55),
   "华锐POMS：让千亿资金拥有Aladdin级的数字化管理能力", size=14, color=BRAND_YELLOW)
tb(slide, Inches(0.8), Inches(2.0), Inches(12), Inches(0.85),
   "POMS不是IT投资\n这是投资能力的投资", size=30, bold=True, color=WHITE)
tb(slide, Inches(0.8), Inches(3.05), Inches(11), Inches(0.6),
   "千亿自营资金，每年交易3000亿，管理500亿杠杆——值得拥有世界级的管理工具。", size=14, color=WHITE)

for i, line in enumerate([
    "年化直接价值：1-1.8亿 · 极端事件避损：1-2亿/次 · 3年不行动代价：3.3-8亿",
    "18个月追赶头部3-5年 · Phase 1六个月CEO看到全貌 · 分步建设风险可控",
    "华锐：唯一同时具备FICC+买方组合管理双重基因的国产厂商",
]):
    tb(slide, Inches(0.8), Inches(3.85) + i * Inches(0.45), Inches(12), Inches(0.42),
       "▸ " + line, size=12, color=BRAND_YELLOW if i == 0 else WHITE)

box(slide, Inches(0.8), Inches(5.8), Inches(12), Inches(0.04), fill=BRAND_RED)
tb(slide, Inches(0.8), Inches(5.9), Inches(6), Inches(0.55),
   "华锐科技  ·  上海+深圳+香港+长沙", size=12, color=MED_GRAY)
tb(slide, Inches(7.5), Inches(5.9), Inches(5.3), Inches(0.55),
   "感谢聆听  ·  期待下一步合作", size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
notes(slide, "Closing. POMS = investment in investment capability. Not IT project. 1-1.8B annual direct value.")

# ─── SAVE ────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"✅ Saved: {OUTPUT}")
print(f"   Slides: {len(prs.slides)}")
