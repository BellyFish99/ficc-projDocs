"""
CEO Deck v2.2 — Built on Huarui's actual company PPT template.
Uses the M9 PPT as template source to inherit theme, backgrounds, logos, and fonts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# === Load company template ===
template_path = "MKT2026_AIM_M9.pptx"
# First, copy the template to preserve its theme
import shutil
shutil.copy("/mnt/d/Work/国元/MKT2026_AIM_华锐投资管理平台 M9介绍PPT(现券管理)_20260316.pptx", template_path)

prs = Presentation(template_path)

# Delete all existing slides (keep theme/master/layouts)
xml_slides = prs.slides._sldIdLst
for sldId in list(xml_slides):
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId:
        try:
            prs.part.drop_rel(rId)
        except (KeyError, Exception):
            pass
    xml_slides.remove(sldId)

# === Color constants from Huarui theme ===
DARK_NAVY = RGBColor(0x03, 0x0E, 0x42)      # Primary dark bg
MEDIUM_BLUE = RGBColor(0x15, 0x42, 0x8F)     # Secondary blue
BRAND_RED = RGBColor(0xF5, 0x4D, 0x61)       # Huarui accent red
BRAND_PINK = RGBColor(0xEC, 0x4C, 0x62)      # Text accent
BRAND_YELLOW = RGBColor(0xFF, 0xD0, 0x51)    # Highlight yellow
BRIGHT_YELLOW = RGBColor(0xFF, 0xFF, 0x00)   # Strong highlight
LIGHT_BLUE = RGBColor(0x37, 0x9D, 0xFF)      # Accent blue
BRIGHT_BLUE = RGBColor(0x32, 0x70, 0xFC)     # Links/accent
TEAL = RGBColor(0x56, 0xC4, 0xD0)            # Secondary accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x88, 0x88, 0x88)
TITLE_BLUE = RGBColor(0x28, 0x36, 0x63)      # Title text color

W = prs.slide_width   # 13.333"
H = prs.slide_height  # 7.5"

# === Layout references ===
LAYOUT_COVER = prs.slide_layouts[14]       # '封面'
LAYOUT_TOC = prs.slide_layouts[12]         # '目录'
LAYOUT_DARK = prs.slide_layouts[15]        # '3_内页深色渐变-Logo右上'
LAYOUT_LIGHT = prs.slide_layouts[11]       # '纯色内页（旧版兼容）-右上Logo'
LAYOUT_BLANK = prs.slide_layouts[16]       # '空白页'
LAYOUT_END = prs.slide_layouts[10]         # '末尾幻灯片'


def tb(slide, left, top, width, height, text="", size=14,
       bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT, font="Arial"):
    """Add textbox shorthand."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return txBox


def ap(text_frame, text, size=14, bold=False, color=DARK_GRAY,
       align=PP_ALIGN.LEFT, font="Arial", space=Pt(4)):
    """Add paragraph shorthand."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    p.space_before = space
    return p


def box(slide, left, top, width, height, fill=None, line=None):
    """Add rectangle shape."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def add_table(slide, left, top, width, rh, headers, rows, hdr_color=DARK_NAVY, fs=11):
    """Add formatted table."""
    nr = len(rows) + 1
    nc = len(headers)
    ts = slide.shapes.add_table(nr, nc, left, top, width, Emu(rh * nr))
    t = ts.table
    cw = width // nc
    for i in range(nc):
        t.columns[i].width = cw
    for i, h in enumerate(headers):
        c = t.cell(0, i)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = hdr_color
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(fs)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Arial"
    for r, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            c = t.cell(r + 1, c_idx)
            c.text = str(val)
            if r % 2 == 1:
                c.fill.solid()
                c.fill.fore_color.rgb = LIGHT_GRAY
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(fs)
                p.font.color.rgb = DARK_GRAY
                p.font.name = "Arial"
    return ts


def notes(slide, text):
    """Set speaker notes."""
    slide.notes_slide.notes_text_frame.text = text


def dark_slide(title_text, subtitle_text=""):
    """Create a dark-bg content slide with title."""
    slide = prs.slides.add_slide(LAYOUT_DARK)
    # Set title placeholder if available
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title_text
            for p in ph.text_frame.paragraphs:
                p.font.color.rgb = WHITE
                p.font.name = "Arial"
                p.font.bold = True
            break
    if subtitle_text:
        tb(slide, Inches(0.6), Inches(0.75), Inches(10), Inches(0.4),
           subtitle_text, size=13, color=BRAND_YELLOW, font="Arial")
    return slide


def light_slide(title_text, subtitle_text=""):
    """Create a light-bg content slide with title bar."""
    slide = prs.slides.add_slide(LAYOUT_LIGHT)
    # Title bar
    box(slide, Inches(0), Inches(0), W, Inches(1.1), fill=DARK_NAVY)
    tb(slide, Inches(0.6), Inches(0.12), Inches(11), Inches(0.5),
       title_text, size=22, bold=True, color=WHITE)
    if subtitle_text:
        tb(slide, Inches(0.6), Inches(0.6), Inches(11), Inches(0.35),
           subtitle_text, size=12, color=BRAND_YELLOW)
    return slide


def divider_slide(act_text, subtitle):
    """Create act divider on dark layout."""
    slide = prs.slides.add_slide(LAYOUT_DARK)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""
    # Red accent bar
    box(slide, Inches(5.5), Inches(2.8), Inches(2.3), Inches(0.05), fill=BRAND_RED)
    tb(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1),
       act_text, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Arial")
    tb(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.6),
       subtitle, size=16, color=BRAND_YELLOW, align=PP_ALIGN.CENTER, font="Arial")
    return slide


# ============================================================
# SLIDE 1 — COVER (using company cover layout)
# ============================================================
slide = prs.slides.add_slide(LAYOUT_COVER)
# Add our content on top of the company cover background
tb(slide, Inches(0.8), Inches(2.0), Inches(8), Inches(0.5),
   "/ / /", size=16, color=WHITE, font="Arial")
tb(slide, Inches(0.8), Inches(2.7), Inches(8), Inches(0.8),
   "POMS", size=48, bold=True, color=WHITE, font="Arial")
t = tb(slide, Inches(0.8), Inches(3.4), Inches(8), Inches(0.7),
       "Archforce Investment Management", size=18, color=WHITE, font="Arial")
# Red accent bar
box(slide, Inches(0.88), Inches(4.3), Inches(1.2), Inches(0.025), fill=BRAND_RED)
tb(slide, Inches(0.8), Inches(4.5), Inches(8), Inches(0.4),
   "2026.Q2", size=16, color=WHITE, font="Arial")
notes(slide, "Cover page. Confirm attendees and schedule (~55 min + Q&A).")


# ============================================================
# SLIDE 2 — Executive Summary
# ============================================================
slide = light_slide("CEO Deck V2.2",
                    "Guoyuan POMS")

headers = ["Dimension", "Key Points"]
rows = [
    ["Current State", "Prop trading = #1 revenue (43% in 2025H1), ~100B AUM, clear strategic direction"],
    ["Strategic Need", "Four capability leaps: Optimal Allocation / Fast Computing / Stable Control / Seamless Flow"],
    ["Solution", "Huarui POMS: Portfolio Mgmt + Quant Strategy + Real-time Risk + Order Management + Cost Accounting"],
    ["Ecosystem", "Symbiotic with Kingdom IBOR + Kingstar Market Center -- IBOR = engine, POMS = autopilot"],
    ["Roadmap", "3 phases / 18 months, Phase 1 delivers in 6 months"],
    ["Conservative Value", "Annual direct value 100-180M + Extreme event avoidance 100-200M/event"],
    ["Investment Protection", "Phase 1 Lighthouse: CEO sees full real-time proprietary trading panorama in 6 months"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(380000),
          headers, rows, fs=10)

# Core judgment box
b = box(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.4),
        fill=RGBColor(0x05, 0x12, 0x4A), line=MEDIUM_BLUE)
t = tb(slide, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.35),
       "Core thesis: POMS is not an IT system -- it is investment capability infrastructure.",
       size=13, bold=True, color=WHITE)
ap(t.text_frame,
   "Value realization requires 4 pillars: System + People + Process + Management. System is necessary but not sufficient.",
   size=11, color=BRAND_YELLOW)
ap(t.text_frame,
   "POMS and existing tech ecosystem are symbiotic: Kingdom IBOR = data engine, Kingstar = information highway, POMS = autonomous driving system.",
   size=11, color=RGBColor(0xAA, 0xBB, 0xDD))

notes(slide, "Executive summary in one page. POMS complements the existing IBOR and market center investments.")


# ============================================================
# ACT 1 DIVIDER
# ============================================================
divider_slide("Act 1: Deep Understanding of Guoyuan",
              "We are here to solve a strategic problem, not sell a system")


# ============================================================
# SLIDE 3 — 4x Growth
# ============================================================
slide = light_slide("Guoyuan Prop Trading: 4x Growth in 4 Years, Now #1 Revenue Pillar",
                    "Revenue grew 4x to 43% share; 'De-directional + FICC + Capital Intermediary' strategy confirmed")

headers = ["Year", "Prop Revenue", "Share", "Wealth Mgmt", "Key Milestone"]
rows = [
    ["2022", "448M", "8%", "1.93B", "Strategic pivot: established de-directional strategy"],
    ["2023", "1.31B", "21%", "1.51B", "+192% explosion, FICC intermediary rapid growth"],
    ["2024", "2.31B", "29%", "1.70B", "First time surpassing Wealth as #1 revenue"],
    ["2025H1", "1.46B", "43%", "998M", "Half-year exceeds full-year 2023, 1.46x Wealth"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(380000),
          headers, rows, fs=11)

tb(slide, Inches(0.6), Inches(3.6), Inches(12), Inches(0.35),
   "Aug 2025 Org Restructure -- Three Strategic Signals", size=15, bold=True, color=TITLE_BLUE)

headers2 = ["Change", "Strategic Signal"]
rows2 = [
    ["Prop Trading Committee -> Financial Markets Committee", "From 'investment-centric' to 'market service' -- needs unified portfolio platform"],
    ["Fixed Income Dept -> FICC Business HQ", "From bonds-only to full FICC -- needs cross-asset portfolio + quant"],
    ["Equity Investment -> Securities Investment", "De-directional + multi-strategy -- needs real-time risk + relative value"],
    ["Innovation Finance -> Innovation Finance HQ (elevated)", "Derivatives from margin to core -- needs unified pricing + cost accounting"],
]
add_table(slide, Inches(0.6), Inches(4.1), Inches(12.1), Emu(330000),
          headers2, rows2, fs=10)

notes(slide, "We have studied Guoyuan deeply. Prop trading grew 4x in 4 years, now 43% of revenue. The Aug 2025 org restructure signals the need for a unified portfolio management platform.")


# ============================================================
# SLIDE 4 — Requirements
# ============================================================
slide = light_slide("We Understand Guoyuan's Needs: Licenses Ready, Missing the Unified Platform",
                    "5 years of FICC licenses assembled -- licenses are the passport, POMS is the driving ability")

t = tb(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(1.8), "", size=13)
tf = t.text_frame
tf.word_wrap = True
for i, (title, desc) in enumerate([
    ("1. Stable return target", "Low-vol 6%, from directional to precision management"),
    ("2. Full asset coverage", "Not just fixed income -- equities, gold, bonds, derivatives, true multi-asset"),
    ("3. Rich investment strategies", "Relative value (spread/basis/on-off-the-run) + multi-strategy multi-timeframe"),
    ("4. Portfolio management as core", "Not trade-centric -- this is Aladdin-level thinking"),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"{title} -- {desc}"
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Arial"
    p.space_before = Pt(8)

headers = ["License Acquired", "Business Capability Unlocked", "What Unified Platform Enables"]
rows = [
    ["IRS Trading (2020)", "Interest rate derivatives hedging", "Unified pricing engine for spread strategies"],
    ["OTC Options L2 Dealer (2021)", "Options hedging + structured products", "Cross-asset portfolio risk + full cost accounting"],
    ["CRM Tool Dealer (2022)", "Credit derivatives + bond lending", "Full asset correlation analysis + multi-strategy"],
    ["Listed Securities Market Making (2024)", "Equity market making + liquidity", "Real-time market making risk + limit management"],
    ["Carbon Trading (2024)", "Green finance + FICC extension", "Multi-asset unified portfolio management"],
]
add_table(slide, Inches(0.6), Inches(3.8), Inches(12.1), Emu(330000),
          headers, rows, fs=10)

notes(slide, "Guoyuan has assembled the complete FICC license puzzle over 5 years. POMS is the platform that unlocks the full business value of these licenses.")


# ============================================================
# SLIDE 5 — Four Capability Leaps
# ============================================================
slide = light_slide("Strategic Execution Requires Four Core Capability Leaps",
                    "Optimal Allocation / Fast Computing / Stable Control / Seamless Flow")

# Four colored boxes
cap_data = [
    ("Optimal\nAllocation", TEAL, ["Full asset coverage", "Cross-asset correlation", "Data-driven what-if", "Auto rebalancing", "IFRS9 simulation"]),
    ("Fast\nComputing", LIGHT_BLUE, ["Real-time risk metrics", "Autonomous pricing", "Spread/basis strategies", "Multi-strategy backtesting", "Full-cost accounting"]),
    ("Stable\nControl", BRAND_RED, ["Real-time drawdown alerts", "Extreme stress testing", "Liquidity impact modeling", "Pre-trade compliance", "Real-time limit mgmt"]),
    ("Seamless\nFlow", BRAND_YELLOW, ["Research-to-portfolio", "Auto order generation", "Multi-leg execution", "TCA analysis", "Real-time attribution"]),
]

for i, (title, color, items) in enumerate(cap_data):
    x = Inches(0.5) + i * Inches(3.15)
    y = Inches(1.5)
    bw = Inches(2.9)
    bh = Inches(4.5)
    box(slide, x, y, bw, bh, line=color)
    box(slide, x, y, bw, Inches(0.9), fill=color)
    tb(slide, x, y + Inches(0.05), bw, Inches(0.8),
       title, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    t = tb(slide, x + Inches(0.15), y + Inches(1.0), bw - Inches(0.3), bh - Inches(1.1), "", size=11)
    tf = t.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Arial"
        p.space_before = Pt(8)

tb(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
   "Not about what is wrong -- it is about what the strategy demands. No single tool can solve this -- requires a unified platform.",
   size=11, color=MED_GRAY, align=PP_ALIGN.CENTER)

notes(slide, "These four capability targets are systemic requirements for strategic execution. No single-point tool can address them all.")


# ============================================================
# ACT 2 DIVIDER
# ============================================================
divider_slide("Act 2: Why Act Now",
              "924 Real Lesson + A Day in the Life = Urgency")


# ============================================================
# SLIDE 6 — 924 Case
# ============================================================
slide = light_slide("924 Event Proves: Prop Desks Without Stress Testing Are Flying Blind",
                    "Sep 24, 2024 PBOC bond market intervention -- tools vs. no tools, completely different outcomes")

headers = ["Time", "Without POMS", "With POMS"]
rows = [
    ["09:30", "Learns from news", "CEP engine captures rate anomaly"],
    ["09:31", "--", "Stress test auto-triggers, real-time P&L: -120M"],
    ["09:35", "--", "System generates reduction plan + compliance check"],
    ["09:38", "--", "One-click execute, slippage only 3bp"],
    ["09:45", "Starts manual calc, data is yesterday's", "--"],
    ["10:15", "Estimates ~150M loss, uncertain", "--"],
    ["11:30", "Executes reduction, price dropped further", "--"],
    ["Close", "Loss: 200M", "Loss: 80M -> Diff: 120M in one day"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(320000),
          headers, rows, fs=10)

# Liquidity lesson
b = box(slide, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.0),
        fill=RGBColor(0x2A, 0x0A, 0x12), line=BRAND_RED)
t = tb(slide, Inches(0.9), Inches(4.9), Inches(11.5), Inches(0.35),
       "The Real Lesson: Liquidity Shock Is More Lethal Than Price Decline",
       size=14, bold=True, color=BRAND_RED)
ap(t.text_frame,
   "The real danger was not knowing which bonds could be sold and at what market impact cost.",
   size=12, color=WHITE)
ap(t.text_frame,
   "POMS differentiator: Real-time Liquidity Impact Cost calculation -- evaluates liquidation days and impact cost for top 20 holdings, selects the lowest-impact reduction path.",
   size=12, color=BRAND_YELLOW)

notes(slide, "This is a real industry event from Sep 24, 2024. The key insight: liquidity shock was more lethal than the price decline itself. POMS calculates real-time liquidity impact cost for each position.")


# ============================================================
# SLIDE 7 — Day in Life
# ============================================================
slide = light_slide("Zhang's Day: Same Person, Same Market -- Different Tools, Different Outcomes",
                    "POMS does not replace PMs -- it gives great people better tools")

headers = ["Time", "Without POMS", "With POMS"]
rows = [
    ["08:30", "3 systems, data in Excel/emails", "One screen, correlation suggests +5% gold"],
    ["09:15", "Waiting for yesterday's VaR email", "Real-time VaR 120M, concentration warning"],
    ["10:00", "Excel estimate, gives up due to uncertainty", "Virtual portfolio: volatility -32%"],
    ["11:30", "Spread at 93rd percentile, no tools to act", "Alert + backtest: 78% win rate"],
    ["14:00", "Manual order -> sign -> price changed", "One-click -> compliance -> execute in seconds"],
    ["15:30", "No idea of today's contribution", "Real-time attribution: spread +8bp"],
    ["17:00", "Learns peer ranking only year-end", "Annualized 5.8%, vol 3.1%, top 30%"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(340000),
          headers, rows, fs=10)

# Knowledge retention
b = box(slide, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.2),
        fill=RGBColor(0x0A, 0x1A, 0x0A), line=TEAL)
t = tb(slide, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.35),
       "Knowledge Institutionalization: If Zhang leaves, his Excel models and judgment experience disappear.",
       size=13, bold=True, color=TEAL)
ap(t.text_frame,
   "POMS systematically retains strategy logic, risk parameters, and decision history -- upgrading from 'dependent on individuals' to 'institutional capability'.",
   size=12, color=WHITE)

notes(slide, "Day-in-life comparison. The key addition: knowledge institutionalization -- if a star PM leaves, the system retains the institutional knowledge.")


# ============================================================
# SLIDE 8 — Why Now
# ============================================================
slide = light_slide("Why Act Now: Four Drivers Converging Simultaneously",
                    "Not 'whether to do it' but 'can you afford not to'")

headers = ["Performance Driver", "Regulatory Driver", "Competition Driver", "Localization Window"]
rows = [
    ["No tools -> miss opportunities", "15th Five-Year Plan: digital + local", "Huatai/Ping An have full platforms", "Dependent on Calypso/Bloomberg"],
    ["IFRS9 effective 2026", "2024: 162 broker penalties", "Gap becomes generation gap", "Localization window: 2-3 years"],
    ["Prop is 43% of revenue", "Pre-trade compliance: mandatory", "Top PMs recruited away", "Benchmark: Aladdin, 3yr roadmap"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(430000),
          headers, rows, fs=10)

tb(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(0.35),
   "Benchmark Validation: Unified Portfolio Platform Is the Endgame", size=15, bold=True, color=TITLE_BLUE)

headers2 = ["BlackRock Aladdin", "Huatai 'Elephant'"]
rows2 = [
    ["Manages USD 11.6T AUM", "Started 2020, built in 3 years"],
    ["Core: portfolio mgmt + risk, not just data", "2023: opened platform to external clients"],
    ["Guoyuan POMS vision = Mini-Aladdin", "Guoyuan: 18 months to catch up 3-5 years"],
]
add_table(slide, Inches(0.6), Inches(4.0), Inches(12.1), Emu(360000),
          headers2, rows2, fs=11)

notes(slide, "Four drivers converging. Competition is the most urgent -- Huatai built their platform in 3 years, already exporting capabilities externally.")


# ============================================================
# ACT 3 DIVIDER
# ============================================================
divider_slide("Act 3: How POMS Supports Four Capability Leaps",
              "Architecture -> Capabilities -> Cost -> Ecosystem")


# ============================================================
# SLIDE 9 — Architecture
# ============================================================
slide = light_slide("Huarui POMS: Six Engines Powering Investment Decision Intelligence",
                    "Portfolio Mgmt + Quant Strategy + Real-time Risk + Order Mgmt + Cost Accounting + Compliance")

# Architecture as labeled boxes
# Top layer
box(slide, Inches(3.5), Inches(1.4), Inches(6.3), Inches(0.65), fill=LIGHT_BLUE)
tb(slide, Inches(3.5), Inches(1.43), Inches(6.3), Inches(0.55),
   "Investment Workstation: PM / Trading / Risk / Strategy", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Left
box(slide, Inches(0.5), Inches(2.5), Inches(2.5), Inches(2.8), fill=RGBColor(0xE8, 0xEE, 0xF5))
tb(slide, Inches(0.5), Inches(2.6), Inches(2.5), Inches(2.6),
   "Research\nSystems\n\nStrategy signals\nMacro factors\nAI models", size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Center - POMS
box(slide, Inches(3.5), Inches(2.3), Inches(6.3), Inches(3.5), fill=RGBColor(0x05, 0x14, 0x4E))
tb(slide, Inches(3.5), Inches(2.35), Inches(6.3), Inches(0.35),
   "H U A R U I   P O M S   P L A T F O R M", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

engines = [
    "1  Full Asset Portfolio Management Engine",
    "2  Quant Strategy Engine (Relative Value / Backtesting)",
    "3  Real-time Risk + Extreme Stress Testing",
    "4  Order & Execution Management",
    "5  Performance Attribution + Full Cost Accounting",
    "6  Pre-trade Compliance + Real-time Limit Management",
]
for i, eng in enumerate(engines):
    y = Inches(2.85) + i * Inches(0.42)
    tb(slide, Inches(4.2), y, Inches(5), Inches(0.35),
       eng, size=12, color=BRAND_YELLOW if i < 3 else WHITE)

# Right
box(slide, Inches(10.3), Inches(2.5), Inches(2.5), Inches(2.8), fill=RGBColor(0xE8, 0xEE, 0xF5))
tb(slide, Inches(10.3), Inches(2.6), Inches(2.5), Inches(2.6),
   "OEMS\nSystems\n\nTrade execution\nMulti-leg\nAlgo trading", size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Bottom
box(slide, Inches(3.5), Inches(6.1), Inches(6.3), Inches(0.65), fill=MED_GRAY)
tb(slide, Inches(3.5), Inches(6.13), Inches(6.3), Inches(0.55),
   "Existing Infrastructure: Kingdom IBOR | Kingstar Market Center | Data Bus",
   size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

notes(slide, "Six engines at the core. Connects down to existing Kingdom IBOR and Kingstar infrastructure, up to investment workstation, left to research, right to trade execution.")


# ============================================================
# SLIDES 10-12 — Capabilities (3 slides)
# ============================================================

# Slide 10: Optimal + Fast
slide = light_slide("Optimal Allocation + Fast Computing: Best Portfolio + Every bp Accounted",
                    "Portfolio Management + Quant Strategy + Cost Accounting = The Foundation for Alpha")

headers = ["Capability Need", "POMS Solution", "Key Feature"]
rows = [
    ["Full asset allocation", "Unified multi-asset portfolio mgmt", "Equities/bonds/commodities/derivatives/gold"],
    ["Cross-asset analysis", "Correlation matrix + risk diversification", "Find low-vol, low-correlation combinations"],
    ["Data-driven decisions", "What-if virtual portfolio simulation", "100+ parallel simulations, results in seconds"],
    ["Real-time risk metrics", "Distributed real-time computation", "VaR/Greeks output in seconds"],
    ["Autonomous pricing", "Full-spectrum pricing engine", "Yield curves / volatility surfaces"],
    ["Full-cost accounting", "Precise cost accounting engine", "Financing + commission + CFETS + settlement = true P&L"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(340000),
          headers, rows, fs=10)

# Cost example
b = box(slide, Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.5),
        fill=RGBColor(0x20, 0x15, 0x05), line=BRAND_YELLOW)
t = tb(slide, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.35),
       "Hidden Truth: A 100M Leveraged Trade", size=14, bold=True, color=BRAND_YELLOW)
ap(t.text_frame, "Apparent profit: Sell spread = +400K", size=12, color=TEAL)
ap(t.text_frame, "All-in cost: Financing 444K + Commission 10K + CFETS 1K + Settlement 0.5K = 455.5K", size=12, color=BRAND_RED)
ap(t.text_frame, "True P&L: -55.5K -> Looks like +400K profit, actually -55.5K loss!", size=13, bold=True, color=BRAND_RED)

notes(slide, "Optimal allocation and fast computing are the foundation for alpha generation. The cost accounting example shows hidden losses in leveraged trades.")


# Slide 11: Stable Control
slide = light_slide("Stable Control: Real-time Risk + 924-grade Stress Testing + Compliance",
                    "From T+1 manual to real-time auto-alert; from no playbook to 3-layer stress testing")

headers = ["Capability Need", "POMS Solution", "Key Feature"]
rows = [
    ["Real-time drawdown alert", "CEP real-time monitoring", "Sub-second detection + threshold auto-alert"],
    ["Extreme scenario playbook", "3-layer stress testing engine", "Historical replay / Hypothetical / Reverse stress test"],
    ["Liquidity management", "Liquidity stress testing", "Days-to-liquidate / market impact cost"],
    ["Pre-trade compliance", "Embedded compliance engine", "Pre-trade block + in-flight interception"],
    ["Real-time limit mgmt", "Real-time limit dashboard", "See limit impact before placing order"],
    ["Cash flow forecast", "T+N forward projection", "Funding gap early warning"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(340000),
          headers, rows, fs=10)

tb(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(0.35),
   "Three-Layer Stress Testing", size=15, bold=True, color=TITLE_BLUE)
headers2 = ["Layer", "Method", "Example Scenarios"]
rows2 = [
    ["Layer 1", "Historical Replay", "924 PBOC intervention | Baoshang Bank | 2013 Cash Crunch"],
    ["Layer 2", "Hypothetical", "Rates +100bp | Spreads widen 50bp | Equities -10%"],
    ["Layer 3", "Reverse Stress Test", "'Max I can lose is X' -> What scenario causes it?"],
]
add_table(slide, Inches(0.6), Inches(4.8), Inches(12.1), Emu(360000),
          headers2, rows2, fs=11)

notes(slide, "924 taught us that standard risk management is not enough. Three layers of stress testing. Plus pre-trade compliance -- 162 broker penalties in 2024.")


# Slide 12: Seamless Flow
slide = light_slide("Seamless Flow: Research to Trade Completion, End-to-End",
                    "Research -> Decision -> Order -> Execution -> Attribution, fully connected")

headers = ["Capability Need", "POMS Solution", "Key Feature"]
rows = [
    ["Research to portfolio", "Strategy signal direct trigger", "Factor library / spread signals auto-drive"],
    ["Decision to order", "Auto order generation", "Compliance auto-cleared"],
    ["Multi-leg execution", "Multi-leg trade engine", "Spread trade: buy + sell simultaneously"],
    ["Execution quality", "TCA + execution analysis", "Slippage attribution / TWAP / VWAP"],
    ["Performance evaluation", "Real-time attribution", "Campisi model + cost attribution"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(380000),
          headers, rows, fs=11)

# Before/After flow
b = box(slide, Inches(0.6), Inches(4.2), Inches(12.1), Inches(1.5), line=MEDIUM_BLUE)
t = tb(slide, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.5),
       "Current (1-2 days): Research finds opportunity -> Tell PM -> Manual order -> Get signatures -> Trader executes -> Price has moved",
       size=12, color=BRAND_RED)
ap(t.text_frame,
   "With POMS (minutes): Strategy signal -> Portfolio adjustment -> Auto order -> Compliance check -> Approval -> Execute in seconds",
   size=12, color=TEAL)

notes(slide, "Seamless flow is where real money is saved. Current 1-2 day cycle compressed to minutes.")


# ============================================================
# SLIDE 13 — IFRS9
# ============================================================
slide = light_slide("IFRS9 + Full Cost: See Accounting Impact and True Cost Before Every Trade",
                    "IFRS9 effective 2026 -- without system support, P&L volatility becomes uncontrollable")

t = tb(slide, Inches(0.6), Inches(1.4), Inches(5.5), Inches(3.5), "", size=13)
tf = t.text_frame
tf.word_wrap = True
ap(tf, "IFRS9 Challenge", size=16, bold=True, color=TITLE_BLUE, space=Pt(0))
for item in [
    "Financial assets classified as FVTPL / FVOCI / AC",
    "Classification directly impacts P&L volatility",
    "FVTPL: fair value changes hit P&L (high volatility)",
    "FVOCI: changes go to OCI (stable P&L, no spread capture)",
    "Must know before trading: which book, what P&L impact",
]:
    ap(tf, f"  {item}", size=11, color=DARK_GRAY)

t2 = tb(slide, Inches(6.5), Inches(1.4), Inches(6), Inches(3.5), "", size=13)
tf2 = t2.text_frame
tf2.word_wrap = True
ap(tf2, "Full Cost Management", size=16, bold=True, color=TITLE_BLUE, space=Pt(0))
for item in [
    "50B leverage, optimize repo rate 10bp = 200-300M/yr savings",
    "Commission + CFETS + settlement ~40M/yr, visible and optimizable",
    "Before every trade, auto-display:",
    "   All-in cost + IFRS9 book recommendation + true net return",
]:
    ap(tf2, f"  {item}", size=11, color=DARK_GRAY)

notes(slide, "IFRS9 becomes effective in 2026. Without system support, PMs cannot see the accounting impact of their trades. POMS auto-displays full cost and IFRS9 book recommendation before every order.")


# ============================================================
# SLIDE 14 — Ecosystem
# ============================================================
slide = light_slide("How POMS Integrates Into Guoyuan's Existing Tech Ecosystem",
                    "Not starting over -- IBOR is the engine, POMS is the autonomous driving system")

# Car analogy
b = box(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.0),
        fill=RGBColor(0x05, 0x12, 0x4A))
t = tb(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.35),
       "Smart Car Analogy for Guoyuan's Tech Ecosystem:", size=14, bold=True, color=WHITE)

analogies = [
    ("Kingdom IBOR", "= Engine (data driving force)", WHITE),
    ("Kingstar Market Center", "= Dashboard sensors (real-time data collection)", WHITE),
    ("Kingstar Data Bus", "= Internal bus (information highway)", WHITE),
    ("Hundsun / Hengtai / Fenghu", "= Functional components (valuation / trading / credit)", WHITE),
    ("Huarui POMS", "= AUTONOMOUS DRIVING SYSTEM (turns all data into decisions and actions)", BRAND_YELLOW),
]
for i, (name, desc, color) in enumerate(analogies):
    y = Inches(2.0) + i * Inches(0.42)
    bold = i == 4
    sz = 13 if i == 4 else 12
    tb(slide, Inches(2), y, Inches(9.5), Inches(0.35),
       f"{name}  {desc}", size=sz, bold=bold, color=color)

# Design principles
tb(slide, Inches(0.6), Inches(4.7), Inches(12), Inches(0.35),
   "Design Principles", size=15, bold=True, color=TITLE_BLUE)

headers = ["Principle", "Meaning"]
rows = [
    ["Complement, not replace", "Consumes data from all systems; adds portfolio optimization, strategy, risk, orders on top"],
    ["Add value, not duplicate", "Does not rebuild any existing capability; focuses purely on business intelligence layer"],
    ["Progressive integration", "Phase 1: IBOR + Market Center; Phase 2: OEMS + Research"],
]
add_table(slide, Inches(0.6), Inches(5.2), Inches(12.1), Emu(340000),
          headers, rows, fs=11)

notes(slide, "This is the most important slide. The car analogy: Kingdom IBOR is the engine, Kingstar is the sensors and bus, Huarui POMS is the autonomous driving system. We do not replace any existing component -- we make them work together to produce investment value.")


# ============================================================
# ACT 4 DIVIDER
# ============================================================
divider_slide("Act 4: Why Huarui",
              "The only vendor that can deliver end-to-end")


# ============================================================
# SLIDE 15 — Why Huarui
# ============================================================
slide = light_slide("Huarui: The Only Vendor That Can Unitarily Support All Four Capability Leaps",
                    "Six engines full coverage + KPMG FinTech Top 50 six consecutive years + 128 patents")

headers = ["#", "Differentiation", "Value for 100B AUM"]
rows = [
    ["1", "Not just one module -- six engines full coverage", "One vendor solves everything, no patchwork integration"],
    ["2", "Not just analytics -- covers trading, quant, market making", "From strategy to execution, end-to-end"],
    ["3", "Not just domestic -- Hong Kong international center", "Cross-border without another vendor"],
    ["4", "Not just replacement -- benchmarks Aladdin/SecDB", "100B AUM deserves world-class systems"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(380000),
          headers, rows, fs=11)

t = tb(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(2.5), "", size=13)
tf = t.text_frame
tf.word_wrap = True
ap(tf, "Credentials", size=16, bold=True, color=TITLE_BLUE, space=Pt(0))
for c in [
    "KPMG China Leading FinTech Top 50 (2020-2025, 6 consecutive years)",
    "National Specialized-Sophisticated-New 'Little Giant' (2024, first batch)",
    "IDC China FinTech 50 (2023-2025, 3 consecutive years)",
    "9 years | 500 people | 1B RMB R&D | 128 patents | Shanghai + Shenzhen + Hong Kong",
    "Delivery: Core team from top brokers and international investment banks, Phase 1 full-time on-site",
]:
    ap(tf, f"  {c}", size=11, color=DARK_GRAY)

notes(slide, "Huarui is the only vendor with six engines full coverage. If the CEO asks about competitors: existing 6 vendors do vertical modules; Huarui does the horizontal portfolio intelligence layer -- complementary, not competitive.")


# ============================================================
# SLIDE 16 — Roadmap
# ============================================================
slide = light_slide("Three Phases, 18 Months: Phase 1 Delivers in 6 Months",
                    "Bottom-up, step-by-step -- Phase 1: CEO sees full real-time proprietary panorama")

# Three phase boxes
phases_data = [
    ("Phase 1\n0-6 months\nFoundation", TEAL,
     ["* Portfolio Mgmt MVP", "* Cross-asset correlation", "  IFRS9 basics", "  Full-cost basics", "* Basic risk controls", "* Connect to IBOR"],
     "CEO opens system\nand sees everything"),
    ("Phase 2\n6-12 months\nEmpowerment", LIGHT_BLUE,
     ["Quant strategy tools", "Real-time risk + stress test", "Order management", "Attribution", "Compliance + limits", "Credit risk"],
     "Real-time risk live\nOrders fully connected"),
    ("Phase 3\n12-18 months\nEvolution", BRAND_YELLOW,
     ["Strategy backtesting", "Multi-leg optimization", "TCA continuous improvement", "AI engine", "", ""],
     "Full platform live\nIndustry-leading"),
]

for i, (title, color, items, milestone) in enumerate(phases_data):
    x = Inches(0.5) + i * Inches(4.2)
    box(slide, x, Inches(1.4), Inches(3.8), Inches(4.8), line=color)
    box(slide, x, Inches(1.4), Inches(3.8), Inches(1.0), fill=color)
    tb(slide, x, Inches(1.45), Inches(3.8), Inches(0.9),
       title, size=12, bold=True, color=WHITE if color != BRAND_YELLOW else DARK_NAVY, align=PP_ALIGN.CENTER)

    t = tb(slide, x + Inches(0.15), Inches(2.5), Inches(3.5), Inches(2.5), "", size=11)
    tf = t.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        if not item: continue
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Arial"
        p.space_before = Pt(4)

    box(slide, x + Inches(0.15), Inches(5.0), Inches(3.5), Inches(0.6), fill=LIGHT_GRAY)
    tb(slide, x + Inches(0.15), Inches(5.05), Inches(3.5), Inches(0.5),
       f"Milestone: {milestone}", size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Bottom commitment
tb(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
   "Bottom-line commitment: Phase 1 must pass acceptance before Phase 2 starts. Each phase independently validated. Guoyuan controls the pace.",
   size=12, bold=True, color=BRAND_RED, align=PP_ALIGN.CENTER)

notes(slide, "Three phases, bottom-up. Phase 1 is the key -- 6 months to first business value. Phase 1 must pass before Phase 2 starts.")


# ============================================================
# ACT 5 DIVIDER
# ============================================================
divider_slide("Act 5: Value, Risk, Next Steps",
              "How much / Risk is manageable / Inaction costs more")


# ============================================================
# SLIDE 17 — ROI
# ============================================================
slide = light_slide("Conservative Value: Annual Direct Value 100-180M + Extreme Event Avoidance 100-200M",
                    "Every number is conservative, only counting system-enabled portion -- 9x safety margin")

# Four pillars bar
b = box(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.7),
        fill=RGBColor(0x05, 0x12, 0x4A))
tb(slide, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.6),
   "4 Pillars: System Platform (POMS) + People Capability + Process Optimization + Management Mechanisms -> System is necessary, without it other pillars cannot function",
   size=11, color=BRAND_YELLOW)

# Two value tables
headers_l = ["Alpha Enhancement", "Conservative Annual"]
rows_l = [
    ["Allocation optimization (1/10 of industry benchmark)", "20-50M"],
    ["Quant strategy", "10-20M"],
    ["Rebalancing efficiency", "20-30M"],
    ["Subtotal", "50-100M/yr"],
]
add_table(slide, Inches(0.6), Inches(2.3), Inches(5.8), Emu(320000),
          headers_l, rows_l, hdr_color=TEAL, fs=10)

headers_r = ["Cost Savings", "Conservative Annual"]
rows_r = [
    ["Slippage optimization", "10-15M"],
    ["Financing optimization", "20-30M"],
    ["TCA + cost accounting + efficiency + localization + idle funds", "~15-30M"],
    ["Subtotal", "50-80M/yr"],
]
add_table(slide, Inches(6.7), Inches(2.3), Inches(5.8), Emu(320000),
          headers_r, rows_r, hdr_color=MEDIUM_BLUE, fs=10)

# Summary
headers_s = ["Metric", "Value"]
rows_s = [
    ["Annual direct value (alpha + cost savings)", "100-180M/yr"],
    ["Extreme event conservative avoidance", "100-200M/event"],
    ["Routine risk control avoidance", "30-50M/event"],
    ["5-year cumulative", "500-900M"],
]
add_table(slide, Inches(2.5), Inches(4.5), Inches(8.3), Emu(320000),
          headers_s, rows_s, hdr_color=DARK_NAVY, fs=11)

# Safety margin
b = box(slide, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.65),
        fill=RGBColor(0x20, 0x15, 0x05), line=BRAND_YELLOW)
tb(slide, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.55),
   "Safety margin: Our allocation alpha assumption (0.02%-0.05%) is only 1/10 of industry benchmark (0.3%-0.5%) -> 9x safety buffer even if challenged",
   size=12, bold=True, color=BRAND_YELLOW)

notes(slide, "All numbers are conservative -- allocation alpha is 1/10 of industry benchmark. Even if challenged, we have 9x safety margin. Value realization requires all four pillars: system + people + process + management.")


# ============================================================
# SLIDE 18 — Risk
# ============================================================
slide = light_slide("Risk & Mitigation: Anticipate and Resolve in Advance",
                    "Every risk has a corresponding safeguard -- the biggest protection is Phase 1 independent validation")

headers = ["Risk", "Mitigation"]
rows = [
    ["Data quality & integration", "Phase 1 Month 1: dedicated data governance, field-by-field validation"],
    ["Business fit", "On-site requirements + bi-weekly iteration, business team full UAT participation"],
    ["Integration complexity", "Standard API + adapter layer; Phase 1 prioritizes IBOR + Market Center"],
    ["User adoption", "2-3 seed PMs deeply involved in design; 'Excel replacement first, then capability upgrade'"],
    ["Existing vendor coordination", "POMS is the intelligence layer -- complementary to Kingdom/Kingstar/Hundsun"],
    ["Project timeline", "Three phases decoupled; Phase 1 must pass before Phase 2 starts"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(400000),
          headers, rows, fs=11)

notes(slide, "Every risk has a mitigation plan. Most importantly: POMS is complementary to the existing 6 vendors, not competitive. And Phase 1 must pass acceptance before Phase 2.")


# ============================================================
# SLIDE 19 — Cost of Inaction
# ============================================================
slide = light_slide("The Cost of Inaction Far Exceeds the Cost of Action",
                    "3 years of inaction: conservative cost 330-600M, plus extreme event exposure")

headers = ["Year", "Cost of Inaction", "Cumulative"]
rows = [
    ["Year 1", "Missed allocation (50-100M) + Hidden costs (20-50M) + Efficiency loss (15-30M)", "~100-200M"],
    ["Year 2", "Y1 costs continue + IFRS9 P&L volatility + Extreme event if occurs (100-200M)", "130-250M+"],
    ["Year 3", "Competitors launch platforms -> Gap becomes 'generation gap' -> Top PMs recruited away", "Irreversible"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Emu(450000),
          headers, rows, fs=11)

# Summary box
b = box(slide, Inches(2), Inches(3.7), Inches(9.3), Inches(2.8),
        fill=RGBColor(0x2A, 0x0A, 0x0A), line=BRAND_RED)
t = tb(slide, Inches(2.3), Inches(3.8), Inches(8.7), Inches(0.4),
       "3-Year Cumulative Cost of Inaction", size=20, bold=True, color=BRAND_RED, align=PP_ALIGN.CENTER)
for line in [
    "Direct cost: Conservative 330-600M",
    "+ One extreme event: 100-200M",
    "+ Talent loss + competitiveness decline -> Irreversible",
    "",
    "vs Platform investment -> Cost of inaction far exceeds investment",
]:
    ap(t.text_frame, line, size=14, color=WHITE, align=PP_ALIGN.CENTER)

notes(slide, "Reframe the conversation: this is not about whether to invest, but whether you can afford NOT to. 3-year cumulative cost is 330-600M minimum, not counting extreme events or talent loss.")


# ============================================================
# SLIDE 20 — Next Steps
# ============================================================
slide = light_slide("Recommended Next Steps",
                    "Start Phase 1 immediately -- see results in 6 months")

steps = [
    ("1. Joint Task Force (Week 1)", "Establish Business + IT joint working group\nAlign project owners, communication cadence, decision process", TEAL),
    ("2. Requirements Workshop (Weeks 2-3)", "Confirm asset scope, Kingdom IBOR interface plan, priorities\nSelect 2-3 seed PMs for deep Phase 1 involvement", LIGHT_BLUE),
    ("3. Phase 1 Detailed Plan (Month 1)", "Technical architecture (incl. integration with 6 existing vendors)\nData integration plan + seed user training plan", BRIGHT_BLUE),
    ("4. Phase 1 Launch -> 6 Month Delivery", "CEO Dashboard live with 3 core KPIs:\n(1) Total Prop Real-time P&L  (2) Portfolio VaR  (3) Cross-Asset Correlation Heatmap", BRAND_YELLOW),
]

for i, (title, desc, color) in enumerate(steps):
    y = Inches(1.4) + i * Inches(1.35)
    box(slide, Inches(0.6), y, Inches(0.08), Inches(1.1), fill=color)
    tb(slide, Inches(1.0), y, Inches(11), Inches(0.35),
       title, size=15, bold=True, color=TITLE_BLUE)
    tb(slide, Inches(1.0), y + Inches(0.4), Inches(11), Inches(0.8),
       desc, size=12, color=DARK_GRAY)

# Contact
tb(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
   "Huarui Technology  |  [Contact]  |  [Phone]  |  [Email]",
   size=14, color=MED_GRAY, align=PP_ALIGN.CENTER)

notes(slide, "Four steps. Step 1 is most critical: joint task force in week 1. Step 4 delivers 3 CEO dashboard KPIs in 6 months: real-time P&L, portfolio VaR, cross-asset correlation heatmap. Thank you.")


# ============================================================
# END SLIDE
# ============================================================
slide = prs.slides.add_slide(LAYOUT_END)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "Thank You"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Arial"
    elif ph.placeholder_format.idx == 13:
        ph.text = "Huarui Technology | Archforce Investment Management"
        for p in ph.text_frame.paragraphs:
            p.font.color.rgb = BRAND_YELLOW
            p.font.name = "Arial"


# ============================================================
# SAVE
# ============================================================
import os
os.remove(template_path)  # Clean up temp copy

output = "CEO_Deck_v2.2.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
